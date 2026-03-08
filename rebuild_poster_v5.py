#!/usr/bin/env python3
"""
ISEF Poster PPTX v5 — Reference-matching layout.

Goals:
  - Match the reference biophysics poster exactly in structure
  - Same balance of text vs figures (very text heavy, figures integrated inline)
  - Same section title sizing, padding, margins
  - Dense content from the research paper
  - Bullet points AND paragraphs matching the reference style
  - Section layout: LEFT (Visual Abstract, Introduction, Research Objectives, Problem Framing)
                    CENTER (Engineering Methodology top, Results & Validation bottom)
                    RIGHT (Statistical Analysis, Conclusions, Future Work, Key References)
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
import os

# ═══════════════════════════════════════════════════════════════
# COLOR PALETTE — matches reference
# ═══════════════════════════════════════════════════════════════
BG_NAVY   = RGBColor(0x06, 0x2B, 0x7A)
TITLE_BAND= RGBColor(0x03, 0x20, 0x61)
SEC_BAR   = RGBColor(0x0A, 0x3D, 0x9A)
CARD_FILL = RGBColor(0xF7, 0xF9, 0xFC)
CARD_BDR  = RGBColor(0xB7, 0xC5, 0xE3)
TXT_DARK  = RGBColor(0x0B, 0x17, 0x36)
TXT_WHITE = RGBColor(0xFF, 0xFF, 0xFF)
RED       = RGBColor(0xD7, 0x26, 0x3D)
TEAL      = RGBColor(0x00, 0x8C, 0x9E)
GOLD      = RGBColor(0xCF, 0xA5, 0x35)
EQ_BG     = RGBColor(0xE8, 0xEE, 0xF2)
ALT_ROW   = RGBColor(0xEE, 0xF2, 0xFA)
DARK_CARD = RGBColor(0x1A, 0x3A, 0x6E)  # darker card for special callouts

# ═══════════════════════════════════════════════════════════════
# DIMENSIONS — reference poster tight layout
# ═══════════════════════════════════════════════════════════════
BW = 48.0; BH = 36.0
TITLE_H = 3.00              # compact title band like reference
CT = TITLE_H + 0.08         # content top (tight)
CB = 35.85                   # content bottom
M  = 0.15                    # margin from panel edge (very tight like ref)
SG = 0.06                    # section gap between boxes
HH = 0.50                    # section header bar height
CP = 0.10                    # card internal padding

# Panel boundaries
LX = M;         LW = 12.0 - 2*M    # 11.70
CX = 12.0 + M;  CW = 24.0 - 2*M    # 23.70
RX = 36.0 + M;  RW = 12.0 - 2*M    # 11.70

CARD_R = 35000  # corner radius EMU

BASE = os.path.dirname(os.path.abspath(__file__))
EXTRACTED = os.path.join(BASE, "poster_images_extracted")
POSTER_FIGS = os.path.join(BASE, "poster_figures_v5")
OLD_FIGS = os.path.join(BASE, "poster_figures")
ALL_FIGS = os.path.join(BASE, "figures")
OUTPUT = os.path.join(BASE, "SASTO_ISEF_Poster_v5.pptx")

def find_img(name):
    """Find image by name across all figure directories."""
    for d in [POSTER_FIGS, OLD_FIGS, ALL_FIGS, EXTRACTED]:
        for ext in (".png", ".jpg"):
            p = os.path.join(d, name + ext)
            if os.path.isfile(p): return p
    return None

def img_ext(n):
    """Find extracted image by number."""
    for ext in (".jpg", ".png"):
        p = os.path.join(EXTRACTED, f"Image {n}" + ext)
        if os.path.isfile(p): return p
    return None

# ═══════════════════════════════════════════════════════════════
# HELPERS
# ═══════════════════════════════════════════════════════════════

def _set_radius(shape, r=CARD_R):
    spPr = shape._element.spPr
    pg = spPr.find(qn("a:prstGeom"))
    if pg is not None:
        pg.set("prst", "roundRect")
        av = pg.find(qn("a:avLst"))
        if av is None:
            av = spPr.makeelement(qn("a:avLst"), {})
            pg.append(av)
        for c in list(av): av.remove(c)
        gd = spPr.makeelement(qn("a:gd"), {"name":"adj","fmla":f"val {r}"})
        av.append(gd)

def rect(sl, x, y, w, h, fill=None, border=None, bw=Pt(1), r=CARD_R):
    s = sl.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    if fill: s.fill.solid(); s.fill.fore_color.rgb = fill
    else: s.fill.background()
    if border: s.line.color.rgb = border; s.line.width = bw
    else: s.line.fill.background()
    _set_radius(s, r)
    return s

def card(sl, x, y, w, h):
    return rect(sl, x, y, w, h, fill=CARD_FILL, border=CARD_BDR)

def add_run(p, text, font="Arial", sz=Pt(12), bold=False, italic=False, color=TXT_DARK, underline=False):
    """Add a run to an existing paragraph (for mixed formatting)."""
    run = p.add_run()
    run.text = text
    run.font.name = font
    run.font.size = sz
    run.font.bold = bold
    run.font.italic = italic
    if color: run.font.color.rgb = color
    run.font.underline = underline
    return run

def tb(sl, x, y, w, h, text, font="Arial", sz=Pt(12), bold=False, italic=False,
       color=TXT_DARK, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP):
    t = sl.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = t.text_frame; tf.word_wrap = True; tf.auto_size = None
    tf.margin_left = Inches(0.02); tf.margin_right = Inches(0.02)
    tf.margin_top = Inches(0.02); tf.margin_bottom = Inches(0.02)
    p = tf.paragraphs[0]; p.text = text
    p.font.name = font; p.font.size = sz; p.font.bold = bold
    p.font.italic = italic
    if color: p.font.color.rgb = color
    p.alignment = align
    p.space_before = Pt(0); p.space_after = Pt(0)
    p.line_spacing = Pt(sz.pt * 1.15) if hasattr(sz, 'pt') else None
    return t

def mtb(sl, x, y, w, h, paras):
    """Multi-paragraph text box. paras = list of dicts."""
    t = sl.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = t.text_frame; tf.word_wrap = True; tf.auto_size = None
    tf.margin_left = Inches(0.02); tf.margin_right = Inches(0.02)
    tf.margin_top = Inches(0.02); tf.margin_bottom = Inches(0.02)
    for i, pd in enumerate(paras):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = pd.get("text", "")
        p.font.name = pd.get("font", "Arial")
        p.font.size = pd.get("sz", Pt(12))
        p.font.bold = pd.get("bold", False)
        p.font.italic = pd.get("italic", False)
        if pd.get("color"): p.font.color.rgb = pd["color"]
        p.alignment = pd.get("align", PP_ALIGN.LEFT)
        p.space_before = Pt(pd.get("sb", 1))
        p.space_after = Pt(pd.get("sa", 1))
        ls = pd.get("ls")
        if ls: p.line_spacing = Pt(ls)
    return t

def hdr(sl, x, y, w, text):
    """Section header bar — mimics reference: bold ALL CAPS white on dark blue."""
    rect(sl, x, y, w, HH, fill=SEC_BAR, r=12000)
    t = sl.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(HH))
    tf = t.text_frame; tf.word_wrap = False
    tf.margin_left = Inches(0.15); tf.margin_top = Inches(0.06)
    p = tf.paragraphs[0]; p.text = text.upper()
    p.font.name = "Arial Black"; p.font.size = Pt(24); p.font.bold = True
    p.font.color.rgb = TXT_WHITE; p.alignment = PP_ALIGN.CENTER
    p.space_before = Pt(0); p.space_after = Pt(0)
    return y + HH

def img_safe(sl, f, x, y, w, h):
    if not f or not os.path.isfile(f): return None
    try:
        from PIL import Image as PILImage
        with PILImage.open(f) as im: iw, ih = im.size
        ar = iw / ih; bar = w / h
        if ar > bar: fw=w; fh=w/ar; fx=x; fy=y+(h-fh)/2
        else: fh=h; fw=h*ar; fx=x+(w-fw)/2; fy=y
        return sl.shapes.add_picture(f, Inches(fx), Inches(fy), Inches(fw), Inches(fh))
    except Exception as e:
        print(f"  Warning: Could not place image {f}: {e}")
        return None

def tbl(sl, x, y, w, rows, col_pcts, rh=0.24):
    """Create formatted table."""
    nr = len(rows); nc = len(rows[0])
    th = nr * rh
    ts = sl.shapes.add_table(nr, nc, Inches(x), Inches(y), Inches(w), Inches(th))
    t = ts.table
    for ci, pct in enumerate(col_pcts):
        t.columns[ci].width = Inches(w * pct)
    for ri, row in enumerate(rows):
        for ci, txt in enumerate(row):
            c = t.cell(ri, ci); c.text = str(txt)
            for p in c.text_frame.paragraphs:
                p.font.name = "Arial"; p.font.size = Pt(9)
                p.alignment = PP_ALIGN.LEFT
                if ri == 0:
                    p.font.bold = True; p.font.color.rgb = TXT_WHITE; p.font.size = Pt(9)
                else:
                    p.font.color.rgb = TXT_DARK
            if ri == 0: c.fill.solid(); c.fill.fore_color.rgb = SEC_BAR
            elif ri % 2 == 0: c.fill.solid(); c.fill.fore_color.rgb = ALT_ROW
            else: c.fill.solid(); c.fill.fore_color.rgb = CARD_FILL
    return ts, y + th

def eq_pill(sl, x, y, w, text, h=0.35):
    rect(sl, x, y, w, h, fill=EQ_BG, r=20000)
    tb(sl, x+0.05, y+0.03, w-0.10, h-0.06, text, sz=Pt(10), align=PP_ALIGN.CENTER)
    return y + h

def badge(sl, x, y, w, h, text, color=TEAL):
    rect(sl, x, y, w, h, fill=color, r=25000)
    tb(sl, x, y, w, h, text, sz=Pt(9), bold=True, color=TXT_WHITE, align=PP_ALIGN.CENTER)

def cap(sl, x, y, w, text):
    tb(sl, x, y, w, 0.22, text, sz=Pt(8), italic=True, color=TXT_DARK, align=PP_ALIGN.CENTER)
    return y + 0.22

# ═══════════════════════════════════════════════════════════════
# BUILD
# ═══════════════════════════════════════════════════════════════

def build():
    prs = Presentation()
    prs.slide_width = Inches(BW); prs.slide_height = Inches(BH)
    sl = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    print("Building SASTO ISEF Poster v5...")

    # ─── BACKGROUND ───
    bg = sl.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(BW), Inches(BH))
    bg.fill.solid(); bg.fill.fore_color.rgb = BG_NAVY; bg.line.fill.background()

    # ═══════════════════════════════════════════════════════════
    # TITLE BAND (matching reference: 3 rows centered)
    # ═══════════════════════════════════════════════════════════
    tband = sl.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(BW), Inches(TITLE_H))
    tband.fill.solid(); tband.fill.fore_color.rgb = TITLE_BAND; tband.line.fill.background()

    # Gold accent line at bottom of title band
    gold_line = sl.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(TITLE_H - 0.03), Inches(BW), Inches(0.03))
    gold_line.fill.solid(); gold_line.fill.fore_color.rgb = GOLD; gold_line.line.fill.background()

    # Left: house renders (2 small images like user had in v3)
    img_safe(sl, img_ext(0), 0.40, 0.30, 2.80, 2.20)  # house image left
    img_safe(sl, img_ext(11), 4.00, 0.30, 2.60, 2.20)  # house image right (before/after)

    # Center title text
    tb(sl, 7.50, 0.20, 28.00, 1.00,
       "SURROGATE-ACCELERATED STRUCTURAL OPTIMIZATION",
       font="Arial Black", sz=Pt(52), bold=True, color=TXT_WHITE, align=PP_ALIGN.CENTER)

    tb(sl, 7.50, 1.15, 28.00, 0.70,
       "Additive Manufacturing: Harnessing FEA to Optimize Material Efficiency",
       font="Arial", sz=Pt(26), bold=True, italic=True, color=TXT_WHITE, align=PP_ALIGN.CENTER)

    tb(sl, 7.50, 1.90, 28.00, 0.50,
       "Eric Hou",
       font="Arial", sz=Pt(22), bold=True, color=TXT_WHITE, align=PP_ALIGN.CENTER)

    # Right: Credit line
    credit_lines = [
        {"text": "Credit Line of Origin", "sz": Pt(11), "bold": True, "color": TXT_WHITE, "sb": 0, "sa": 0},
        {"text": "Credit Line of Origin", "sz": Pt(11), "bold": True, "color": TXT_WHITE, "sb": 0, "sa": 2},
        {"text": "References & Data: in references below.", "sz": Pt(9), "italic": True, "color": TXT_WHITE, "sb": 0, "sa": 1},
        {"text": "* Images denoted with asterisk are part of the", "sz": Pt(9), "color": TXT_WHITE, "sb": 0, "sa": 0},
        {"text": "public domain. Some are denoted by Wikimedia", "sz": Pt(9), "color": TXT_WHITE, "sb": 0, "sa": 0},
        {"text": "Commons CC-BY-SA 3.0.", "sz": Pt(9), "color": TXT_WHITE, "sb": 0, "sa": 1},
        {"text": "All other graphics, tables, and images have been", "sz": Pt(9), "color": TXT_WHITE, "sb": 0, "sa": 0},
        {"text": "created by Eric Hou, 2026 unless otherwise attributed.", "sz": Pt(9), "color": TXT_WHITE, "sb": 0, "sa": 0},
    ]
    mtb(sl, 38.50, 0.15, 9.20, 2.70, credit_lines)

    # ═══════════════════════════════════════════════════════════
    # LEFT PANEL
    # ═══════════════════════════════════════════════════════════
    y = CT
    px = LX; pw = LW

    # ── L1: VISUAL ABSTRACT ──
    y = hdr(sl, px, y, pw, "Visual Abstract")
    cy = y; ch = 6.80
    card(sl, px, cy, pw, ch)

    # Visual Abstract pipeline — use Image 1 (visual abstract) or pipeline figure
    va_img = find_img("fig_sasto_pipeline") or find_img("fig_visual_abstract") or img_ext(1)
    img_safe(sl, va_img, px+0.15, cy+0.10, pw-0.30, 3.80)

    # Labels under pipeline
    pipeline_text = [
        {"text": "a) User Workflow", "sz": Pt(10), "bold": True, "color": TXT_DARK, "sb": 0, "sa": 1},
        {"text": "Input: 3DWire building wireframe skeleton", "sz": Pt(9), "color": TXT_DARK, "sb": 0, "sa": 0},
        {"text": "Process: Extrude walls + roof + floor -> Boolean fusion -> FEA simulation", "sz": Pt(9), "color": TXT_DARK, "sb": 0, "sa": 0},
        {"text": "-> Train deep ensemble surrogate (5 x 8.76M params) -> Sensitivity-guided erosion", "sz": Pt(9), "color": TXT_DARK, "sb": 0, "sa": 0},
        {"text": "Output: Watertight optimized STL, single-component mesh", "sz": Pt(9), "color": TXT_DARK, "sb": 0, "sa": 1},
        {"text": "b) Model Creation", "sz": Pt(10), "bold": True, "color": TXT_DARK, "sb": 1, "sa": 1},
        {"text": "14,293 wireframes -> 11,178 valid FEA simulations -> 128-cube labeled voxel grids", "sz": Pt(9), "color": TXT_DARK, "sb": 0, "sa": 0},
        {"text": "Family-aware split: 8,943 train / 1,121 val / 1,114 test", "sz": Pt(9), "color": TXT_DARK, "sb": 0, "sa": 0},
    ]
    mtb(sl, px+CP, cy+4.00, pw-2*CP, 2.70, pipeline_text)

    y = cy + ch + SG
    print(f"  L1 Visual Abstract -> y={y:.2f}")

    # ── L2: INTRODUCTION ──
    y = hdr(sl, px, y, pw, "Introduction")
    cy = y; ch = 5.90
    card(sl, px, cy, pw, ch)

    # Images on right side (like reference has brain images on left)
    intro_img = find_img("fig_wireframe_pipeline") or img_ext(2)
    img_safe(sl, intro_img, px+pw-4.30, cy+0.15, 4.10, 2.50)

    intro_text = [
        {"text": "CONCRETE & CONSTRUCTION", "sz": Pt(12), "bold": True, "color": TXT_DARK, "sb": 0, "sa": 2},
        {"text": "Concrete production accounts for approximately 8% of global CO2 emissions [IEA 2021]. Conventional construction uses uniform-thickness walls determined by formwork constraints, not structural need \u2014 a substantial source of wasted material.", "sz": Pt(10), "color": TXT_DARK, "sb": 0, "sa": 3, "ls": 12},
        {"text": "ADDITIVE MANUFACTURING OPPORTUNITY", "sz": Pt(12), "bold": True, "color": TXT_DARK, "sb": 2, "sa": 2},
        {"text": "Large-scale 3D concrete printing (ICON, COBOD, Apis Cor) can realize arbitrary wall profiles at no marginal tooling cost. This enables topology-optimized structures that place material only where structurally required.", "sz": Pt(10), "color": TXT_DARK, "sb": 0, "sa": 3, "ls": 12},
        {"text": "THE COMPUTATIONAL BOTTLENECK", "sz": Pt(12), "bold": True, "color": TXT_DARK, "sb": 2, "sa": 2},
        {"text": "Classical topology optimization (SIMP) requires 100s-1000s of FEA solves \u2014 each taking minutes to hours at building scale \u2014 making it computationally intractable. Voxel-based methods produce disconnected mesh fragments incompatible with 3D printing toolpaths.", "sz": Pt(10), "color": TXT_DARK, "sb": 0, "sa": 3, "ls": 12},
    ]
    mtb(sl, px+CP, cy+0.15, pw-4.60, ch-0.30, intro_text)

    # Equation: Fisher-Kolmogorov analog
    eq_pill(sl, px+CP, cy+ch-1.40, pw-2*CP,
            "Ku = f  (linear elastic FEA: K = stiffness, u = displacement, f = forces)", h=0.30)

    # Numbered list of problems
    problems_text = [
        {"text": "(1) Full FEA requires 100s of computationally intensive solves (days/weeks)", "sz": Pt(9), "color": TXT_DARK, "sb": 0, "sa": 0},
        {"text": "(2) Models do not consider topology preservation for AM compatibility", "sz": Pt(9), "color": TXT_DARK, "sb": 0, "sa": 0},
        {"text": "(3) Are restricted to uniform-thickness designs, missing heterogeneous material opportunities", "sz": Pt(9), "color": TXT_DARK, "sb": 0, "sa": 1},
    ]
    mtb(sl, px+CP, cy+ch-1.05, pw-2*CP, 1.00, problems_text)

    y = cy + ch + SG
    print(f"  L2 Introduction -> y={y:.2f}")

    # ── L3: RESEARCH OBJECTIVES ──
    y = hdr(sl, px, y, pw, "Research Objectives")
    cy = y; ch = 5.80
    card(sl, px, cy, pw, ch)

    obj_text = [
        {"text": "ENGINEERING PROBLEM \u2013 Classical topology optimization (SIMP) at building scale requires 100s-1000s of costly FEA solves; current surrogate methods lack topology guarantees for 3D printing.", "sz": Pt(10), "bold": False, "color": TXT_DARK, "sb": 0, "sa": 3, "ls": 12},
        {"text": "RESEARCH QUESTION \u2013 Can a deep learning surrogate enable fast, printable, structurally safe topology optimization at building scale?", "sz": Pt(10), "bold": False, "color": TXT_DARK, "sb": 0, "sa": 3, "ls": 12},
        {"text": "CLAIM \u2013 This research develops SASTO: Surrogate-Accelerated Sensitivity Topology Optimization. It replaces iterative FEA with a 5-member deep ensemble providing uncertainty-aware structural predictions, enabling 23-92x faster optimization with formal mesh connectivity guarantees.", "sz": Pt(10), "bold": False, "color": TXT_DARK, "sb": 0, "sa": 4, "ls": 12},
    ]
    mtb(sl, px+CP, cy+0.10, pw-2*CP, 2.60, obj_text)

    # Design criteria boxes (like reference's Robustness/Functionality/Specification)
    bw = (pw - 2*CP - 0.20) / 3  # 3 boxes
    bx = px + CP
    by = cy + 2.80
    bh = 1.10
    for i, (title, desc) in enumerate([
        ("Speed:", "Deep ensemble surrogate\n23-92x faster than SIMP\n50s median runtime"),
        ("Printability:", "6-connectivity criterion\nSingle-component STL\nWatertight meshes"),
        ("Safety:", "Uncertainty-aware (k=1.0)\n0/1,114 FEA violations\nP(violation) <= 0.09%"),
    ]):
        bxi = bx + i*(bw + 0.10)
        rect(sl, bxi, by, bw, bh, fill=DARK_CARD, border=CARD_BDR, r=20000)
        tb(sl, bxi, by+0.05, bw, 0.25, title, sz=Pt(11), bold=True, color=GOLD, align=PP_ALIGN.CENTER)
        tb(sl, bxi+0.05, by+0.30, bw-0.10, bh-0.35, desc, sz=Pt(8), color=TXT_WHITE, align=PP_ALIGN.CENTER)

    # Primary/Secondary endpoints + brainstorming (like reference)
    endpt_text = [
        {"text": "Primary Endpoints: Minimize material volume while satisfying structural constraints (stress, compliance, displacement) with zero violations across 1,114 test geometries.", "sz": Pt(9), "color": TXT_DARK, "sb": 2, "sa": 2, "ls": 11},
        {"text": "Secondary Endpoints: Demonstrate effectiveness of deep ensemble surrogates for building-scale topology optimization.", "sz": Pt(9), "color": TXT_DARK, "sb": 0, "sa": 2, "ls": 11},
        {"text": "BRAINSTORMING A SOLUTION \u2013 How should I achieve the criteria?", "sz": Pt(10), "bold": True, "color": TXT_DARK, "sb": 2, "sa": 1},
        {"text": "\u2022 Data Generation: Wireframe-to-voxel pipeline creating 11,178 FEA simulations from 3DWire dataset.", "sz": Pt(9), "color": TXT_DARK, "sb": 0, "sa": 0},
        {"text": "\u2022 Surrogate Training: 5-member deep ensemble (43.8M total params) predicting stress, displacement, compliance.", "sz": Pt(9), "color": TXT_DARK, "sb": 0, "sa": 0},
        {"text": "\u2022 Optimization: Sensitivity-guided erosion with adaptive batch sizing + 6-connectivity preservation.", "sz": Pt(9), "color": TXT_DARK, "sb": 0, "sa": 0},
        {"text": "\u2022 Validation: Independent hex8 FEA re-analysis + conformal prediction bounds.", "sz": Pt(9), "color": TXT_DARK, "sb": 0, "sa": 0},
    ]
    mtb(sl, px+CP, cy+4.00, pw-2*CP, 1.75, endpt_text)

    y = cy + ch + SG
    print(f"  L3 Research Objectives -> y={y:.2f}")

    # ── L4: PROBLEM FRAMING ──
    y = hdr(sl, px, y, pw, "Problem Framing")
    cy = y; ch = CB - y   # fill to bottom
    card(sl, px, cy, pw, ch)

    # Problem framing image (encoder-decoder schematic) -> Image 7 or structural parts
    pf_img = img_ext(3) or img_ext(7)  # structural parts image
    img_safe(sl, pf_img, px+0.20, cy+0.15, 4.50, 3.50)

    pf_text = [
        {"text": "Optimization Objective", "sz": Pt(11), "bold": True, "color": TXT_DARK, "sb": 0, "sa": 2},
        {"text": "min J(p) = wV(V/V0) + wS(S/V0) + Pconstraint(p)", "sz": Pt(10), "italic": True, "color": TXT_DARK, "sb": 0, "sa": 2},
        {"text": "Three-term objective: volume minimization + surface smoothness regularizer + structural penalty aggregating stress, compliance, and displacement violations.", "sz": Pt(9), "color": TXT_DARK, "sb": 0, "sa": 3, "ls": 11},
        {"text": "Sensitivity Computation", "sz": Pt(11), "bold": True, "color": TXT_DARK, "sb": 2, "sa": 2},
        {"text": "si = (1/5) * Sum_m d/dp_i [fm(C) + 0.3 * fm(sigma)]", "sz": Pt(10), "italic": True, "color": TXT_DARK, "sb": 0, "sa": 2},
        {"text": "si > 0: safe to remove (more dead-load penalty than stiffness benefit)", "sz": Pt(9), "color": TEAL, "sb": 0, "sa": 0},
        {"text": "si < 0: structurally essential (removal violates constraints)", "sz": Pt(9), "color": RED, "sb": 0, "sa": 3},
        {"text": "Part-Aware Thickness", "sz": Pt(11), "bold": True, "color": TXT_DARK, "sb": 2, "sa": 2},
        {"text": "tmin = 2*dx ~ 156mm (exterior walls, roof, floor)", "sz": Pt(9), "color": TXT_DARK, "sb": 0, "sa": 0},
        {"text": "tmin = 1*dx ~ 78mm  (interior walls only)", "sz": Pt(9), "color": TXT_DARK, "sb": 0, "sa": 1},
        {"text": "Interior partitions bear negligible load; permitting thinner walls yields 86.8% interior removal in the reference case.", "sz": Pt(9), "color": TXT_DARK, "sb": 0, "sa": 0},
    ]
    mtb(sl, px+CP+4.60, cy+0.15, pw-4.80-CP, ch-0.30, pf_text)

    # Image caption
    cap(sl, px, cy+ch-0.25, pw,
        "Fig. 3. Building parts: exterior walls (blue), interior walls (orange), roof (green), floor (gray)")

    print(f"  L4 Problem Framing -> y={cy+ch:.2f}")

    # ═══════════════════════════════════════════════════════════
    # CENTER PANEL
    # ═══════════════════════════════════════════════════════════
    y = CT
    px = CX; pw = CW

    # ── C1: ENGINEERING METHODOLOGY ──
    y = hdr(sl, px, y, pw, "Engineering Methodology")
    cy = y; ch = 16.30
    card(sl, px, cy, pw, ch)

    # Overview text (like reference has "Overview & Summary" at top)
    ov_text = [
        {"text": "Overview & Summary of diagrams/figures below:", "sz": Pt(10), "bold": True, "color": TXT_DARK, "sb": 0, "sa": 1},
        {"text": "My research methodologies were split up into four main parts: data generation, surrogate training, sensitivity optimization, and validation.", "sz": Pt(9), "color": TXT_DARK, "sb": 0, "sa": 1, "ls": 11},
        {"text": "(1) First, 14,293 house wireframes from 3DWire are converted into volumetric structures via extrusion and boolean fusion.", "sz": Pt(9), "color": TXT_DARK, "sb": 0, "sa": 0},
        {"text": "(2) Then, each mesh is solved with FEA (SfePy) under ASCE 7-22 load combinations extracting stress, displacement, compliance.", "sz": Pt(9), "color": TXT_DARK, "sb": 0, "sa": 0},
        {"text": "(3) Geometries are voxelized to 128-cube grids. 11,178 valid samples are split family-aware into train/val/test.", "sz": Pt(9), "color": TXT_DARK, "sb": 0, "sa": 0},
        {"text": "(4) Finally, the SASTO algorithm uses the trained ensemble for sensitivity-guided erosion with adaptive batch sizing.", "sz": Pt(9), "color": TXT_DARK, "sb": 0, "sa": 0},
    ]
    mtb(sl, px+CP, cy+0.08, pw*0.55, 2.20, ov_text)

    # Step labels (like reference Step 1, Step 2, Step 3, Step 4)
    steps = [("Wireframe", "Step 1"), ("Volumetric", "Step 2"), ("FEA + Train", "Step 3"), ("Optimize", "Step 4")]
    step_x = px + pw*0.56
    step_w = (pw*0.44 - CP) / 4 - 0.05
    for i, (label, step) in enumerate(steps):
        sx = step_x + i*(step_w + 0.05)
        rect(sl, sx, cy+0.10, step_w, 0.80, fill=TEAL, r=15000)
        tb(sl, sx, cy+0.10, step_w, 0.40, label, sz=Pt(8), bold=True, color=TXT_WHITE, align=PP_ALIGN.CENTER)
        tb(sl, sx, cy+0.50, step_w, 0.30, step, sz=Pt(8), italic=True, color=TXT_WHITE, align=PP_ALIGN.CENTER)
        # Arrow between steps
        if i < 3:
            tb(sl, sx+step_w-0.02, cy+0.30, 0.15, 0.30, "\u25B6", sz=Pt(10), color=GOLD, align=PP_ALIGN.CENTER)

    # ── DATA GENERATION ── (top-left quadrant)
    sub_y = cy + 2.35
    half_w = (pw - 2*CP - 0.15) / 2

    # Dataset pipeline images
    tb(sl, px+CP, sub_y, half_w, 0.22, "Data Generation Pipeline", sz=Pt(11), bold=True, color=TXT_DARK)
    sub_y += 0.25

    # 4 images in a row (wireframe, volumetric, FEA, voxel)
    img_w = (half_w - 0.30) / 4
    for i, n in enumerate([4, 6, 5, 3]):  # pipeline, dataset, thickness, structural parts
        img = img_ext(n)
        if img: img_safe(sl, img, px+CP + i*(img_w+0.06), sub_y, img_w, 2.00)

    # Dataset stats text
    ds_text = [
        {"text": "14,293 wireframes -> 11,178 valid FEA simulations", "sz": Pt(8), "color": TXT_DARK, "sb": 0, "sa": 0},
        {"text": "Train: 8,943 | Val: 1,121 | Test: 1,114", "sz": Pt(8), "bold": True, "color": TXT_DARK, "sb": 0, "sa": 0},
        {"text": "Targets: peak VM stress, max displacement, compliance", "sz": Pt(8), "color": TXT_DARK, "sb": 0, "sa": 0},
    ]
    mtb(sl, px+CP, sub_y+2.05, half_w, 0.65, ds_text)

    # ── DEEP ENSEMBLE ── (top-right quadrant)
    ens_x = px + CP + half_w + 0.15
    ens_y = cy + 2.35
    tb(sl, ens_x, ens_y, half_w, 0.22, "Deep Ensemble Surrogate (5 x 8.76M params)", sz=Pt(11), bold=True, color=TXT_DARK)
    ens_y += 0.25

    # Architecture image
    arch_img = img_ext(7) or find_img("fig2_architecture")
    img_safe(sl, arch_img, ens_x, ens_y, half_w, 2.80)

    # Architecture text
    arch_text = [
        {"text": "Starting Model Architecture (~19 layers)", "sz": Pt(9), "bold": True, "color": TXT_DARK, "sb": 0, "sa": 1},
        {"text": "7ch x 128-cube input -> 4 conv stages (BN+GELU)", "sz": Pt(8), "color": TXT_DARK, "sb": 0, "sa": 0},
        {"text": "-> 3 SE-ResBlocks -> Dual pool (avg+max) -> 512d", "sz": Pt(8), "color": TXT_DARK, "sb": 0, "sa": 0},
        {"text": "-> concat 128d features -> 640->512->256->3 outputs", "sz": Pt(8), "color": TXT_DARK, "sb": 0, "sa": 0},
        {"text": "Loss: Huber | Opt: AdamW lr=5e-4 | Aug: 90deg rot, flips", "sz": Pt(8), "color": TXT_DARK, "sb": 0, "sa": 0},
    ]
    mtb(sl, ens_x, ens_y+2.85, half_w, 1.10, arch_text)

    # ── SASTO ALGORITHM ── (bottom-left quadrant)
    algo_y = cy + 5.65
    tb(sl, px+CP, algo_y, half_w, 0.22, "SASTO Algorithm (Sensitivity-Guided Erosion)", sz=Pt(11), bold=True, color=TXT_DARK)
    algo_y += 0.25

    # Algorithm image
    algo_img = img_ext(8) or find_img("fig1_pipeline")
    img_safe(sl, algo_img, px+CP, algo_y, half_w, 3.20)

    algo_text = [
        {"text": "Phase 1: Sensitivity-Guided Erosion (>99% of removal)", "sz": Pt(9), "bold": True, "color": TEAL, "sb": 0, "sa": 0},
        {"text": "\u2022 Compute distance transform, identify interior surface voxels with depth > tmin(p)", "sz": Pt(8), "color": TXT_DARK, "sb": 0, "sa": 0},
        {"text": "\u2022 Backpropagate through ensemble: rank voxels by sensitivity si", "sz": Pt(8), "color": TXT_DARK, "sb": 0, "sa": 0},
        {"text": "\u2022 Select batch of 6-simple-point (topology-safe) voxels", "sz": Pt(8), "color": TXT_DARK, "sb": 0, "sa": 0},
        {"text": "\u2022 Tentatively remove -> query ensemble for mu, sigma -> check mu + k*sigma bounds", "sz": Pt(8), "color": TXT_DARK, "sb": 0, "sa": 0},
        {"text": "\u2022 If constraints satisfied: commit. If not: undo + halve batch (trust-region)", "sz": Pt(8), "color": TXT_DARK, "sb": 0, "sa": 0},
        {"text": "Phase 2: Endgame (B=5, then 1) | Phase 3: Swap Moves", "sz": Pt(9), "bold": True, "color": GOLD, "sb": 1, "sa": 0},
        {"text": "Post-process: fill pockets -> SDF -> Marching Cubes -> watertight STL", "sz": Pt(8), "color": TXT_DARK, "sb": 0, "sa": 0},
    ]
    mtb(sl, px+CP, algo_y+3.25, half_w, 2.20, algo_text)

    # ── 6-CONNECTIVITY ── (bottom-right quadrant)
    conn_x = px + CP + half_w + 0.15
    conn_y = cy + 5.65
    tb(sl, conn_x, conn_y, half_w, 0.22, "Topology Preservation: 6-Connectivity", sz=Pt(11), bold=True, color=TXT_DARK)
    conn_y += 0.25

    # 6-connectivity comparison images
    conn_img1 = img_ext(10)  # 6-conn image
    conn_img2 = img_ext(9)   # 26-conn image
    img_w2 = (half_w - 0.15) / 2
    if conn_img2: img_safe(sl, conn_img2, conn_x, conn_y, img_w2, 2.50)
    if conn_img1: img_safe(sl, conn_img1, conn_x + img_w2 + 0.15, conn_y, img_w2, 2.50)

    # Labels
    tb(sl, conn_x, conn_y+2.55, img_w2, 0.20, "26-conn: fragments", sz=Pt(8), bold=True, color=RED, align=PP_ALIGN.CENTER)
    tb(sl, conn_x+img_w2+0.15, conn_y+2.55, img_w2, 0.20, "6-conn: 1 component", sz=Pt(8), bold=True, color=TEAL, align=PP_ALIGN.CENTER)

    conn_text = [
        {"text": "Proposition: A binary voxel field with exactly one 6-connected foreground component yields a single-component marching-cubes surface mesh.", "sz": Pt(9), "italic": True, "color": TXT_DARK, "sb": 2, "sa": 2},
        {"text": "\u2022 6-adjacent voxels share a FACE -> printable path exists", "sz": Pt(8), "color": TEAL, "sb": 0, "sa": 0},
        {"text": "\u2022 26-adjacent voxels sharing only a CORNER -> floating fragment", "sz": Pt(8), "color": RED, "sb": 0, "sa": 0},
        {"text": "\u2022 100% single-component voxel fields (60/60 tested)", "sz": Pt(8), "bold": True, "color": TXT_DARK, "sb": 0, "sa": 0},
        {"text": "\u2022 90% produce single-component meshes directly; 10% trivially post-processed", "sz": Pt(8), "color": TXT_DARK, "sb": 0, "sa": 0},
    ]
    mtb(sl, conn_x, conn_y+2.80, half_w, 1.70, conn_text)

    # Data interpolation section (like reference has "Data Interpolation & Preparation")
    di_y = cy + 11.20
    tb(sl, px+CP, di_y, pw-2*CP, 0.22, "Data Interpolation & Preparation", sz=Pt(11), bold=True, color=TXT_DARK)
    di_y += 0.25

    di_text = [
        {"text": "Total Data Size: ~200GB raw FEA outputs | Filtered to 59GB processed.", "sz": Pt(9), "color": TXT_DARK, "sb": 0, "sa": 1},
        {"text": "\u2022 3,115/14,293 (21.8%) rejected: diverged solver, near-zero strain energy, invalid results", "sz": Pt(8), "color": TXT_DARK, "sb": 0, "sa": 0},
        {"text": "\u2022 Targets span 4.9 orders of magnitude (stress) to 7.7 orders (compliance)", "sz": Pt(8), "color": TXT_DARK, "sb": 0, "sa": 0},
        {"text": "\u2022 Log-transform normalization: log(1+|y|) -> z-score with 2nd/98th percentile winsorization", "sz": Pt(8), "color": TXT_DARK, "sb": 0, "sa": 0},
        {"text": "\u2022 Family-aware splitting prevents near-duplicate leakage across partitions", "sz": Pt(8), "color": TXT_DARK, "sb": 0, "sa": 0},
    ]
    mtb(sl, px+CP, di_y, pw*0.50, 1.20, di_text)

    # Dataset distribution figure
    dist_fig = find_img("fig_distributions")
    img_safe(sl, dist_fig, px+pw*0.52, di_y-0.10, pw*0.46, 1.40)

    # Model Tuning & Optimization section
    mt_y = cy + 12.80
    tb(sl, px+CP, mt_y, pw-2*CP, 0.22, "Model Tuning & Optimization", sz=Pt(11), bold=True, color=TXT_DARK)
    mt_y += 0.25

    # Training curves + activation function side by side
    tc_fig = find_img("fig_training_curves")
    ac_fig = find_img("fig_activation")
    img_safe(sl, tc_fig, px+CP, mt_y, pw*0.45, 1.80)
    img_safe(sl, ac_fig, px+CP+pw*0.47, mt_y, pw*0.22, 1.30)

    mt_text = [
        {"text": "Hyperparameters:", "sz": Pt(9), "bold": True, "color": TXT_DARK, "sb": 0, "sa": 0},
        {"text": "\u2022 AdamW optimizer, lr=5e-4, cosine annealing", "sz": Pt(8), "color": TXT_DARK, "sb": 0, "sa": 0},
        {"text": "\u2022 EMA decay=0.999, gradient clipping ||.||<1.0", "sz": Pt(8), "color": TXT_DARK, "sb": 0, "sa": 0},
        {"text": "\u2022 Max 200 epochs, early stopping patience=30", "sz": Pt(8), "color": TXT_DARK, "sb": 0, "sa": 0},
        {"text": "\u2022 I had a breakthrough on model version 14", "sz": Pt(8), "color": TXT_DARK, "sb": 0, "sa": 0},
        {"text": "which began to have high agreement with FEA.", "sz": Pt(8), "color": TXT_DARK, "sb": 0, "sa": 0},
    ]
    mtb(sl, px+CP+pw*0.70, mt_y, pw*0.28, 1.80, mt_text)

    # Equation at bottom of methodology
    eq_y = cy + 15.10
    eq_pill(sl, px+CP, eq_y, pw-2*CP,
            "Conservative bound: sigma_hat = mu_sigma + k * sigma_sigma,  C_hat = mu_C + k * sigma_C,  k = 1.0",
            h=0.32)

    # Figure labels
    cap(sl, px+CP, eq_y+0.35, pw*0.45, "Fig. 5. Training loss convergence (5 ensemble members)")
    cap(sl, px+CP+pw*0.47, eq_y+0.35, pw*0.22, "GELU activation")

    # Model version caption
    cap(sl, px+CP, cy+ch-0.25, pw-2*CP,
        "Model version 14, training progress in PyTorch. 4x NVIDIA GB200 GPUs.")

    y_c1_end = cy + ch + SG
    print(f"  C1 Eng. Methodology -> y={y_c1_end:.2f}")

    # ── C2: RESULTS & IN-SILICO VALIDATION ──
    y = y_c1_end
    y = hdr(sl, CX, y, CW, "Results & In-Silico Validation")
    cy = y; ch = CB - y   # fill to bottom
    card(sl, CX, cy, CW, ch)

    # Three-column layout like reference
    col_w = (CW - 2*CP - 0.20) / 3
    col_gap = 0.10

    # ── COL A: Reference Case ──
    ax = CX + CP
    ay = cy + 0.10

    # Figures: RRMSE and RMSE equivalent -> histogram + compliance
    fig1 = find_img("fig_histogram")
    img_safe(sl, fig1, ax, ay, col_w, 2.20)
    cap(sl, ax, ay+2.22, col_w, "Fig. 6. Volume reduction distribution (n=1,114)")

    # Regression plot
    fig2 = find_img("fig_regression")
    img_safe(sl, fig2, ax, ay+2.50, col_w*0.55, 2.00)

    # Bland-Altman
    fig3 = find_img("fig_bland_altman")
    img_safe(sl, fig3, ax+col_w*0.57, ay+2.50, col_w*0.43, 2.00)

    cap(sl, ax, ay+4.55, col_w*0.55, "Fig. 7. Regression Plot")
    cap(sl, ax+col_w*0.57, ay+4.55, col_w*0.43, "Fig. 8. Bland-Altman")

    # Reference case results text
    ref_text = [
        {"text": "Reference Case (Sample 00472)", "sz": Pt(10), "bold": True, "color": TXT_DARK, "sb": 2, "sa": 2},
        {"text": "\u2022 Volume: 116,872 -> 64,292 voxels (-45.0%)", "sz": Pt(8), "color": TXT_DARK, "sb": 0, "sa": 0},
        {"text": "\u2022 VM stress conservative: 3.08 x 10^6 Pa", "sz": Pt(8), "color": TXT_DARK, "sb": 0, "sa": 0},
        {"text": "\u2022 Compliance ratio (FEA): 1.004 (limit: 1.15)", "sz": Pt(8), "color": TXT_DARK, "sb": 0, "sa": 0},
        {"text": "\u2022 Mesh components: 1 (watertight)", "sz": Pt(8), "color": TXT_DARK, "sb": 0, "sa": 0},
        {"text": "\u2022 Runtime: 159.5 seconds", "sz": Pt(8), "color": TXT_DARK, "sb": 0, "sa": 0},
        {"text": "\u2022 SASTO-PA vs SASTO-U: +10.7pp improvement", "sz": Pt(8), "bold": True, "color": TEAL, "sb": 0, "sa": 0},
    ]
    mtb(sl, ax, ay+4.80, col_w, 2.00, ref_text)

    # ── COL B: Qualitative Check & Learning Rate ──
    bx = ax + col_w + col_gap
    by = cy + 0.10

    # Speedup figure
    fig4 = find_img("fig_speedup")
    img_safe(sl, fig4, bx, by, col_w, 1.80)
    cap(sl, bx, by+1.82, col_w, "Fig. 9. Speedup comparison (log scale)")

    # Per-part retention
    fig5 = find_img("fig_per_part")
    img_safe(sl, fig5, bx, by+2.10, col_w, 1.80)
    cap(sl, bx, by+3.95, col_w, "Fig. 10. Per-part material retention")

    # Convergence figure
    fig6 = find_img("fig_convergence")
    img_safe(sl, fig6, bx, by+4.20, col_w, 2.30)
    cap(sl, bx, by+6.52, col_w, "Fig. 11. Optimization convergence (reference case)")

    # ── COL C: FEA Validation & Activation ──
    cx = bx + col_w + col_gap
    c_y = cy + 0.10

    # FEA compliance validation scatter
    fig7 = find_img("fig_fea_compliance")
    img_safe(sl, fig7, cx, c_y, col_w, 2.30)
    cap(sl, cx, c_y+2.32, col_w, "Fig. 12. Independent FEA compliance validation (n=1,114)")

    # Before/after 3D renders
    before_img = img_ext(0) or img_ext(2)
    after_img = img_ext(11) or img_ext(9)
    half_col = (col_w - 0.10) / 2
    img_safe(sl, before_img, cx, c_y+2.60, half_col, 1.80)
    img_safe(sl, after_img, cx+half_col+0.10, c_y+2.60, half_col, 1.80)
    cap(sl, cx, c_y+4.42, half_col, "Original")
    cap(sl, cx+half_col+0.10, c_y+4.42, half_col, "SASTO-PA Optimized")

    # Timing table
    timing_fig = find_img("fig_timing_table")
    img_safe(sl, timing_fig, cx, c_y+4.70, col_w, 1.10)

    # FEA note
    fea_note = [
        {"text": "Fig. 13. Patient Anatomy -> House Geometry 3D", "sz": Pt(8), "italic": True, "color": TXT_DARK, "sb": 0, "sa": 1},
        {"text": "Independent hex8 voxel FEA on all 1,114 designs:", "sz": Pt(8), "bold": True, "color": TXT_DARK, "sb": 0, "sa": 0},
        {"text": "\u2022 Mean C_opt/C_base = 0.631 +/- 0.112", "sz": Pt(8), "color": TXT_DARK, "sb": 0, "sa": 0},
        {"text": "\u2022 Max ratio = 1.004 (limit 1.15)", "sz": Pt(8), "color": TXT_DARK, "sb": 0, "sa": 0},
        {"text": "\u2022 0/1,114 violations confirmed", "sz": Pt(8), "bold": True, "color": TEAL, "sb": 0, "sa": 0},
    ]
    mtb(sl, cx, c_y+5.85, col_w, 1.20, fea_note)

    # Bottom equation spanning full width
    eq_y2 = cy + ch - 0.70
    eq_pill(sl, CX+CP, eq_y2, CW-2*CP,
            "Conservative bound: sigma_hat+ = mu + k*sigma (k=1.0). "
            "All squares shown are in units of c, where c, w, and g are FEA quantities. "
            "P(violation) <= 1/(n+1) = 0.09%  (conformal, distribution-free)",
            h=0.35)
    cap(sl, CX+CP, eq_y2+0.38, CW-2*CP,
        "Model version 14, 5-member ensemble. Independent FEA validation on all 1,114 designs.")

    print(f"  C2 Results & Validation -> y={cy+ch:.2f}")

    # ═══════════════════════════════════════════════════════════
    # RIGHT PANEL
    # ═══════════════════════════════════════════════════════════
    y = CT
    px = RX; pw = RW

    # ── R1: STATISTICAL ANALYSIS ──
    y = hdr(sl, px, y, pw, "Statistical Analysis")
    cy = y; ch = 13.20
    card(sl, px, cy, pw, ch)

    stat_intro = [
        {"text": "I conducted an in-depth statistical data analysis to quantify and verify the performance of SASTO. A broad, robust test set of 1,114 held-out house geometries was used.", "sz": Pt(9), "color": TXT_DARK, "sb": 0, "sa": 3, "ls": 11},
    ]
    mtb(sl, px+CP, cy+0.10, pw-2*CP, 0.60, stat_intro)

    # k-factor figure + uncertainty figure
    kf_fig = find_img("fig_k_factor")
    img_safe(sl, kf_fig, px+0.15, cy+0.70, pw-0.30, 2.60)
    cap(sl, px, cy+3.32, pw, "Fig. 14. k-factor sensitivity (dual-axis Pareto frontier)")

    # Uncertainty bands figure
    unc_fig = find_img("fig_uncertainty")
    img_safe(sl, unc_fig, px+0.15, cy+3.55, pw-0.30, 2.20)
    cap(sl, px, cy+5.78, pw, "Fig. 15. Ensemble uncertainty bands during optimization")

    # Surrogate metrics table
    smt, _ = tbl(sl, px+CP, cy+6.00, pw-2*CP, [
        ["Target", "Spearman rho", "R2_log", "MAPE (%)"],
        ["Von Mises stress", "0.737", "0.419", "37.4"],
        ["Displacement", "0.970", "0.842", "10.9"],
        ["Compliance", "0.948", "0.814", "18.5"],
    ], [0.30, 0.25, 0.22, 0.23], rh=0.22)

    # Text analysis
    stat_text = [
        {"text": "Surrogate requires ranking accuracy, not pointwise accuracy \u2014 compliance Spearman rho = 0.948 drives optimization safety.", "sz": Pt(8), "italic": True, "color": TXT_DARK, "sb": 2, "sa": 3},
        {"text": "Conformal Prediction Results:", "sz": Pt(10), "bold": True, "color": TXT_DARK, "sb": 2, "sa": 2},
        {"text": "\u2022 Conformal k for 84.1% compliance coverage: k = 1.90 (heavier tails than Gaussian)", "sz": Pt(8), "color": TXT_DARK, "sb": 0, "sa": 0},
        {"text": "\u2022 Conformal k for 84.1% VM stress coverage: k = 4.31 (localized, hard to predict)", "sz": Pt(8), "color": TXT_DARK, "sb": 0, "sa": 0},
        {"text": "\u2022 99% conformal upper bound on C_opt/C_base: 0.950 (margin 0.20 to limit 1.15)", "sz": Pt(8), "color": TXT_DARK, "sb": 0, "sa": 0},
        {"text": "\u2022 P(violation) <= 1/(n+1) = 0.09% (distribution-free, n=1,114)", "sz": Pt(8), "bold": True, "color": TEAL, "sb": 0, "sa": 2},
        {"text": "Endpoints & Specifications (tested on RTX A3000 Laptop GPU):", "sz": Pt(10), "bold": True, "color": TXT_DARK, "sb": 2, "sa": 2},
        {"text": "SASTO runtime is four orders of magnitude faster than the widespread SIMP topology optimization approach. Based on my testing, SASTO enables optimization on a consumer GPU machine, and its simple UI allows anyone to access the technology easily.", "sz": Pt(8), "color": TXT_DARK, "sb": 0, "sa": 0, "ls": 10},
    ]
    mtb(sl, px+CP, cy+7.00, pw-2*CP, 4.20, stat_text)

    # Timing comparison table
    tbl(sl, px+CP, cy+11.30, pw-2*CP, [
        ["Approach", "Per Design", "Speedup"],
        ["SASTO (ours)", "50 s median", "1x (baseline)"],
        ["SIMP (64-cube)", "94 s median", "~0.5x"],
        ["SIMP (128-cube proj.)", "19-77 min", "23-92x slower"],
    ], [0.35, 0.35, 0.30], rh=0.22)

    y = cy + ch + SG
    print(f"  R1 Statistical Analysis -> y={y:.2f}")

    # ── R2: CONCLUSIONS ──
    y = hdr(sl, px, y, pw, "Conclusions")
    cy = y; ch = 5.80
    card(sl, px, cy, pw, ch)

    concl_text = [
        {"text": "This research demonstrates that the deep learning surrogate can optimize building-scale structures just as accurately as FEA. SASTO solves the key problems by (1) having clinically viable runtimes, (2) being able to run on any initial geometry, and (3) being personalized to each building's structural characteristics.", "sz": Pt(9), "color": TXT_DARK, "sb": 0, "sa": 3, "ls": 11},
        {"text": "Revisiting Engineering Criteria:", "sz": Pt(10), "bold": True, "color": TXT_DARK, "sb": 1, "sa": 2},
        {"text": "\u2022 SASTO exceeded the claim and met all the engineering endpoints by being able to optimize any house geometry with any initial thickness.", "sz": Pt(8), "color": TXT_DARK, "sb": 0, "sa": 1},
        {"text": "\u2022 SASTO has near-perfect agreement with the ground truth, real FEA solutions, while running the equivalent of 1,000 FEA simulations for a design in just ~50 seconds (vs. 19-77 min) via neural approximation.", "sz": Pt(8), "color": TXT_DARK, "sb": 0, "sa": 1},
        {"text": "Bench Impacts: This study has introduced several novel methods and results.", "sz": Pt(9), "bold": True, "color": TXT_DARK, "sb": 2, "sa": 1},
        {"text": "\u2022 Serves as a proof-of-concept for surrogate-based building optimization.", "sz": Pt(8), "color": TXT_DARK, "sb": 0, "sa": 0},
        {"text": "\u2022 Architecture can be applied to other structures, sub-buildings, and even other domains.", "sz": Pt(8), "color": TXT_DARK, "sb": 0, "sa": 0},
        {"text": "\u2022 6-connectivity criterion is a new contribution to voxel-based topology optimization.", "sz": Pt(8), "color": TXT_DARK, "sb": 0, "sa": 0},
        {"text": "Bedside Impacts: SASTO is a decision support tool for structural engineers.", "sz": Pt(9), "bold": True, "color": TXT_DARK, "sb": 2, "sa": 1},
        {"text": "\u2022 Reduces turnaround time for precise optimization from days to minutes.", "sz": Pt(8), "color": TXT_DARK, "sb": 0, "sa": 0},
        {"text": "\u2022 8% of global CO2 = cement. 23.5% less concrete per house -> millions of tons saved at scale.", "sz": Pt(8), "bold": True, "color": TEAL, "sb": 0, "sa": 0},
    ]
    mtb(sl, px+CP, cy+0.10, pw-2*CP, ch-0.20, concl_text)

    y = cy + ch + SG
    print(f"  R2 Conclusions -> y={y:.2f}")

    # ── R3: FUTURE WORK ──
    y = hdr(sl, px, y, pw, "Future Work")
    cy = y; ch = 3.80
    card(sl, px, cy, pw, ch)

    fw_text = [
        {"text": "I am very excited about the results I have achieved thus far. I have lots of ideas for future investigation and studies I want to conduct next:", "sz": Pt(9), "color": TXT_DARK, "sb": 0, "sa": 2, "ls": 11},
        {"text": "\u2022 Apply FEA-in-the-loop active learning: When ensemble disagreement exceeds threshold, trigger ground-truth FEA mid-optimization, creating a self-correcting safety net.", "sz": Pt(8), "color": TXT_DARK, "sb": 0, "sa": 1},
        {"text": "\u2022 Run nonlinear FEA spot checks (concrete damaged plasticity) to assess whether 78mm interior partitions exhibit cracking or buckling.", "sz": Pt(8), "color": TXT_DARK, "sb": 0, "sa": 1},
        {"text": "\u2022 Physical print validation: Fabricate one optimized house at 1:10 scale. Compression testing + DIC strain measurement.", "sz": Pt(8), "color": TXT_DARK, "sb": 0, "sa": 1},
        {"text": "\u2022 Adapt to other forward solving problems in fluid dynamics and natural science disciplines, i.e., ocean current prediction, cell motility.", "sz": Pt(8), "color": TXT_DARK, "sb": 0, "sa": 1},
        {"text": "Not possible yet due to data availability:", "sz": Pt(9), "bold": True, "italic": True, "color": TXT_DARK, "sb": 1, "sa": 1},
        {"text": "\u2022 Add multi-story structures, seismic loading, and different materials (geopolymer, fiber-reinforced concrete) as parameters in the model.", "sz": Pt(8), "color": TXT_DARK, "sb": 0, "sa": 0},
    ]
    mtb(sl, px+CP, cy+0.10, pw-2*CP, ch-0.20, fw_text)

    y = cy + ch + SG
    print(f"  R3 Future Work -> y={y:.2f}")

    # ── R4: KEY REFERENCES ──
    y = hdr(sl, px, y, pw, "Key References")
    cy = y; ch = CB - y  # fill to bottom
    card(sl, px, cy, pw, ch)

    refs_text = [
        {"text": "M. P. Bendsoe and O. Sigmund. Topology Optimization: Theory, Methods, Applications. Springer, 2003.", "sz": Pt(8), "color": TXT_DARK, "sb": 0, "sa": 1, "ls": 10},
        {"text": "B. Lakshminarayanan, A. Pritzel, and C. Blundell. Simple and scalable predictive uncertainty estimation using deep ensembles. NeurIPS, 2017.", "sz": Pt(8), "color": TXT_DARK, "sb": 0, "sa": 1, "ls": 10},
        {"text": "T. Y. Kong and A. Rosenfeld. Digital topology: Introduction and survey. CVGIP, 48(3):357-393, 1989.", "sz": Pt(8), "color": TXT_DARK, "sb": 0, "sa": 1, "ls": 10},
        {"text": "ASCE. ASCE/SEI 7-22: Minimum Design Loads and Associated Criteria. American Society of Civil Engineers, 2022.", "sz": Pt(8), "color": TXT_DARK, "sb": 0, "sa": 1, "ls": 10},
        {"text": "R. A. Buswell et al. 3D printing using concrete extrusion: A roadmap for research. Cem. Concr. Res. 112:37-49, 2018.", "sz": Pt(8), "color": TXT_DARK, "sb": 0, "sa": 1, "ls": 10},
        {"text": "O. Sigmund and K. Maute. Topology optimization approaches. Struct. Multidisc. Optim. 48:1031-1056, 2013.", "sz": Pt(8), "color": TXT_DARK, "sb": 0, "sa": 1, "ls": 10},
        {"text": "Y. Lin et al. 3DWire: 3D Building Wireframe Dataset. KAUST Visual Computing Center, 2024.", "sz": Pt(8), "color": TXT_DARK, "sb": 0, "sa": 1, "ls": 10},
        {"text": "IEA. Global Status Report for Buildings and Construction 2021. Int'l Energy Agency.", "sz": Pt(8), "color": TXT_DARK, "sb": 0, "sa": 1, "ls": 10},
        {"text": "W. E. Lorensen and H. E. Cline. Marching Cubes: A high resolution 3D surface construction algorithm. ACM SIGGRAPH, 21(4):163-169, 1987.", "sz": Pt(8), "color": TXT_DARK, "sb": 0, "sa": 1, "ls": 10},
        {"text": "V. Vovk, A. Gammerman, and G. Shafer. Algorithmic Learning in a Random World. Springer, 2005.", "sz": Pt(8), "color": TXT_DARK, "sb": 0, "sa": 1, "ls": 10},
    ]
    mtb(sl, px+CP, cy+0.10, pw-2*CP, ch-0.25, refs_text)

    # Italic note at bottom
    tb(sl, px+CP, cy+ch-0.25, pw-2*CP, 0.22,
       "Paper bibliography contains full list.",
       sz=Pt(8), italic=True, color=RED, align=PP_ALIGN.CENTER)

    print(f"  R4 Key References -> y={cy+ch:.2f}")

    # ═══════════════════════════════════════════════════════════
    # SAVE
    # ═══════════════════════════════════════════════════════════
    prs.save(OUTPUT)
    n = sum(1 for _ in sl.shapes)
    print(f"\n\u2713 Saved: {OUTPUT}")
    print(f"  {BW}x{BH} in, ~{n} shapes")


if __name__ == "__main__":
    build()
