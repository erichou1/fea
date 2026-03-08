#!/usr/bin/env python3
"""
ISEF Poster PPTX v6 — Updated production build.
- Removes "Engineering Design Criteria" section (4 sections per side panel)
- Folds key equations into Problem Framing
- Uses improved figures with actual 3D renders
- Visual Abstract uses reference-style two-column layout

Usage:
    cd poster_final && python3 build_poster_v6.py

Output:  SASTO_ISEF_Poster_v6.pptx (48×36 in)
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
import os

# ═══════════════════════════════════════════════════════════════
# COLOR PALETTE
# ═══════════════════════════════════════════════════════════════
BG_NAVY   = RGBColor(0x06, 0x2B, 0x7A)
TITLE_BAND= RGBColor(0x03, 0x20, 0x61)
SEC_BAR   = RGBColor(0x0A, 0x3D, 0x9A)
CARD_FILL = RGBColor(0xF7, 0xF9, 0xFC)
CARD_BDR  = RGBColor(0xB7, 0xC5, 0xE3)
TXT_DARK  = RGBColor(0x0B, 0x17, 0x36)
TXT_WHITE = RGBColor(0xFF, 0xFF, 0xFF)
RED       = RGBColor(0xD7, 0x26, 0x3D)
TEAL      = RGBColor(0x00, 0x8C, 0x9E)
GOLD      = RGBColor(0xCF, 0xA5, 0x35)
EQ_BG     = RGBColor(0xE8, 0xEE, 0xF8)
ALT_ROW   = RGBColor(0xEE, 0xF2, 0xFA)

# ═══════════════════════════════════════════════════════════════
# DIMENSIONS — 48 × 36 tri-fold (12 | 24 | 12)
# ═══════════════════════════════════════════════════════════════
BW = 48.0;  BH = 36.0
TITLE_H = 3.25
CT = TITLE_H + 0.10   # content top
CB = 35.82             # content bottom
M  = 0.18              # margin
SG = 0.04              # section gap
HH = 0.45              # header height
CP = 0.10              # card padding

LX = M;            LW = 12.0 - 2*M    # left panel
CX = 12.0 + M;     CW = 24.0 - 2*M    # center panel
RX = 36.0 + M;     RW = 12.0 - 2*M    # right panel

CARD_R = 40000  # corner radius EMU

BASE = os.path.dirname(os.path.abspath(__file__))
PARENT = os.path.dirname(BASE)
EXTRACTED = os.path.join(PARENT, "poster_images_extracted")
RENDERS = os.path.join(BASE, "renders_hq")
FIG_DIR = BASE
OUTPUT = os.path.join(BASE, "SASTO_ISEF_Poster_v6.pptx")


def fig(name):
    p = os.path.join(FIG_DIR, name + ".png")
    return p if os.path.isfile(p) else None


def render(name):
    p = os.path.join(RENDERS, name + ".png")
    return p if os.path.isfile(p) else None


def extracted_img(name):
    for ext in (".jpg", ".png"):
        p = os.path.join(EXTRACTED, name + ext)
        if os.path.isfile(p):
            return p
    return None


# ═══════════════════════════════════════════════════════════════
# HELPERS
# ═══════════════════════════════════════════════════════════════

def _set_radius(shape, r=CARD_R):
    spPr = shape._element.spPr
    pg = spPr.find(qn("a:prstGeom"))
    if pg is not None:
        pg.set("prst", "roundRect")
        av = pg.find(qn("a:avLst"))
        if av is None:
            av = spPr.makeelement(qn("a:avLst"), {})
            pg.append(av)
        for c in list(av):
            av.remove(c)
        gd = spPr.makeelement(qn("a:gd"), {"name": "adj", "fmla": f"val {r}"})
        av.append(gd)


def rect(sl, x, y, w, h, fill=None, border=None, bw=Pt(1), r=CARD_R):
    s = sl.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                            Inches(x), Inches(y), Inches(w), Inches(h))
    if fill:
        s.fill.solid(); s.fill.fore_color.rgb = fill
    else:
        s.fill.background()
    if border:
        s.line.color.rgb = border; s.line.width = bw
    else:
        s.line.fill.background()
    _set_radius(s, r)
    return s


def card(sl, x, y, w, h):
    return rect(sl, x, y, w, h, fill=CARD_FILL, border=CARD_BDR)


def tb(sl, x, y, w, h, text, font="Arial", sz=Pt(12), bold=False, italic=False,
       color=TXT_DARK, align=PP_ALIGN.LEFT):
    t = sl.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = t.text_frame; tf.word_wrap = True; tf.auto_size = None
    p = tf.paragraphs[0]; p.text = text
    p.font.name = font; p.font.size = sz; p.font.bold = bold; p.font.italic = italic
    if color:
        p.font.color.rgb = color
    p.alignment = align
    p.space_before = Pt(0); p.space_after = Pt(0)
    return t


def mtb(sl, x, y, w, h, paras):
    t = sl.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = t.text_frame; tf.word_wrap = True; tf.auto_size = None
    for i, pd in enumerate(paras):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = pd.get("text", "")
        p.font.name = pd.get("font", "Arial")
        p.font.size = pd.get("sz", Pt(12))
        p.font.bold = pd.get("bold", False)
        p.font.italic = pd.get("italic", False)
        if pd.get("color"):
            p.font.color.rgb = pd["color"]
        p.alignment = pd.get("align", PP_ALIGN.LEFT)
        p.space_before = Pt(pd.get("sb", 1))
        p.space_after = Pt(pd.get("sa", 1))
    return t


def hdr(sl, x, y, w, text):
    rect(sl, x, y, w, HH, fill=SEC_BAR, r=15000)
    tb(sl, x, y, w, HH, text, font="Arial Black", sz=Pt(28), bold=True,
       color=TXT_WHITE, align=PP_ALIGN.CENTER)
    return y + HH


def img_safe(sl, f, x, y, w, h):
    if not f or not os.path.isfile(f):
        return None
    try:
        from PIL import Image as PILImage
        with PILImage.open(f) as im:
            iw, ih = im.size
    except ImportError:
        return sl.shapes.add_picture(f, Inches(x), Inches(y), Inches(w), Inches(h))
    ar = iw / ih; bar = w / h
    if ar > bar:
        fw = w; fh = w / ar; fx = x; fy = y + (h - fh) / 2
    else:
        fh = h; fw = h * ar; fx = x + (w - fw) / 2; fy = y
    return sl.shapes.add_picture(f, Inches(fx), Inches(fy), Inches(fw), Inches(fh))


def tbl(sl, x, y, w, rows, col_pcts):
    nr = len(rows); nc = len(rows[0])
    rh = 0.22; th = nr * rh
    ts = sl.shapes.add_table(nr, nc, Inches(x), Inches(y), Inches(w), Inches(th))
    t = ts.table
    for ci, pct in enumerate(col_pcts):
        t.columns[ci].width = Inches(w * pct)
    for ri, row in enumerate(rows):
        for ci, txt in enumerate(row):
            c = t.cell(ri, ci); c.text = txt
            for p in c.text_frame.paragraphs:
                p.font.name = "Arial"; p.font.size = Pt(10)
                p.alignment = PP_ALIGN.LEFT
                if ri == 0:
                    p.font.bold = True; p.font.color.rgb = TXT_WHITE; p.font.size = Pt(10)
                else:
                    p.font.color.rgb = TXT_DARK
            if ri == 0:
                c.fill.solid(); c.fill.fore_color.rgb = SEC_BAR
            elif ri % 2 == 0:
                c.fill.solid(); c.fill.fore_color.rgb = ALT_ROW
            else:
                c.fill.solid(); c.fill.fore_color.rgb = CARD_FILL
    return ts, y + th


def eq_pill(sl, x, y, w, text, h=0.38):
    rect(sl, x, y, w, h, fill=EQ_BG, r=25000)
    tb(sl, x + 0.06, y + 0.03, w - 0.12, h - 0.06, text, sz=Pt(10), align=PP_ALIGN.CENTER)
    return y + h


def badge(sl, x, y, w, h, text, color=TEAL):
    rect(sl, x, y, w, h, fill=color, r=30000)
    tb(sl, x, y, w, h, text, sz=Pt(9), bold=True, color=TXT_WHITE, align=PP_ALIGN.CENTER)


def cap(sl, x, y, w, text):
    tb(sl, x, y, w, 0.20, text, sz=Pt(8), italic=True, color=TXT_DARK, align=PP_ALIGN.CENTER)
    return y + 0.20


# ═══════════════════════════════════════════════════════════════
# BUILD
# ═══════════════════════════════════════════════════════════════

def build():
    prs = Presentation()
    prs.slide_width = Inches(BW)
    prs.slide_height = Inches(BH)
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    print("Building SASTO ISEF Poster v6...")

    # ── BACKGROUND ──
    bg = sl.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(BW), Inches(BH))
    bg.fill.solid(); bg.fill.fore_color.rgb = BG_NAVY; bg.line.fill.background()

    # ── TITLE BAND ──
    tb_bg = sl.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(BW), Inches(TITLE_H))
    tb_bg.fill.solid(); tb_bg.fill.fore_color.rgb = TITLE_BAND; tb_bg.line.fill.background()
    gl = sl.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(TITLE_H - 0.03), Inches(BW), Inches(0.03))
    gl.fill.solid(); gl.fill.fore_color.rgb = GOLD; gl.line.fill.background()

    # Left: house renders — use HQ renders extracted from render_figures.py composites
    r_orig = render("original_solid")
    r_opt = render("sasto_pa_solid")
    if r_orig:
        img_safe(sl, r_orig, 0.40, 0.25, 3.20, 2.60)
    else:
        h0 = extracted_img("Image 0")
        if h0:
            img_safe(sl, h0, 0.50, 0.30, 2.80, 2.40)

    if r_opt:
        img_safe(sl, r_opt, 4.00, 0.25, 3.20, 2.60)
    else:
        h11 = extracted_img("Image 11")
        if h11:
            img_safe(sl, h11, 4.00, 0.30, 3.00, 2.40)

    # Title text
    tb(sl, 8.50, 0.10, 30.00, 0.90,
       "SURROGATE-ACCELERATED STRUCTURAL OPTIMIZATION",
       font="Arial Black", sz=Pt(44), bold=True, color=TXT_WHITE, align=PP_ALIGN.CENTER)
    tb(sl, 8.50, 1.10, 30.00, 0.65,
       "Additive Manufacturing: Harnessing FEA to Optimize Material Efficiency",
       font="Arial", sz=Pt(26), bold=True, italic=True, color=TXT_WHITE, align=PP_ALIGN.CENTER)
    tb(sl, 8.50, 1.85, 30.00, 0.45,
       "Eric Hou",
       font="Arial", sz=Pt(22), bold=True, color=TXT_WHITE, align=PP_ALIGN.CENTER)

    # Credit block (right)
    mtb(sl, 39.00, 0.15, 8.60, 2.80, [
        {"text": "Intel ISEF 2026 - Engineering Mechanics", "sz": Pt(13), "bold": True, "color": TXT_WHITE},
        {"text": "", "sz": Pt(4)},
        {"text": "Credit Line of Origin", "sz": Pt(12), "bold": True, "color": TXT_WHITE},
        {"text": "References & data in references below.", "sz": Pt(10), "color": TXT_WHITE},
        {"text": "", "sz": Pt(4)},
        {"text": "Images denoted with an asterisk (*) are from", "sz": Pt(10), "color": TXT_WHITE},
        {"text": "public-domain or adapted from publicly available", "sz": Pt(10), "color": TXT_WHITE},
        {"text": "sources. All other figures, tables, and images", "sz": Pt(10), "color": TXT_WHITE},
        {"text": "created by Eric Hou, 2026.", "sz": Pt(10), "color": TXT_WHITE},
    ])
    print("  Title band done")

    # ═══════════════════════════════════════════════════════════
    # LEFT PANEL — 4 sections (no Engineering Design Criteria)
    # ═══════════════════════════════════════════════════════════
    # Available height = CB - CT = 35.82 - 3.35 = 32.47
    # 4 sections + 4 headers + 3 gaps = 4*HH + 3*SG = 1.80 + 0.12 = 1.92
    # Content height available ≈ 30.55 for 4 cards
    y = CT

    # ─── L1: VISUAL ABSTRACT ──────────────────────────────────
    y = hdr(sl, LX, y, LW, "VISUAL ABSTRACT")
    va_h = 5.60  # expanded from 4.00
    card(sl, LX, y, LW, va_h)
    va_fig = fig("fig01_visual_abstract_pipeline")
    if va_fig:
        img_safe(sl, va_fig, LX + CP, y + CP, LW - 2*CP, 4.50)
    cap(sl, LX + CP, y + 4.60, LW - 2*CP,
        "Fig. 1. SASTO pipeline: two-column workflow overview")
    tb(sl, LX + CP, y + 4.82, LW - 2*CP, 0.66,
       "Left: user workflow — input building wireframe, run SASTO optimization (~50s), "
       "output watertight STL with 23.5% less material. Right: model creation — "
       "dataset generation, FEA, deep ensemble training, and conformal calibration.",
       sz=Pt(10), color=TXT_DARK)
    y += va_h + SG
    print("  L1 Visual Abstract")

    # ─── L2: INTRODUCTION ─────────────────────────────────────
    y = hdr(sl, LX, y, LW, "INTRODUCTION")
    intro_h = 7.30  # expanded from 5.80
    card(sl, LX, y, LW, intro_h)

    tw = 6.80; iw = LW - tw - 3*CP
    ix = LX + tw + 2*CP

    intro_paras = [
        {"text": "CONCRETE & CO2 CRISIS", "sz": Pt(12), "bold": True, "color": TEAL, "sa": 1},
        {"text": "Concrete production accounts for ~8% of global CO2 emissions [IEA 2021]. "
                 "In conventional construction, walls are built at uniform thickness - a practice driven "
                 "by formwork constraints, not structural need. Interior partitions bear negligible loads "
                 "compared to exterior shear walls, creating a substantial opportunity for material reduction.",
         "sz": Pt(10), "color": TXT_DARK, "sa": 3},
        {"text": "ADDITIVE MANUFACTURING", "sz": Pt(12), "bold": True, "color": TEAL, "sa": 1},
        {"text": "Large-scale 3D printing (ICON, COBOD, Apis Cor) can realize arbitrary wall profiles "
                 "at no marginal tooling cost. However, exploiting this geometric freedom requires optimized "
                 "3D models that minimize volume, satisfy structural constraints under ASCE 7-22 load cases, "
                 "and produce watertight meshes compatible with printer toolpath generation.",
         "sz": Pt(10), "color": TXT_DARK, "sa": 3},
        {"text": "COMPUTATIONAL BOTTLENECK", "sz": Pt(12), "bold": True, "color": TEAL, "sa": 1},
        {"text": "Classical SIMP requires hundreds-thousands of FEA evaluations, each costing minutes-hours "
                 "at building scale. Voxel-based implementations using 26-connectivity produce disconnected "
                 "floating fragments incompatible with AM.",
         "sz": Pt(10), "color": TXT_DARK, "sa": 3},
        {"text": "SASTO CONTRIBUTION", "sz": Pt(12), "bold": True, "color": RED, "sa": 1},
        {"text": "SASTO replaces FEA with a deep ensemble surrogate, achieving 23-92x speedup. "
                 "A 6-connectivity criterion eliminates floating fragments. Part-aware thickness enables "
                 "differential thinning. Across 1,114 geometries: 23.5% +/- 7.8% mean reduction, zero violations.",
         "sz": Pt(10), "bold": True, "color": TXT_DARK, "sa": 2},
    ]
    mtb(sl, LX + CP, y + CP, tw, intro_h - 2*CP, intro_paras)

    # Uniform vs optimized figure on right
    f02 = fig("fig02_uniform_vs_optimized")
    if f02:
        img_safe(sl, f02, ix, y + CP, iw, 6.30)
    cap(sl, ix, y + CP + 6.40, iw, "Uniform vs. SASTO*")

    y += intro_h + SG
    print("  L2 Introduction")

    # ─── L3: RESEARCH OBJECTIVES ──────────────────────────────
    y = hdr(sl, LX, y, LW, "RESEARCH OBJECTIVES")
    obj_h = 3.00  # expanded from 2.50
    card(sl, LX, y, LW, obj_h)

    objs = [
        ("1", "SPEED", "23-92x speedup via deep ensemble surrogate replacing FEA during optimization", TEAL),
        ("2", "PRINTABILITY", "6-connected topology -> single-mesh watertight STL for AM toolpath", TEAL),
        ("3", "EFFICIENCY", "23.5% +/- 7.8% mean material reduction (up to 45.0%) across 1,114 geometries", GOLD),
        ("4", "SAFETY", "0/1,114 FEA violations; conformal P(violation) <= 0.09%", RED),
    ]
    oy = y + CP
    for num, title, desc, bc in objs:
        rect(sl, LX + CP, oy, 0.26, 0.26, fill=bc, r=50000)
        tb(sl, LX + CP, oy, 0.26, 0.26, num, sz=Pt(11), bold=True, color=TXT_WHITE, align=PP_ALIGN.CENTER)
        tb(sl, LX + CP + 0.34, oy - 0.01, 1.40, 0.24, title, sz=Pt(11), bold=True, color=bc)
        tb(sl, LX + CP + 1.76, oy - 0.01, LW - 2*CP - 1.76, 0.24, desc, sz=Pt(10), color=TXT_DARK)
        oy += 0.36

    # Pipeline boxes
    pipe_y = oy + 0.12
    plabels = ["Data\nGeneration", "Surrogate\nTraining", "SASTO\nOptimize", "Validate\n& Certify"]
    pw = 2.50; pg = 0.28
    ptot = len(plabels) * pw + (len(plabels) - 1) * pg
    ps = LX + (LW - ptot) / 2
    for i, lab in enumerate(plabels):
        px = ps + i * (pw + pg)
        rect(sl, px, pipe_y, pw, 0.52, fill=TEAL, r=25000)
        tb(sl, px, pipe_y + 0.02, pw, 0.48, lab, sz=Pt(9), bold=True, color=TXT_WHITE, align=PP_ALIGN.CENTER)
        if i < len(plabels) - 1:
            tb(sl, px + pw + 0.02, pipe_y + 0.10, pg - 0.04, 0.30, "->",
               sz=Pt(14), bold=True, color=GOLD, align=PP_ALIGN.CENTER)

    y += obj_h + SG
    print("  L3 Objectives")

    # ─── L4: PROBLEM FRAMING (expanded — absorbs Design Criteria) ──
    y = hdr(sl, LX, y, LW, "PROBLEM FRAMING")
    pf_h = CB - y  # fill remaining space
    card(sl, LX, y, LW, pf_h)

    py = y + CP

    # Material & Design Criteria (folded in from removed section)
    tb(sl, LX + CP, py, LW - 2*CP, 0.20, "Material & Design Constraints",
       sz=Pt(12), bold=True, color=TEAL)
    py += 0.24

    crit_rows = [
        ["Parameter", "Value", "Source"],
        ["E (Young's modulus)", "25 GPa", "Isotropic concrete"],
        ["v (Poisson's ratio)", "0.20", "--"],
        ["f'c (compressive)", "30 MPa", "Standard mix"],
        ["sigma_VM,allow", "5.0 MPa", "f'c / (gamma_m * gamma_f)"],
        ["Dead load", "Self-weight", "ASCE 7-22 ASD"],
        ["Live load", "1.92 kPa", "ASCE 7-22 Table 4-1"],
        ["Wind load", "0.96 kPa lateral", "ASCE 7-22 Ch. 26-30"],
    ]
    _, cty = tbl(sl, LX + CP, py, LW - 2*CP, crit_rows, [0.40, 0.30, 0.30])
    py = cty + 0.08

    eq_pill(sl, LX + CP, py, LW - 2*CP,
            "sigma_VM,allow = f'c / (gamma_m x gamma_f) = 30 / (3.0 x 2.0) = 5.0 MPa")
    py += 0.42
    eq_pill(sl, LX + CP, py, LW - 2*CP,
            "Conservative bound: mu_sigma + k*sigma_sigma,  k=1.0 -> P(violation) <= 0.09%")
    py += 0.48

    # Optimization objective
    tb(sl, LX + CP, py, LW - 2*CP, 0.20, "Optimization Objective",
       sz=Pt(12), bold=True, color=TXT_DARK)
    py += 0.22
    f16 = fig("fig16_optimization_objective")
    if f16:
        img_safe(sl, f16, LX + CP, py, LW - 2*CP, 1.20)
    py += 1.24

    # Sensitivity formula
    tb(sl, LX + CP, py, LW - 2*CP, 0.20, "Sensitivity via Surrogate Backpropagation",
       sz=Pt(12), bold=True, color=TXT_DARK)
    py += 0.22
    f17 = fig("fig17_sensitivity_formula")
    if f17:
        img_safe(sl, f17, LX + CP, py, LW - 2*CP, 1.20)
    py += 1.24

    # Part-aware thickness
    tb(sl, LX + CP, py, LW - 2*CP, 0.20, "Part-Aware Thickness Constraint",
       sz=Pt(12), bold=True, color=TXT_DARK)
    py += 0.22
    f14 = fig("fig14_part_aware_thickness")
    remaining = pf_h - (py - y) - CP - 0.30
    if f14 and remaining > 0.5:
        img_safe(sl, f14, LX + CP, py, LW - 2*CP, remaining)
    cap(sl, LX + CP, y + pf_h - CP - 0.22, LW - 2*CP,
        "Fig. 5. Part-aware: 86.8% interior wall removal, no exterior degradation")

    print("  L4 Problem Framing (with design criteria)")

    # ═══════════════════════════════════════════════════════════
    # CENTER PANEL
    # ═══════════════════════════════════════════════════════════
    y = CT

    # ─── C1: ENGINEERING METHODOLOGY ──────────────────────────
    y = hdr(sl, CX, y, CW, "ENGINEERING METHODOLOGY")
    c1_h = 16.50
    card(sl, CX, y, CW, c1_h)

    gut = 0.14
    sw = (CW - 2*CP - gut) / 2
    th = 7.80
    bh_c = c1_h - 2*CP - th - gut

    sx1 = CX + CP
    sx2 = sx1 + sw + gut
    sy1 = y + CP
    sy2 = sy1 + th + gut

    # ── Top-left: Dataset Pipeline ──
    rect(sl, sx1, sy1, sw, th, fill=CARD_FILL, border=CARD_BDR, bw=Pt(0.5), r=30000)
    tb(sl, sx1 + 0.08, sy1 + 0.06, sw - 0.16, 0.22,
       "Dataset Generation Pipeline", sz=Pt(12), bold=True, color=TXT_DARK)

    f03 = fig("fig03_dataset_pipeline")
    if f03:
        img_safe(sl, f03, sx1 + 0.08, sy1 + 0.30, sw - 0.16, 3.00)
    cap(sl, sx1 + 0.08, sy1 + 3.32, sw - 0.16,
        "Fig. 6. 14,293 wireframes -> 11,178 FEA sims -> 128^3 voxel grids")

    ds_rows = [
        ["Split", "n", "Targets"],
        ["Train", "8,943", "sigma_VM, u_max, C"],
        ["Validation", "1,121", "--"],
        ["Test", "1,114", "Ground-truth FEA"],
    ]
    _, dty = tbl(sl, sx1 + 0.08, sy1 + 3.56, sw - 0.16, ds_rows, [0.30, 0.25, 0.45])

    tb(sl, sx1 + 0.08, dty + 0.06, sw - 0.16, 1.80,
       "Pipeline: (1) 3DWire wireframe -> volumetric parts via extrusion (ext walls 4 voxels ~316 mm, "
       "int walls 2 voxels ~158 mm); (2) boolean fusion via FreeCAD; (3) tetrahedral meshing via Gmsh "
       "(50k-200k elements); (4) SfePy FEA under ASCE 7-22 ASD loads; (5) voxelization onto 128^3 grid.\n\n"
       "Data filtering: 3,115/14,293 (21.8%) rejected for diverged displacement (>1.0 m), degenerate "
       "compliance (<1e-6 J), or invalid stress (<=0 Pa). Retained data spans 4.9 orders of magnitude "
       "in stress and 7.7 orders in compliance.",
       sz=Pt(9), color=TXT_DARK)

    # ── Top-right: Deep Ensemble Architecture ──
    rect(sl, sx2, sy1, sw, th, fill=CARD_FILL, border=CARD_BDR, bw=Pt(0.5), r=30000)
    tb(sl, sx2 + 0.08, sy1 + 0.06, sw - 0.16, 0.22,
       "Deep Ensemble Surrogate (5x8.76M params)", sz=Pt(12), bold=True, color=TXT_DARK)
    badge(sl, sx2 + sw - 1.50, sy1 + 0.06, 1.40, 0.22, "x5 ENSEMBLE", RED)

    f04 = fig("fig04_architecture")
    if f04:
        img_safe(sl, f04, sx2 + 0.08, sy1 + 0.32, sw - 0.16, 4.00)
    cap(sl, sx2 + 0.08, sy1 + 4.36, sw - 0.16,
        "Fig. 7. Dual pooling -> 512-d -> 3 scalar outputs per member")

    hp_rows = [
        ["Parameter", "Value"],
        ["Ensemble size", "M=5 (43.8M total)"],
        ["Input", "7-ch 128^3 + 10-d feature"],
        ["Regularization", "dropout 0.15, WD 1e-4"],
        ["Normalization", "log(1+|y|) -> z-score"],
        ["Outputs", "peak sigma_VM, max u, C"],
    ]
    _, hty = tbl(sl, sx2 + 0.08, sy1 + 4.60, sw - 0.16, hp_rows, [0.40, 0.60])

    tb(sl, sx2 + 0.08, hty + 0.06, sw - 0.16, 1.20,
       "Architecture: 4 conv stages (128->64->32->16->8), 3 SE-ResBlocks with squeeze-excitation, "
       "dual pooling (avg+max) -> 512-d embed + 128-d feature MLP. Head: 640->512->256->3 with skip.\n\n"
       "Training: Huber loss, AdamW lr=5e-4, cosine anneal, EMA 0.999, AMP, grad clip ||.||<=1. "
       "Aug: 90 deg rot, flips, noise sigma=0.02. Dropout 0.15, stochastic depth 0-0.1.",
       sz=Pt(9), color=TXT_DARK)

    # ── Bottom-left: SASTO Algorithm ──
    rect(sl, sx1, sy2, sw, bh_c, fill=CARD_FILL, border=CARD_BDR, bw=Pt(0.5), r=30000)
    tb(sl, sx1 + 0.08, sy2 + 0.06, sw - 0.16, 0.22,
       "SASTO Algorithm - Sensitivity-Guided Erosion", sz=Pt(12), bold=True, color=TXT_DARK)

    f05 = fig("fig05_sasto_flowchart")
    if f05:
        img_safe(sl, f05, sx1 + 0.08, sy2 + 0.32, sw - 0.16, 3.20)
    cap(sl, sx1 + 0.08, sy2 + 3.54, sw - 0.16,
        "Fig. 8. Three-phase optimization. Phase 1 removes >99% of material.")

    phases = [
        ("Phase 1: Sensitivity-Guided Erosion",
         "Sort by descending s_i; batch-remove 6-simple-points; if violated -> undo, halve B -> max(B/2,10). >99% of removal.",
         TEAL),
        ("Phase 2: Fine-Grained Endgame",
         "Re-run Phase 1 with B in {5, 1} to squeeze remaining feasible removals near constraint boundary.",
         GOLD),
        ("Phase 3: Swap Refinement",
         "Thick interior voxels (dist >= 3) swapped with removed neighbors; accept if volume decreases & constraints hold.",
         RED),
    ]
    phy = sy2 + 3.78
    for label, desc, clr in phases:
        rect(sl, sx1 + 0.08, phy, sw - 0.16, 0.20, fill=clr, r=15000)
        tb(sl, sx1 + 0.16, phy, sw - 0.32, 0.20, label, sz=Pt(9), bold=True, color=TXT_WHITE)
        phy += 0.22
        tb(sl, sx1 + 0.16, phy, sw - 0.32, 0.40, desc, sz=Pt(8), color=TXT_DARK)
        phy += 0.42

    tb(sl, sx1 + 0.08, phy, sw - 0.16, 0.48,
       "Post-processing: fill enclosed air pockets (<=50 voxels), remove shard voxels "
       "(<2 face-neighbors), SDF -> marching cubes -> Laplacian smoothing -> watertight STL.",
       sz=Pt(8), color=TXT_DARK)

    # ── Bottom-right: 6-Connectivity ──
    rect(sl, sx2, sy2, sw, bh_c, fill=CARD_FILL, border=CARD_BDR, bw=Pt(0.5), r=30000)
    tb(sl, sx2 + 0.08, sy2 + 0.06, sw - 0.16, 0.22,
       "Topology: 6-Connectivity Guarantee", sz=Pt(12), bold=True, color=TXT_DARK)

    f06 = fig("fig06_connectivity")
    if f06:
        img_safe(sl, f06, sx2 + 0.08, sy2 + 0.32, sw - 0.16, 2.80)
    cap(sl, sx2 + 0.08, sy2 + 3.16, sw - 0.16,
        "Fig. 9. 6-conn eliminates floating fragments incompatible with AM toolpaths")

    ppy = sy2 + 3.40
    rect(sl, sx2 + 0.08, ppy, sw - 0.16, 0.40, fill=EQ_BG, border=TEAL, bw=Pt(1), r=20000)
    tb(sl, sx2 + 0.14, ppy + 0.04, sw - 0.28, 0.32,
       "Proposition: A binary voxel field with exactly one 6-connected foreground component "
       "yields a single-component marching-cubes surface mesh.",
       sz=Pt(8), color=TXT_DARK, align=PP_ALIGN.CENTER)
    ppy += 0.44

    tb(sl, sx2 + 0.08, ppy, sw - 0.16, 0.20,
       "6-Simple Point Test (Kong & Rosenfeld 1989)", sz=Pt(10), bold=True, color=TXT_DARK)
    ppy += 0.22
    eq_pill(sl, sx2 + 0.08, ppy, sw - 0.16,
            "SP6(v) = 1 iff |C6(rho' & N26(v))| = 1 AND |C26(rho_bar' & N26(v) U {v})| = 1")
    ppy += 0.42
    tb(sl, sx2 + 0.08, ppy, sw - 0.16, 0.50,
       "Foreground uses 6-connectivity; background uses complementary 26-connectivity. "
       "This (6,26) pairing prevents diagonal-only attachments that violate marching cubes "
       "assumptions. Standard (26,6) pairing produces thousands of floating fragments.",
       sz=Pt(8), color=TXT_DARK)

    print("  C1 Methodology")

    y += c1_h + SG

    # ─── C2: RESULTS & IN-SILICO VALIDATION ──────────────────
    y = hdr(sl, CX, y, CW, "RESULTS & IN-SILICO VALIDATION")
    c2_h = CB - y
    card(sl, CX, y, CW, c2_h)

    cg = 0.14
    cw3 = (CW - 2*CP - 2*cg) / 3
    cx_arr = [CX + CP, CX + CP + cw3 + cg, CX + CP + 2*(cw3 + cg)]
    ry = y + CP

    # ── Col A: Reference Case — with 3D renders ──
    tb(sl, cx_arr[0], ry, cw3, 0.20, "Reference Case (Sample 00472)",
       sz=Pt(11), bold=True, color=TXT_DARK)

    # Reference case: HQ original vs SASTO-PA comparison
    ref_comp = fig("fig_ref_comparison")
    if ref_comp:
        img_safe(sl, ref_comp, cx_arr[0], ry + 0.24, cw3, 3.20)
    else:
        r_iso_comp = render("stl_pa_solid")
        if r_iso_comp:
            img_safe(sl, r_iso_comp, cx_arr[0], ry + 0.24, cw3, 2.80)
        else:
            ri_img = extracted_img("Image 11")
            if ri_img:
                img_safe(sl, ri_img, cx_arr[0], ry + 0.24, cw3, 2.80)

    badge(sl, cx_arr[0] + cw3/2 - 0.70, ry + 2.60, 1.40, 0.22, "-45.0% material", RED)
    cap(sl, cx_arr[0], ry + 3.48, cw3, "Before/after 3D renders")

    # Reference table
    f19 = fig("fig19_reference_table")
    if f19:
        img_safe(sl, f19, cx_arr[0], ry + 3.72, cw3, 2.00)

    tb(sl, cx_arr[0], ry + 5.76, cw3, 0.50,
       "SASTO-PA achieves 10.7pp more reduction than SASTO-U (34.3%) by permitting "
       "1-voxel interior walls.",
       sz=Pt(8), color=TXT_DARK)

    # ── Col B: Multi-Geometry ──
    tb(sl, cx_arr[1], ry, cw3, 0.20, "1,114-Geometry Generalization",
       sz=Pt(11), bold=True, color=TXT_DARK)

    f07 = fig("fig07_histogram")
    if f07:
        img_safe(sl, f07, cx_arr[1], ry + 0.24, cw3, 2.40)
    cap(sl, cx_arr[1], ry + 2.68, cw3, "Fig. 10. Volume reduction | n=1,114 | mu=23.5%+/-7.8%")

    tb(sl, cx_arr[1], ry + 2.90, cw3, 0.46,
       "Across all 1,114 held-out test geometries, SASTO-PA achieves mean 23.5% +/- 7.8% "
       "material reduction. Maximum 45.0% on individual designs.",
       sz=Pt(8), color=TXT_DARK)

    f08 = fig("fig08_per_part")
    if f08:
        img_safe(sl, f08, cx_arr[1], ry + 3.40, cw3, 2.10)
    cap(sl, cx_arr[1], ry + 5.54, cw3, "Fig. 11. Per-part retention: load-bearing >91%")

    tb(sl, cx_arr[1], ry + 5.76, cw3, 0.50,
       "Exterior walls and roof retain >91% of material (load-bearing). Interior walls are the "
       "primary target, consistent with non-structural classification.",
       sz=Pt(8), color=TXT_DARK)

    # ── Col C: Speedup & FEA ──
    tb(sl, cx_arr[2], ry, cw3, 0.20, "Speedup vs. SIMP",
       sz=Pt(11), bold=True, color=TXT_DARK)

    f09 = fig("fig09_speedup")
    if f09:
        img_safe(sl, f09, cx_arr[2], ry + 0.24, cw3, 2.00)
    cap(sl, cx_arr[2], ry + 2.28, cw3, "Fig. 12. SIMP 64^3: 94s vs SASTO 128^3: 50s")

    tb(sl, cx_arr[2], ry + 2.50, cw3, 0.40,
       "SASTO at 128^3 (median 50s, consumer GPU) vs. SIMP at 64^3 (94s) and 128^3 (77 min).",
       sz=Pt(8), color=TXT_DARK)

    tb(sl, cx_arr[2], ry + 2.94, cw3, 0.20, "Independent FEA Re-analysis (n=1,114)",
       sz=Pt(11), bold=True, color=TXT_DARK)

    f10 = fig("fig10_fea_compliance")
    if f10:
        img_safe(sl, f10, cx_arr[2], ry + 3.18, cw3, 2.30)
    badge(sl, cx_arr[2] + 0.06, ry + 3.28, 1.50, 0.20, "0/1,114 violations", TEAL)
    badge(sl, cx_arr[2] + 0.06, ry + 3.52, 1.60, 0.20, "P(violation) <= 0.09%", TEAL)
    cap(sl, cx_arr[2], ry + 5.52, cw3, "Fig. 13. All C_opt/C_base <= 1.15. Max: 1.004")

    tb(sl, cx_arr[2], ry + 5.74, cw3, 0.50,
       "Same-method FEA reanalysis on all 1,114 optimized designs: zero violations. "
       "Max compliance ratio = 1.004 (vs 1.15 limit). Confirms surrogate conservatism.",
       sz=Pt(8), color=TXT_DARK)

    # ── Bottom Stats Banner ──
    ban_h = 1.20
    by = y + c2_h - ban_h - CP
    rect(sl, CX + CP, by, CW - 2*CP, ban_h, fill=SEC_BAR, r=30000)

    f20 = fig("fig20_stats_banner")
    if f20:
        img_safe(sl, f20, CX + CP, by, CW - 2*CP, ban_h)
    else:
        stats = [
            ("23.5%", "Mean material\nreduction"),
            ("23-92x", "Speedup\nvs. SIMP"),
            ("0 / 1,114", "FEA constraint\nviolations"),
            ("50 sec", "Median\nruntime"),
        ]
        stw = (CW - 2*CP) / 4
        for i, (num, lab) in enumerate(stats):
            sx = CX + CP + i * stw
            tb(sl, sx, by + 0.08, stw, 0.55, num, font="Arial Black", sz=Pt(28),
               bold=True, color=GOLD, align=PP_ALIGN.CENTER)
            tb(sl, sx, by + 0.72, stw, 0.40, lab, sz=Pt(9), color=TXT_WHITE, align=PP_ALIGN.CENTER)
            if i > 0:
                sep = sl.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                          Inches(sx), Inches(by + 0.18), Inches(0.02), Inches(ban_h - 0.36))
                sep.fill.solid(); sep.fill.fore_color.rgb = TXT_WHITE; sep.line.fill.background()

    print("  C2 Results")

    # ═══════════════════════════════════════════════════════════
    # RIGHT PANEL — 4 sections
    # ═══════════════════════════════════════════════════════════
    y = CT

    # ─── R1: STATISTICAL ANALYSIS ─────────────────────────────
    y = hdr(sl, RX, y, RW, "STATISTICAL ANALYSIS")
    r1_h = 14.20  # slightly expanded
    card(sl, RX, y, RW, r1_h)

    ry2 = y + CP

    # Surrogate table
    tb(sl, RX + CP, ry2, RW - 2*CP, 0.20, "Surrogate Model Performance",
       sz=Pt(12), bold=True, color=TXT_DARK)
    ry2 += 0.22
    f15 = fig("fig15_surrogate_table")
    if f15:
        img_safe(sl, f15, RX + CP, ry2, RW - 2*CP, 1.40)
    ry2 += 1.44

    rect(sl, RX + CP, ry2, RW - 2*CP, 0.34, fill=EQ_BG, r=20000)
    tb(sl, RX + CP + 0.06, ry2 + 0.04, RW - 2*CP - 0.12, 0.26,
       "Ranking accuracy (Spearman rho), not pointwise prediction, drives optimization safety. "
       "Compliance rho=0.948 -> correct voxel ordering for removal.",
       sz=Pt(8), color=TXT_DARK, align=PP_ALIGN.CENTER)
    ry2 += 0.38

    # Convergence
    tb(sl, RX + CP, ry2, RW - 2*CP, 0.20, "Optimization Convergence",
       sz=Pt(12), bold=True, color=TXT_DARK)
    ry2 += 0.22
    f11 = fig("fig11_convergence")
    if f11:
        img_safe(sl, f11, RX + CP, ry2, RW - 2*CP, 2.30)
    ry2 += 2.34
    cap(sl, RX + CP, ry2, RW - 2*CP,
        "Fig. 14. SASTO-PA (teal) vs. SASTO-U (gold)")
    ry2 += 0.22
    tb(sl, RX + CP, ry2, RW - 2*CP, 0.36,
       "SASTO-PA converges to 45.0% reduction vs SASTO-U at 34.3%. Phase 1 removes >99% of "
       "material in large batches; Phases 2-3 squeeze remaining feasible voxels.",
       sz=Pt(8), color=TXT_DARK)
    ry2 += 0.40

    # k-Factor
    tb(sl, RX + CP, ry2, RW - 2*CP, 0.20, "k-Factor Sensitivity (Pareto Frontier)",
       sz=Pt(12), bold=True, color=TXT_DARK)
    ry2 += 0.22
    f12 = fig("fig12_k_factor")
    if f12:
        img_safe(sl, f12, RX + CP, ry2, RW - 2*CP, 2.00)
    badge(sl, RX + RW/2 - 0.80, ry2 + 1.20, 1.60, 0.20, "k=1.0 Operating Point", GOLD)
    ry2 += 2.04
    cap(sl, RX + CP, ry2, RW - 2*CP,
        "Fig. 15. Non-monotonic Pareto: both gate stringency and budget depend on k")
    ry2 += 0.22
    tb(sl, RX + CP, ry2, RW - 2*CP, 0.36,
       "k<0.5: under-conservative, risks violations. k>2.0: over-conservative, less removal. "
       "k=1.0 achieves Pareto-optimal: max reduction with zero violations.",
       sz=Pt(8), color=TXT_DARK)
    ry2 += 0.40

    # Conformal / UQ
    tb(sl, RX + CP, ry2, RW - 2*CP, 0.20, "Conformal Prediction & UQ",
       sz=Pt(12), bold=True, color=TXT_DARK)
    ry2 += 0.22
    f13 = fig("fig13_uncertainty")
    if f13:
        img_safe(sl, f13, RX + CP, ry2, RW - 2*CP, 1.80)
    ry2 += 1.84
    cap(sl, RX + CP, ry2, RW - 2*CP,
        "Fig. 16. Uncertainty bands. Gamma_D ~ 0.184 (sub-linear)")
    ry2 += 0.22
    tb(sl, RX + CP, ry2, RW - 2*CP, 0.50,
       "Ensemble disagreement D(phi) tracks distribution shift as material is removed. "
       "Gamma_D ~ 0.184 = sub-linear growth. Distribution-free conformal: P(violation) <= 0.09% "
       "(Clopper-Pearson, alpha=0.05, n=1,114). Ensemble residuals heavier-tailed than Gaussian.",
       sz=Pt(8), color=TXT_DARK)

    y += r1_h + SG
    print("  R1 Statistics")

    # ─── R2: CONCLUSIONS ──────────────────────────────────────
    y = hdr(sl, RX, y, RW, "CONCLUSIONS")
    r2_h = 5.80  # slightly expanded
    card(sl, RX, y, RW, r2_h)

    conclusions = [
        ("1", "SASTO achieves 23.5% +/- 7.8% mean material reduction across 1,114 held-out geometries, "
              "up to 45.0% on individual designs - the first surrogate-accelerated, topology-preserving "
              "optimizer at building scale."),
        ("2", "Deep ensemble surrogate (5x8.76M params) provides 23-92x speedup vs SIMP: median 50s on "
              "consumer GPU vs 19-77 min for SIMP."),
        ("3", "6-connectivity criterion eliminates thousands of floating mesh fragments produced by "
              "standard 26-connectivity, guaranteeing watertight single-component STLs for AM."),
        ("4", "Part-aware thickness yields 10.7pp more reduction than uniform baseline by permitting "
              "1-voxel (78 mm) interior walls while protecting 2-voxel (156 mm) load-bearing exteriors."),
        ("5", "Independent FEA re-analysis: zero violations across all 1,114 designs "
              "(max C_opt/C_base = 1.004 vs 1.15 limit). Conformal: P(violation) <= 0.09%."),
    ]
    coy = y + CP
    for num, text in conclusions:
        rect(sl, RX + CP, coy, 0.22, 0.22, fill=TEAL, r=50000)
        tb(sl, RX + CP, coy, 0.22, 0.22, num, sz=Pt(10), bold=True, color=TXT_WHITE, align=PP_ALIGN.CENTER)
        tb(sl, RX + CP + 0.28, coy, RW - 2*CP - 0.28, 0.90, text, sz=Pt(9), color=TXT_DARK)
        coy += 0.96

    imy = y + r2_h - 0.38
    rect(sl, RX + CP, imy, RW - 2*CP, 0.32, fill=EQ_BG, border=GOLD, bw=Pt(1), r=20000)
    tb(sl, RX + CP + 0.06, imy + 0.03, RW - 2*CP - 0.12, 0.26,
       "Impact: 8% of global CO2 = cement. 23.5% less concrete per house -> millions of tons saved.",
       sz=Pt(8), bold=True, color=TXT_DARK, align=PP_ALIGN.CENTER)

    y += r2_h + SG
    print("  R2 Conclusions")

    # ─── R3: FUTURE WORK ──────────────────────────────────────
    y = hdr(sl, RX, y, RW, "FUTURE WORK")
    r3_h = 4.20  # slightly expanded
    card(sl, RX, y, RW, r3_h)

    future = [
        ("1", "FEA-IN-THE-LOOP ACTIVE LEARNING",
         "When ensemble Gamma_D > threshold tau, trigger ground-truth FEA mid-optimization -> "
         "self-correcting safety net."),
        ("2", "NONLINEAR FEA SPOT CHECKS",
         "Concrete damaged plasticity (CDP) on 5 representative designs to assess tension cracking, "
         "compression softening, and buckling."),
        ("3", "PHYSICAL PRINT VALIDATION",
         "1:10 scale print of optimized house. Compression testing + digital image correlation (DIC) "
         "for full-field strain mapping."),
    ]
    fy = y + CP
    for num, title, desc in future:
        rect(sl, RX + CP, fy, 0.22, 0.22, fill=TEAL, r=50000)
        tb(sl, RX + CP, fy, 0.22, 0.22, num, sz=Pt(10), bold=True, color=TXT_WHITE, align=PP_ALIGN.CENTER)
        tb(sl, RX + CP + 0.28, fy, RW - 2*CP - 0.28, 0.20, title, sz=Pt(10), bold=True, color=SEC_BAR)
        tb(sl, RX + CP + 0.28, fy + 0.22, RW - 2*CP - 0.28, 0.70, desc, sz=Pt(9), color=TXT_DARK)
        fy += 1.00

    # Pipeline boxes
    pp_y = fy + 0.06
    ppl = ["Optimized STL", "3D Print 1:10", "Compression + DIC"]
    ppw = 3.30; ppg = 0.30
    ppt = len(ppl) * ppw + (len(ppl) - 1) * ppg
    pps = RX + (RW - ppt) / 2
    for i, lab in enumerate(ppl):
        px = pps + i * (ppw + ppg)
        rect(sl, px, pp_y, ppw, 0.34, fill=TEAL, r=20000)
        tb(sl, px, pp_y + 0.04, ppw, 0.26, lab, sz=Pt(9), bold=True, color=TXT_WHITE, align=PP_ALIGN.CENTER)
        if i < len(ppl) - 1:
            tb(sl, px + ppw + 0.02, pp_y + 0.02, ppg - 0.04, 0.30, "->",
               sz=Pt(13), bold=True, color=GOLD, align=PP_ALIGN.CENTER)

    y += r3_h + SG
    print("  R3 Future Work")

    # ─── R4: KEY REFERENCES ───────────────────────────────────
    y = hdr(sl, RX, y, RW, "KEY REFERENCES")
    r4_h = CB - y
    card(sl, RX, y, RW, r4_h)

    refs = [
        "1. Bendsoe & Sigmund (2003). Topology Optimization. Springer.",
        "2. Lakshminarayanan et al. (2017). Deep ensembles. NeurIPS.",
        "3. Kong & Rosenfeld (1989). Digital topology. CVGIP 48(3).",
        "4. ASCE (2022). ASCE/SEI 7-22 Min Design Loads.",
        "5. Buswell et al. (2018). 3D printing concrete. Cem Concr Res 112.",
        "6. Sigmund & Maute (2013). Topology optimization. SMO 48.",
        "7. Lin et al. (2024). 3DWire dataset. KAUST VCC.",
        "8. IEA (2021). Global Status Report for Buildings.",
        "9. Lorensen & Cline (1987). Marching Cubes. SIGGRAPH.",
        "10. Vovk et al. (2005). Algorithmic Learning. Springer.",
    ]
    ref_paras = [{"text": r, "sz": Pt(9), "color": TXT_DARK, "sb": 0, "sa": 1} for r in refs]
    mtb(sl, RX + CP, y + CP, RW - 2*CP, r4_h - 2*CP, ref_paras)
    print("  R4 References")

    # ── SAVE ──
    prs.save(OUTPUT)
    print(f"\nSaved: {OUTPUT}")
    print(f"  {BW} x {BH} in, ~{len(sl.shapes)} shapes")


if __name__ == "__main__":
    build()
