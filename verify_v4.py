"""Verify v4 poster: check section positions, overlaps, and coverage."""
import sys
sys.stdout.reconfigure(encoding='utf-8')
from pptx import Presentation

prs = Presentation("SASTO_ISEF_Poster_v4.pptx")
sl = prs.slides[0]
print(f"Shapes: {len(sl.shapes)}")

# Get all shapes with positions
shapes = []
for i, s in enumerate(sl.shapes):
    l = s.left/914400; t = s.top/914400; w = s.width/914400; h = s.height/914400
    txt = ""
    if s.has_text_frame:
        txt = s.text_frame.text.replace('\n',' ')[:60]
    shapes.append((i, l, t, w, h, str(s.shape_type), txt))

# Print section headers (text in caps, short)
print("\nSection Headers:")
for i, l, t, w, h, st, txt in shapes:
    if txt.isupper() and 5 < len(txt) < 50 and h < 0.50:
        panel = "LEFT" if l < 12 else ("CENTER" if l < 36 else "RIGHT")
        print(f"  y={t:5.2f}  {panel:7s}  {txt}")

# Print images
print("\nImages:")
imgs = 0
for i, l, t, w, h, st, txt in shapes:
    if "PICTURE" in st:
        imgs += 1
        panel = "LEFT" if l < 12 else ("CENTER" if l < 36 else "RIGHT")
        print(f"  [{l:5.1f},{t:5.1f}] {w:5.2f}x{h:5.2f}  {panel}")
print(f"Total images: {imgs}")

# Check bottom extent
max_bottom = 0
for i, l, t, w, h, st, txt in shapes:
    b = t + h
    if b > max_bottom:
        max_bottom = b
print(f"\nMax bottom extent: {max_bottom:.2f} in (board=36.00)")

# Check for content that extends beyond board
overflows = 0
for i, l, t, w, h, st, txt in shapes:
    if t + h > 36.05 or l + w > 48.05:
        overflows += 1
        print(f"  OVERFLOW S{i}: [{l:.1f},{t:.1f}] {w:.1f}x{h:.1f} '{txt[:40]}'")
if overflows == 0:
    print("No overflows detected ✓")

# Check panel coverage (what % of area is filled)
for panel_name, px, pw in [("LEFT", 0, 12), ("CENTER", 12, 24), ("RIGHT", 36, 12)]:
    panel_area = pw * (36 - 3.35)  # exclude title band
    filled = 0
    for i, l, t, w, h, st, txt in shapes:
        if l >= px and l < px + pw and t > 3.25:
            filled += w * h
    coverage = min(filled / panel_area * 100, 100)
    print(f"{panel_name} coverage: ~{coverage:.0f}%")
