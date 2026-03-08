from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.shapes import MSO_SHAPE_TYPE
import sys
sys.stdout.reconfigure(encoding='utf-8')

prs = Presentation('SASTO_ISEF_Poster_v2.pptx')
slide = prs.slides[0]
for j, shape in enumerate(slide.shapes):
    left = round(shape.left / 914400, 2) if shape.left else 0
    top = round(shape.top / 914400, 2) if shape.top else 0
    w = round(shape.width / 914400, 2) if shape.width else 0
    h = round(shape.height / 914400, 2) if shape.height else 0
    txt = ""
    if shape.has_text_frame:
        txt = shape.text[:80].replace('\n','|')
    img = ""
    try:
        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE or shape.shape_type == 13:
            img = f" IMG:{shape.image.content_type},{len(shape.image.blob)}b"
    except:
        pass
    print(f"{j:3d} ({left:6.2f},{top:6.2f}) {w:6.2f}x{h:5.2f} {shape.name:20s}{img} {txt[:60]}")

print(f"\nTotal shapes: {len(slide.shapes)}")
