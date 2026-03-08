from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.shapes import MSO_SHAPE_TYPE
import sys
sys.stdout.reconfigure(encoding='utf-8')

prs = Presentation('SASTO_ISEF_Poster_v2.pptx')
for i, slide in enumerate(prs.slides):
    print(f'=== Slide {i+1} ===')
    print(f'  Width: {prs.slide_width.inches:.2f}in, Height: {prs.slide_height.inches:.2f}in')
    for j, shape in enumerate(slide.shapes):
        left = round(shape.left / 914400, 2) if shape.left else 0
        top = round(shape.top / 914400, 2) if shape.top else 0
        w = round(shape.width / 914400, 2) if shape.width else 0
        h = round(shape.height / 914400, 2) if shape.height else 0
        print(f'  Shape {j}: type={shape.shape_type}, name="{shape.name}", pos=({left},{top}), size=({w}x{h})')
        if shape.has_text_frame:
            for pi, para in enumerate(shape.text_frame.paragraphs):
                txt = para.text[:120]
                if txt.strip():
                    # Get font info from first run
                    font_info = ""
                    if para.runs:
                        r = para.runs[0]
                        font_info = f" [font={r.font.name}, size={r.font.size}, bold={r.font.bold}, color={r.font.color.rgb if r.font.color and r.font.color.rgb else 'none'}]"
                    print(f'    p{pi}: "{txt}"{font_info}')
        try:
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE or shape.shape_type == 13:
                print(f'    IMAGE: {shape.image.content_type}, {len(shape.image.blob)} bytes')
        except:
            pass
        # Check for fill
        try:
            fill = shape.fill
            if fill.type is not None:
                print(f'    fill: type={fill.type}')
        except:
            pass
