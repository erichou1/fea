#!/usr/bin/env python3
"""
ISEF Poster PPTX v4 — Dense, reference-matching layout.

Key changes from v3:
  - Margins: 0.35 → 0.18"
  - Section gaps: 0.10 → 0.04"
  - Card padding: 0.18 → 0.10"
  - Section header fonts: 22pt → 28pt, height 0.60 → 0.45"
  - MUCH more text content from the paper
  - Images smaller, integrated side-by-side with text
  - Tables fill full card width
  - Every section packed with content (no empty space)
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
import os

# ═══════════════════════════════════════════════════════════════
# COLOR PALETTE
# ═══════════════════════════════════════════════════════════════
BG_NAVY   = RGBColor(0x06, 0x2B, 0x7A)
TITLE_BAND= RGBColor(0x03, 0x20, 0x61)
SEC_BAR   = RGBColor(0x0A, 0x3D, 0x9A)
CARD_FILL = RGBColor(0xF7, 0xF9, 0xFC)
CARD_BDR  = RGBColor(0xB7, 0xC5, 0xE3)
TXT_DARK  = RGBColor(0x0B, 0x17, 0x36)
TXT_WHITE = RGBColor(0xFF, 0xFF, 0xFF)
RED       = RGBColor(0xD7, 0x26, 0x3D)
TEAL      = RGBColor(0x00, 0x8C, 0x9E)
GOLD      = RGBColor(0xCF, 0xA5, 0x35)
EQ_BG     = RGBColor(0xE8, 0xEE, 0xF2)
ALT_ROW   = RGBColor(0xEE, 0xF2, 0xFA)

# ═══════════════════════════════════════════════════════════════
# DIMENSIONS — tight like reference poster
# ═══════════════════════════════════════════════════════════════
BW = 48.0;  BH = 36.0
TITLE_H = 3.25
CT = TITLE_H + 0.10   # content top
CB = 35.82             # content bottom — use nearly all space
M  = 0.18              # margin from panel edge
SG = 0.04              # section gap (reference-tight)
HH = 0.45              # section header height
CP = 0.10              # card internal padding

# Panel x/w
LX = M;            LW = 12.0 - 2*M   # 11.64
CX = 12.0 + M;     CW = 24.0 - 2*M   # 23.64
RX = 36.0 + M;     RW = 12.0 - 2*M   # 11.64

CARD_R = 40000  # corner radius EMU

BASE = os.path.dirname(os.path.abspath(__file__))
EXTRACTED = os.path.join(BASE, "poster_images_extracted")
POSTER_FIGS = os.path.join(BASE, "poster_figures")
OUTPUT = os.path.join(BASE, "SASTO_ISEF_Poster_v4.pptx")

def img_path(name):
    for ext in (".jpg", ".png"):
        p = os.path.join(EXTRACTED, name + ext)
        if os.path.isfile(p): return p
    return None

def pfig(name):
    p = os.path.join(POSTER_FIGS, name + ".png")
    return p if os.path.isfile(p) else None

# ═══════════════════════════════════════════════════════════════
# HELPERS
# ═══════════════════════════════════════════════════════════════

def _set_radius(shape, r=CARD_R):
    spPr = shape._element.spPr
    pg = spPr.find(qn("a:prstGeom"))
    if pg is not None:
        pg.set("prst", "roundRect")
        av = pg.find(qn("a:avLst"))
        if av is None:
            av = spPr.makeelement(qn("a:avLst"), {})
            pg.append(av)
        for c in list(av): av.remove(c)
        gd = spPr.makeelement(qn("a:gd"), {"name":"adj","fmla":f"val {r}"})
        av.append(gd)

def rect(sl, x, y, w, h, fill=None, border=None, bw=Pt(1), r=CARD_R):
    s = sl.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    if fill: s.fill.solid(); s.fill.fore_color.rgb = fill
    else: s.fill.background()
    if border: s.line.color.rgb = border; s.line.width = bw
    else: s.line.fill.background()
    _set_radius(s, r)
    return s

def card(sl, x, y, w, h):
    return rect(sl, x, y, w, h, fill=CARD_FILL, border=CARD_BDR)

def tb(sl, x, y, w, h, text, font="Arial", sz=Pt(12), bold=False, italic=False,
       color=TXT_DARK, align=PP_ALIGN.LEFT):
    t = sl.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = t.text_frame; tf.word_wrap = True; tf.auto_size = None
    p = tf.paragraphs[0]; p.text = text
    p.font.name = font; p.font.size = sz; p.font.bold = bold
    p.font.italic = italic
    if color: p.font.color.rgb = color
    p.alignment = align
    p.space_before = Pt(0); p.space_after = Pt(0)
    return t

def mtb(sl, x, y, w, h, paras):
    """Multi-paragraph text box. paras = list of dicts with text, sz, bold, italic, color, align, sb, sa."""
    t = sl.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = t.text_frame; tf.word_wrap = True; tf.auto_size = None
    for i, pd in enumerate(paras):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = pd.get("text", "")
        p.font.name = pd.get("font", "Arial")
        p.font.size = pd.get("sz", Pt(12))
        p.font.bold = pd.get("bold", False)
        p.font.italic = pd.get("italic", False)
        if pd.get("color"): p.font.color.rgb = pd["color"]
        p.alignment = pd.get("align", PP_ALIGN.LEFT)
        p.space_before = Pt(pd.get("sb", 1))
        p.space_after = Pt(pd.get("sa", 1))
    return t

def hdr(sl, x, y, w, text):
    """Section header bar with centered text."""
    rect(sl, x, y, w, HH, fill=SEC_BAR, r=15000)
    tb(sl, x, y, w, HH, text, font="Arial Black", sz=Pt(28), bold=True,
       color=TXT_WHITE, align=PP_ALIGN.CENTER)
    return y + HH

def img_safe(sl, f, x, y, w, h):
    if not f or not os.path.isfile(f): return None
    from PIL import Image as PILImage
    with PILImage.open(f) as im: iw, ih = im.size
    ar = iw / ih; bar = w / h
    if ar > bar: fw=w; fh=w/ar; fx=x; fy=y+(h-fh)/2
    else: fh=h; fw=h*ar; fx=x+(w-fw)/2; fy=y
    return sl.shapes.add_picture(f, Inches(fx), Inches(fy), Inches(fw), Inches(fh))

def tbl(sl, x, y, w, rows, col_pcts):
    """Table. col_pcts = list of fractions summing to 1.0."""
    nr = len(rows); nc = len(rows[0])
    rh = 0.22
    th = nr * rh
    ts = sl.shapes.add_table(nr, nc, Inches(x), Inches(y), Inches(w), Inches(th))
    t = ts.table
    for ci, pct in enumerate(col_pcts):
        t.columns[ci].width = Inches(w * pct)
    for ri, row in enumerate(rows):
        for ci, txt in enumerate(row):
            c = t.cell(ri, ci); c.text = txt
            for p in c.text_frame.paragraphs:
                p.font.name = "Arial"; p.font.size = Pt(10)
                p.alignment = PP_ALIGN.LEFT
                if ri == 0:
                    p.font.bold = True; p.font.color.rgb = TXT_WHITE; p.font.size = Pt(10)
                else:
                    p.font.color.rgb = TXT_DARK
            if ri == 0: c.fill.solid(); c.fill.fore_color.rgb = SEC_BAR
            elif ri % 2 == 0: c.fill.solid(); c.fill.fore_color.rgb = ALT_ROW
            else: c.fill.solid(); c.fill.fore_color.rgb = CARD_FILL
    return ts, y + th

def eq_pill(sl, x, y, w, text, h=0.38):
    rect(sl, x, y, w, h, fill=EQ_BG, r=25000)
    tb(sl, x+0.06, y+0.03, w-0.12, h-0.06, text, sz=Pt(10), align=PP_ALIGN.CENTER)
    return y + h

def badge(sl, x, y, w, h, text, color=TEAL):
    rect(sl, x, y, w, h, fill=color, r=30000)
    tb(sl, x, y, w, h, text, sz=Pt(9), bold=True, color=TXT_WHITE, align=PP_ALIGN.CENTER)

def cap(sl, x, y, w, text):
    """Figure caption."""
    tb(sl, x, y, w, 0.20, text, sz=Pt(8), italic=True, color=TXT_DARK, align=PP_ALIGN.CENTER)
    return y + 0.20

# ═══════════════════════════════════════════════════════════════
# BUILD
# ═══════════════════════════════════════════════════════════════

def build():
    prs = Presentation()
    prs.slide_width = Inches(BW); prs.slide_height = Inches(BH)
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    print("Building v4 poster...")

    # ── BACKGROUND ──
    bg = sl.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(BW), Inches(BH))
    bg.fill.solid(); bg.fill.fore_color.rgb = BG_NAVY; bg.line.fill.background()

    # ── TITLE BAND ──
    tb_bg = sl.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(BW), Inches(TITLE_H))
    tb_bg.fill.solid(); tb_bg.fill.fore_color.rgb = TITLE_BAND; tb_bg.line.fill.background()
    # Gold line
    gl = sl.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(TITLE_H - 0.03), Inches(BW), Inches(0.03))
    gl.fill.solid(); gl.fill.fore_color.rgb = GOLD; gl.line.fill.background()

    # -- User already fixed the title in their v3(1), preserve their layout approach --
    # House images (user added two in their version): use extracted images
    h0 = img_path("Image 0")
    if h0: img_safe(sl, h0, 0.50, 0.30, 2.80, 2.40)
    # Second house if available (user's modified version had two)
    h11 = img_path("Image 11")
    if h11: img_safe(sl, h11, 4.00, 0.30, 3.00, 2.40)

    # Title text
    tb(sl, 8.50, 0.10, 30.00, 0.90,
       "SURROGATE-ACCELERATED STRUCTURAL OPTIMIZATION",
       font="Arial Black", sz=Pt(44), bold=True, color=TXT_WHITE, align=PP_ALIGN.CENTER)
    tb(sl, 8.50, 1.10, 30.00, 0.65,
       "Additive Manufacturing: Harnessing FEA to Optimize Material Efficiency",
       font="Arial", sz=Pt(26), bold=True, italic=True, color=TXT_WHITE, align=PP_ALIGN.CENTER)
    tb(sl, 8.50, 1.85, 30.00, 0.45,
       "Eric Hou",
       font="Arial", sz=Pt(22), bold=True, color=TXT_WHITE, align=PP_ALIGN.CENTER)
    # Credit
    mtb(sl, 39.00, 0.15, 8.60, 2.80, [
        {"text": "Intel ISEF 2026 — Engineering Mechanics", "sz": Pt(13), "bold": True, "color": TXT_WHITE},
        {"text": "", "sz": Pt(4)},
        {"text": "Credit Line of Origin", "sz": Pt(12), "bold": True, "color": TXT_WHITE},
        {"text": "References & data in references below:", "sz": Pt(10), "color": TXT_WHITE},
        {"text": "", "sz": Pt(4)},
        {"text": "Images denoted with an asterisk (*) are from", "sz": Pt(10), "color": TXT_WHITE},
        {"text": "public-domain or adapted from publicly available", "sz": Pt(10), "color": TXT_WHITE},
        {"text": "sources. All other figures created by Eric Hou.", "sz": Pt(10), "color": TXT_WHITE},
    ])
    print("  Title done")

    # ═══════════════════════════════════════════════════════════
    # LEFT PANEL
    # ═══════════════════════════════════════════════════════════
    y = CT

    # ─── VISUAL ABSTRACT ──────
    y = hdr(sl, LX, y, LW, "VISUAL ABSTRACT")
    va_h = 4.00
    card(sl, LX, y, LW, va_h)
    va_img = img_path("Image 1")
    if va_img: img_safe(sl, va_img, LX+CP, y+CP, LW-2*CP, 2.70)
    cap(sl, LX+CP, y+2.80, LW-2*CP,
        "Fig. 1. SASTO pipeline: 3D wireframe → volumetric parts → FEA dataset → deep ensemble surrogate → sensitivity-guided optimization → watertight STL")
    # Brief description below caption
    tb(sl, LX+CP, y+3.02, LW-2*CP, 0.85,
       "The visual abstract shows the end-to-end SASTO workflow. A 3DWire building wireframe is extruded into four structural parts (exterior walls, interior walls, roof, floor), meshed and simulated with FEA to train a 5-member deep ensemble. At optimization time, the surrogate replaces FEA, enabling 23–92× faster topology optimization with zero constraint violations.",
       sz=Pt(10), color=TXT_DARK)
    y += va_h + SG
    print("  L1 Visual Abstract")

    # ─── INTRODUCTION ─────────
    y = hdr(sl, LX, y, LW, "INTRODUCTION")
    intro_h = 5.80
    card(sl, LX, y, LW, intro_h)

    # Dense text on left, images on right
    tw = 6.80; iw = LW - tw - 3*CP  # ~4.54
    ix = LX + tw + 2*CP

    intro_text = [
        {"text": "GLIOBLASTOMA & CO₂ CRISIS", "sz": Pt(12), "bold": True, "color": TEAL, "sa": 1},
        {"text": "Concrete production accounts for ~8% of global CO₂ emissions [IEA 2021]. In conventional construction, walls are built at uniform thickness—a practice driven by formwork constraints, not structural need. Interior partitions bear negligible loads compared to exterior shear walls, creating a substantial opportunity for material reduction.", "sz": Pt(10), "color": TXT_DARK, "sa": 3},
        {"text": "ADDITIVE MANUFACTURING", "sz": Pt(12), "bold": True, "color": TEAL, "sa": 1},
        {"text": "Large-scale 3D printing (ICON, COBOD, Apis Cor) can realize arbitrary wall profiles at no marginal tooling cost. However, exploiting this geometric freedom requires optimized 3D models that minimize volume, satisfy structural constraints under ASCE 7-22 load cases, and produce watertight meshes compatible with printer toolpath generation.", "sz": Pt(10), "color": TXT_DARK, "sa": 3},
        {"text": "COMPUTATIONAL BOTTLENECK", "sz": Pt(12), "bold": True, "color": TEAL, "sa": 1},
        {"text": "Classical SIMP requires hundreds–thousands of FEA evaluations, each costing minutes–hours at building scale. Voxel-based implementations using 26-connectivity produce disconnected floating fragments incompatible with AM. No prior work combines surrogate-accelerated optimization with formal mesh connectivity guarantees at building scale.", "sz": Pt(10), "color": TXT_DARK, "sa": 3},
        {"text": "SASTO CONTRIBUTION", "sz": Pt(12), "bold": True, "color": RED, "sa": 1},
        {"text": "SASTO replaces FEA with a deep ensemble surrogate, achieving 23–92× speedup. A 6-connectivity criterion eliminates floating fragments. Part-aware thickness enables differential thinning. Across 1,114 geometries: 23.5%±7.8% mean reduction, zero FEA violations.", "sz": Pt(10), "bold": True, "color": TXT_DARK, "sa": 2},
    ]
    mtb(sl, LX+CP, y+CP, tw, intro_h-2*CP, intro_text)

    # Images on right
    i2 = img_path("Image 2")
    i3 = img_path("Image 3")
    if i2: img_safe(sl, i2, ix, y+CP, iw, 2.40)
    if i3: img_safe(sl, i3, ix, y+CP+2.50, iw, 2.40)
    cap(sl, ix, y+CP+4.95, iw, "Before vs. after optimization*")

    y += intro_h + SG
    print("  L2 Introduction")

    # ─── RESEARCH OBJECTIVES ──
    y = hdr(sl, LX, y, LW, "RESEARCH OBJECTIVES")
    obj_h = 2.50
    card(sl, LX, y, LW, obj_h)

    objs = [
        ("1", "SPEED", "23–92× speedup via deep ensemble surrogate replacing FEA during optimization", TEAL),
        ("2", "PRINTABILITY", "6-connected topology → single-mesh watertight STL for AM toolpath", TEAL),
        ("3", "EFFICIENCY", "23.5%±7.8% mean material reduction (up to 45.0%) across 1,114 geometries", GOLD),
        ("4", "SAFETY", "0/1,114 FEA violations; conformal P(violation) ≤ 0.09%", RED),
    ]
    oy = y + CP
    for num, title, desc, bc in objs:
        rect(sl, LX+CP, oy, 0.26, 0.26, fill=bc, r=50000)
        tb(sl, LX+CP, oy, 0.26, 0.26, num, sz=Pt(11), bold=True, color=TXT_WHITE, align=PP_ALIGN.CENTER)
        tb(sl, LX+CP+0.34, oy-0.01, 1.40, 0.24, title, sz=Pt(11), bold=True, color=bc)
        tb(sl, LX+CP+1.76, oy-0.01, LW-2*CP-1.76, 0.24, desc, sz=Pt(10), color=TXT_DARK)
        oy += 0.32

    # Pipeline boxes
    pipe_y = oy + 0.08
    plabels = ["Data\nGeneration", "Surrogate\nTraining", "SASTO\nOptimize", "Validate\n& Certify"]
    pw = 2.50; pg = 0.28
    ptot = len(plabels)*pw + (len(plabels)-1)*pg
    ps = LX + (LW - ptot)/2
    for i,lab in enumerate(plabels):
        px = ps + i*(pw+pg)
        rect(sl, px, pipe_y, pw, 0.52, fill=TEAL, r=25000)
        tb(sl, px, pipe_y+0.02, pw, 0.48, lab, sz=Pt(9), bold=True, color=TXT_WHITE, align=PP_ALIGN.CENTER)
        if i < len(plabels)-1:
            tb(sl, px+pw+0.02, pipe_y+0.10, pg-0.04, 0.30, "►", sz=Pt(14), bold=True, color=GOLD, align=PP_ALIGN.CENTER)

    y += obj_h + SG
    print("  L3 Objectives")

    # ─── ENGINEERING DESIGN CRITERIA ──
    y = hdr(sl, LX, y, LW, "ENGINEERING DESIGN CRITERIA")
    crit_h = 3.50
    card(sl, LX, y, LW, crit_h)

    crit_rows = [
        ["Constraint", "Limit", "Basis"],
        ["Von Mises stress", "σ_VM ≤ 5.0 MPa", "f'c/(γ_m×γ_f) = 30/(3×2)"],
        ["Compliance ratio", "C_opt/C_base ≤ 1.15", "15% stiffness budget"],
        ["Displacement", "u_max ≤ L/360 ≈ 28 mm", "ASCE 7-22 serviceability"],
        ["Wall t (exterior)", "2Δx ≈ 156 mm", "Load-bearing path"],
        ["Wall t (interior)", "1Δx ≈ 78 mm", "Non-structural partition"],
        ["Mesh integrity", "1 connected component", "AM printability req."],
    ]
    _, ty = tbl(sl, LX+CP, y+CP, LW-2*CP, crit_rows, [0.35, 0.35, 0.30])

    # Equations below table
    ty += 0.06
    eq_pill(sl, LX+CP, ty, LW-2*CP,
            "σ_VM,allow = f'c / (γ_m × γ_f) = 30 MPa / (3.0 × 2.0) = 5.0 MPa")
    ty += 0.42
    eq_pill(sl, LX+CP, ty, LW-2*CP,
            "Conservative bound: μ̂_σ + k·σ̂_σ,  k = 1.0  →  P(violation) ≤ 0.09%  (conformal)")
    ty += 0.42
    # Additional text
    tb(sl, LX+CP, ty, LW-2*CP, 0.60,
       "Material: isotropic linear elastic concrete. E = 25 GPa, ν = 0.20, ρ = 2,400 kg/m³, f'c = 30 MPa. Loading: ASCE 7-22 ASD combinations — dead (self-weight), live (1.92 kPa), lateral wind (0.96 kPa). Fixed-base BC at min-x face. Displacement inactive in all 1,114 cases (max ~ 10⁻⁴ m, four orders below limit).",
       sz=Pt(9), color=TXT_DARK)

    y += crit_h + SG
    print("  L4 Design Criteria")

    # ─── PROBLEM FRAMING ──────
    y = hdr(sl, LX, y, LW, "PROBLEM FRAMING")
    pf_h = CB - y
    card(sl, LX, y, LW, pf_h)

    py = y + CP
    # Objective
    tb(sl, LX+CP, py, LW-2*CP, 0.20, "Optimization Objective", sz=Pt(12), bold=True, color=TXT_DARK)
    py += 0.22
    py = eq_pill(sl, LX+CP, py, LW-2*CP,
                 "min J(ρ) = w_V·(V/V₀) + w_S·(S/V₀) + P_constraint(ρ)    [w_V=1.0, w_S=0.01, κ=100]")
    py += 0.04
    tb(sl, LX+CP, py, LW-2*CP, 0.55,
       "V(ρ) = Σᵢρᵢ (total volume), S(ρ) = ½·Σᵢρᵢ·Σⱼ∈N₆(i)(1−ρⱼ) (exposed surface regularizer). Constraint penalty aggregates all structural violations with κ = 100 to approximate hard constraints. The upper-bound ensemble estimates: σ̂⁺ = μ_σ + k·σ_σ, Ĉ⁺ = μ_C + k·σ_C.",
       sz=Pt(9), color=TXT_DARK)
    py += 0.58

    # Sensitivity
    tb(sl, LX+CP, py, LW-2*CP, 0.20, "Sensitivity via Surrogate Backpropagation", sz=Pt(12), bold=True, color=TXT_DARK)
    py += 0.22
    py = eq_pill(sl, LX+CP, py, LW-2*CP,
                 "sᵢ = (1/M)·Σₘ ∂/∂ρᵢ[fₘ(C) + 0.3·fₘ(σ)]     sᵢ > 0 → safe to remove  |  sᵢ < 0 → essential")
    py += 0.04
    tb(sl, LX+CP, py, LW-2*CP, 0.45,
       "Each sensitivity computation requires M=5 forward+backward passes through the CNN (~3–8 s on RTX A3000), replacing a full FEA adjoint solve requiring minutes. Ensemble averaging: Var[s̄ᵢ] ≈ Var[sᵢ⁽¹⁾]/M → √5 ≈ 2.24× reduction in gradient noise.",
       sz=Pt(9), color=TXT_DARK)
    py += 0.50

    # Part-aware thickness — TEXT + smaller image
    tb(sl, LX+CP, py, LW-2*CP, 0.20, "Part-Aware Thickness Constraint", sz=Pt(12), bold=True, color=TXT_DARK)
    py += 0.22
    py = eq_pill(sl, LX+CP, py, LW-2*CP,
                 "t_min(p) = 2Δx = 156 mm for {ext, roof, floor}  |  1Δx = 78 mm for interior walls")
    py += 0.06

    # Image + text side by side
    fig5 = img_path("Image 5")
    img5_w = 5.20
    txt5_w = LW - 2*CP - img5_w - 0.10
    if fig5: img_safe(sl, fig5, LX+CP, py, img5_w, pf_h - (py - y) - CP - 0.22)
    tb(sl, LX+CP+img5_w+0.10, py, txt5_w, pf_h - (py - y) - CP - 0.22,
       "The part-aware formulation differentially thins interior partitions (1-voxel minimum, ~78 mm) while protecting load-bearing exterior walls, roof, and floor (2-voxel minimum, ~156 mm). Exterior walls are the primary lateral and gravity load path; interior walls serve as non-structural partitions.\n\nThis heterogeneous constraint is the key to SASTO-PA achieving 10.7pp more material reduction than the uniform-thickness baseline (SASTO-U). The distance transform identifies interior surface voxels whose depth exceeds t_min(p), making them eligible for removal.\n\nThe 3-voxel exterior shell surface is excluded from modification entirely, protecting the building envelope.",
       sz=Pt(9), color=TXT_DARK)
    cap(sl, LX+CP, y + pf_h - CP - 0.18, img5_w, "Fig. 5. Part-aware thickness schematic")

    print("  L5 Problem Framing")

    # ═══════════════════════════════════════════════════════════
    # CENTER PANEL
    # ═══════════════════════════════════════════════════════════
    y = CT

    # ─── ENGINEERING METHODOLOGY ──
    y = hdr(sl, CX, y, CW, "ENGINEERING METHODOLOGY")
    c1_h = 16.50
    card(sl, CX, y, CW, c1_h)

    # 2×2 grid
    gut = 0.14
    sw = (CW - 2*CP - gut) / 2   # ~11.63
    th = 7.80  # top row height
    bh = c1_h - 2*CP - th - gut  # bottom row

    sx1 = CX + CP
    sx2 = sx1 + sw + gut
    sy1 = y + CP
    sy2 = sy1 + th + gut

    # ── C1-A: Dataset Generation (top-left) ──
    rect(sl, sx1, sy1, sw, th, fill=CARD_FILL, border=CARD_BDR, bw=Pt(0.5), r=30000)
    tb(sl, sx1+0.08, sy1+0.06, sw-0.16, 0.22,
       "Dataset Generation Pipeline", sz=Pt(12), bold=True, color=TXT_DARK)

    d_img = img_path("Image 6")
    if d_img: img_safe(sl, d_img, sx1+0.08, sy1+0.30, sw-0.16, 3.00)
    cap(sl, sx1+0.08, sy1+3.32, sw-0.16,
        "Fig. 6. 14,293 wireframes → 11,178 FEA sims → 128³ voxel grids")

    # Dataset table
    ds_rows = [
        ["Split", "n", "Targets"],
        ["Train", "8,943", "σ_VM, u_max, C"],
        ["Validation", "1,121", "—"],
        ["Test", "1,114", "Ground-truth FEA"],
    ]
    _, dty = tbl(sl, sx1+0.08, sy1+3.56, sw-0.16, ds_rows, [0.30, 0.25, 0.45])

    # Dense text below table
    tb(sl, sx1+0.08, dty+0.06, sw-0.16, 1.80,
       "Pipeline: (1) 3DWire wireframe → volumetric parts via extrusion (ext walls 4 voxels ≈ 316 mm, int walls 2 voxels ≈ 158 mm); (2) boolean fusion via FreeCAD; (3) tetrahedral meshing via Gmsh (50k–200k elements); (4) SfePy FEA under ASCE 7-22 ASD loads; (5) voxelization onto 128³ grid with part labels.\n\nData filtering: 3,115/14,293 (21.8%) rejected for diverged displacement (>1.0 m), degenerate compliance (<10⁻⁶ J), or invalid stress (≤0 Pa). Retained data spans 4.9 orders of magnitude in stress and 7.7 orders in compliance, motivating log-transform normalization.",
       sz=Pt(9), color=TXT_DARK)

    # ── C1-B: Deep Ensemble (top-right) ──
    rect(sl, sx2, sy1, sw, th, fill=CARD_FILL, border=CARD_BDR, bw=Pt(0.5), r=30000)
    tb(sl, sx2+0.08, sy1+0.06, sw-0.16, 0.22,
       "Deep Ensemble Surrogate (5×8.76M params)", sz=Pt(12), bold=True, color=TXT_DARK)
    badge(sl, sx2+sw-1.50, sy1+0.06, 1.40, 0.22, "×5 ENSEMBLE", RED)

    # Image 7 is portrait (0.62 ratio) — put on left side, text on right
    a_img = img_path("Image 7")
    img7w = 4.00
    txt7x = sx2 + 0.08 + img7w + 0.10
    txt7w = sw - 0.16 - img7w - 0.10
    if a_img: img_safe(sl, a_img, sx2+0.08, sy1+0.32, img7w, 4.80)
    mtb(sl, txt7x, sy1+0.32, txt7w, 4.80, [
        {"text": "Architecture:", "sz": Pt(10), "bold": True, "color": TXT_DARK, "sa": 1},
        {"text": "• 4 conv stages: 128³→64³→32³→16³→8³", "sz": Pt(9), "color": TXT_DARK, "sa": 0},
        {"text": "• 3 SE-ResBlocks with squeeze-excitation", "sz": Pt(9), "color": TXT_DARK, "sa": 0},
        {"text": "• Dual pooling (avg+max) → 512-d embed", "sz": Pt(9), "color": TXT_DARK, "sa": 0},
        {"text": "• + 128-d feature MLP (material+loads)", "sz": Pt(9), "color": TXT_DARK, "sa": 0},
        {"text": "• Head: 640→512→256→3 with skip conn", "sz": Pt(9), "color": TXT_DARK, "sa": 2},
        {"text": "Training:", "sz": Pt(10), "bold": True, "color": TXT_DARK, "sa": 1},
        {"text": "• Huber loss (SmoothL1)", "sz": Pt(9), "color": TXT_DARK, "sa": 0},
        {"text": "• AdamW, lr=5×10⁻⁴, cosine anneal", "sz": Pt(9), "color": TXT_DARK, "sa": 0},
        {"text": "• EMA decay 0.999, AMP, grad clip ‖·‖≤1", "sz": Pt(9), "color": TXT_DARK, "sa": 0},
        {"text": "• Aug: 90° rot, flips, noise σ=0.02", "sz": Pt(9), "color": TXT_DARK, "sa": 0},
        {"text": "• Dropout 0.15, stochastic depth 0–0.1", "sz": Pt(9), "color": TXT_DARK, "sa": 0},
        {"text": "• Input: log(1+|y|) → z-score normalize", "sz": Pt(9), "color": TXT_DARK, "sa": 2},
        {"text": "Outputs: peak σ_VM, max u, compliance C", "sz": Pt(10), "bold": True, "color": TXT_DARK, "sa": 1},
        {"text": "Predicts global scalars (not field), enabling fast backprop sensitivity + ensemble UQ.", "sz": Pt(9), "color": TXT_DARK, "sa": 0},
    ])

    # Hyperparameter table below
    hp_rows = [
        ["Parameter", "Value"],
        ["Ensemble size", "M = 5 (43.8M total)"],
        ["Input", "7-ch 128³ + 10-d feature"],
        ["Regularization", "dropout 0.15, WD 10⁻⁴"],
        ["Normalization", "log(1+|y|) → z-score, 2/98‰ clip"],
    ]
    _, hty = tbl(sl, sx2+0.08, sy1+5.20, sw-0.16, hp_rows, [0.40, 0.60])

    cap(sl, sx2+0.08, hty+0.04, sw-0.16,
        "Fig. 7. Dual pooling → 512-d embed → 3 scalar outputs per member")
    tb(sl, sx2+0.08, hty+0.26, sw-0.16, 0.80,
       "Each member: ~8.76M parameters. Ensemble provides epistemic uncertainty estimates via member disagreement. Gradient variance reduced by √M ≈ 2.24× vs. single model. Trained on 4× NVIDIA GB200 GPUs.",
       sz=Pt(9), color=TXT_DARK)

    # ── C1-C: SASTO Algorithm (bottom-left) ──
    rect(sl, sx1, sy2, sw, bh, fill=CARD_FILL, border=CARD_BDR, bw=Pt(0.5), r=30000)
    tb(sl, sx1+0.08, sy2+0.06, sw-0.16, 0.22,
       "SASTO Algorithm — Sensitivity-Guided Erosion", sz=Pt(12), bold=True, color=TXT_DARK)

    a_img8 = img_path("Image 8")
    if a_img8: img_safe(sl, a_img8, sx1+0.08, sy2+0.32, sw-0.16, 3.20)
    cap(sl, sx1+0.08, sy2+3.54, sw-0.16,
        "Fig. 8. Three-phase optimization. Phase 1 removes >99% of material.")

    # Phase labels
    phases = [
        ("Phase 1: Sensitivity-Guided Erosion", "Sort candidates by descending sᵢ; batch-remove 6-simple-points; if constraints violated → undo, halve B → max(B/2, 10). Accounts for >99% of total removal.", TEAL),
        ("Phase 2: Fine-Grained Endgame", "Re-run Phase 1 with B ∈ {5, 1} to squeeze remaining feasible removals near the constraint boundary.", GOLD),
        ("Phase 3: Swap Refinement", "Thick interior voxels (dist ≥ 3) swapped with removed neighbors; accept if volume decreases and constraints hold.", RED),
    ]
    phy = sy2 + 3.78
    for label, desc, clr in phases:
        rect(sl, sx1+0.08, phy, sw-0.16, 0.20, fill=clr, r=15000)
        tb(sl, sx1+0.16, phy, sw-0.32, 0.20, label, sz=Pt(9), bold=True, color=TXT_WHITE)
        phy += 0.22
        tb(sl, sx1+0.16, phy, sw-0.32, 0.42, desc, sz=Pt(8), color=TXT_DARK)
        phy += 0.44

    # Post-processing note
    tb(sl, sx1+0.08, phy, sw-0.16, 0.50,
       "Post-processing: fill enclosed air pockets (≤50 voxels), remove shard voxels (<2 face-neighbors), compute SDF → marching cubes → Laplacian smoothing → watertight STL. 90% produce single-component mesh directly; remaining 10% need trivial fragment removal (mean 1.2 components, max 4).",
       sz=Pt(8), color=TXT_DARK)

    # ── C1-D: 6-Connectivity (bottom-right) ──
    rect(sl, sx2, sy2, sw, bh, fill=CARD_FILL, border=CARD_BDR, bw=Pt(0.5), r=30000)
    tb(sl, sx2+0.08, sy2+0.06, sw-0.16, 0.22,
       "Topology: 6-Connectivity Guarantee", sz=Pt(12), bold=True, color=TXT_DARK)

    # Side-by-side comparison
    hw = (sw - 0.36) / 2  # half width
    cy1 = sy2 + 0.32

    # 26-conn FAILS
    rect(sl, sx2+0.08, cy1, hw, 2.60, fill=CARD_FILL, border=RED, bw=Pt(1.5), r=25000)
    tb(sl, sx2+0.08, cy1+0.02, hw, 0.18, "26-CONN — FAILS", sz=Pt(9), bold=True, color=RED, align=PP_ALIGN.CENTER)
    i9 = img_path("Image 9")
    if i9: img_safe(sl, i9, sx2+0.16, cy1+0.22, hw-0.16, 1.60)
    tb(sl, sx2+0.08, cy1+1.86, hw, 0.70,
       "✗ Thousands of floating fragments\n✗ Diagonal-only connections\n✗ Marching cubes incompatible\n✗ Cannot generate AM toolpath",
       sz=Pt(8), bold=True, color=RED, align=PP_ALIGN.CENTER)

    # vs.
    tb(sl, sx2+0.08+hw, cy1+0.90, 0.20, 0.30, "vs.", sz=Pt(11), bold=True, color=RED, align=PP_ALIGN.CENTER)

    # 6-conn WORKS
    rect(sl, sx2+0.08+hw+0.20, cy1, hw, 2.60, fill=CARD_FILL, border=TEAL, bw=Pt(1.5), r=25000)
    tb(sl, sx2+0.08+hw+0.20, cy1+0.02, hw, 0.18, "6-CONN — WORKS", sz=Pt(9), bold=True, color=TEAL, align=PP_ALIGN.CENTER)
    i10 = img_path("Image 10")
    if i10: img_safe(sl, i10, sx2+0.16+hw+0.20, cy1+0.22, hw-0.16, 1.60)
    tb(sl, sx2+0.08+hw+0.20, cy1+1.86, hw, 0.70,
       "✓ 1 connected component\n✓ Face-share adjacency only\n✓ Watertight STL confirmed\n✓ AM toolpath compatible",
       sz=Pt(8), bold=True, color=TEAL, align=PP_ALIGN.CENTER)

    # Proposition box
    ppy = cy1 + 2.68
    rect(sl, sx2+0.08, ppy, sw-0.16, 0.50, fill=EQ_BG, border=TEAL, bw=Pt(1), r=20000)
    tb(sl, sx2+0.14, ppy+0.04, sw-0.28, 0.42,
       "Proposition: A binary voxel field with exactly one 6-connected foreground component yields a single-component marching-cubes surface mesh. Proof by induction: two 6-adjacent voxels share a face whose 4 dual-grid vertices satisfy ψ ≤ 0 → connected triangle patches.",
       sz=Pt(8), bold=False, color=TXT_DARK, align=PP_ALIGN.CENTER)

    # Definition box
    dpy = ppy + 0.54
    tb(sl, sx2+0.08, dpy, sw-0.16, 0.20,
       "6-Simple Point Test (Kong & Rosenfeld 1989)", sz=Pt(10), bold=True, color=TXT_DARK)
    dpy += 0.22
    eq_pill(sl, sx2+0.08, dpy, sw-0.16,
            "SP₆(v) = 1  iff  |C₆(ρ'∩N₂₆(v))| = 1  ∧  |C₂₆(ρ̄'∩N₂₆(v)∪{v})| = 1")
    dpy += 0.42
    tb(sl, sx2+0.08, dpy, sw-0.16, 0.60,
       "Foreground uses 6-connectivity; background uses complementary 26-connectivity. This (6,26) pairing follows digital topology convention and prevents diagonal-only attachments that violate marching cubes assumptions. The standard (26,6) pairing in topology optimization literature produces thousands of floating fragments.",
       sz=Pt(8), color=TXT_DARK)

    cap(sl, sx2+0.08, sy2+bh-0.22, sw-0.16,
        "Fig. 9. 6-conn eliminates floating fragments incompatible with AM toolpaths")
    print("  C1 Methodology")

    y += c1_h + SG

    # ─── RESULTS & IN-SILICO VALIDATION ──
    y = hdr(sl, CX, y, CW, "RESULTS & IN-SILICO VALIDATION")
    c2_h = CB - y
    card(sl, CX, y, CW, c2_h)

    # 3-column layout
    cg = 0.14
    cw3 = (CW - 2*CP - 2*cg) / 3
    cx = [CX+CP, CX+CP+cw3+cg, CX+CP+2*(cw3+cg)]
    ry = y + CP

    # ── Col A: Reference Case ──
    tb(sl, cx[0], ry, cw3, 0.20, "Reference Case (Sample 00472)", sz=Pt(11), bold=True, color=TXT_DARK)
    ri = img_path("Image 11")
    if ri: img_safe(sl, ri, cx[0], ry+0.24, cw3, 2.80)
    badge(sl, cx[0]+cw3/2-0.70, ry+2.20, 1.40, 0.22, "−45.0% material", RED)
    cap(sl, cx[0], ry+3.08, cw3, "Before/after optimization renders")

    # Results table
    ref_rows = [
        ["Metric", "Baseline", "SASTO-PA"],
        ["Vol. reduction", "—", "45.0% ★"],
        ["VM stress (Pa)", "3.08×10⁶", "3.08×10⁶ ✓"],
        ["Compliance ratio", "1.00", "1.004 ✓"],
        ["Mesh components", "1", "1 ✓"],
        ["Runtime", "—", "160 s"],
        ["EI Index", "—", "0.358 ★"],
    ]
    _, rty = tbl(sl, cx[0], ry+3.32, cw3, ref_rows, [0.38, 0.30, 0.32])
    tb(sl, cx[0], rty+0.04, cw3, 0.60,
       "Table 1. SASTO-PA achieves 10.7pp more reduction than SASTO-U (34.3%) by permitting 1-voxel interior walls. EI Index = ΔV/V₀ / [(σ⁺/σ_allow)(1+C⁺/C_allow)]. Higher = better material efficiency per unit of structural utilization.",
       sz=Pt(8), color=TXT_DARK)

    # ── Col B: Multi-Geometry ──
    tb(sl, cx[1], ry, cw3, 0.20, "1,114-Geometry Generalization", sz=Pt(11), bold=True, color=TXT_DARK)

    hf = pfig("fig10_histogram") or img_path("Image 12")
    if hf: img_safe(sl, hf, cx[1], ry+0.24, cw3, 2.40)
    cap(sl, cx[1], ry+2.68, cw3, "Fig. 10. Volume reduction histogram | n=1,114 | μ=23.5%±7.8%")

    tb(sl, cx[1], ry+2.90, cw3, 0.50,
       "Across all 1,114 held-out test geometries, SASTO-PA achieves mean 23.5%±7.8% (std) material reduction. Maximum 45.0% on individual designs. Part-aware rule yields 10.7pp more than uniform.",
       sz=Pt(8), color=TXT_DARK)

    ppf = pfig("fig11_per_part") or img_path("Image 13")
    if ppf: img_safe(sl, ppf, cx[1], ry+3.44, cw3, 2.10)
    cap(sl, cx[1], ry+5.58, cw3, "Fig. 11. Per-part retention: load-bearing >91%, interior = primary target")

    tb(sl, cx[1], ry+5.80, cw3, 0.50,
       "Exterior walls and roof retain >91% of material (load-bearing). Interior walls are the primary optimization target, consistent with their non-structural classification under ASCE 7-22.",
       sz=Pt(8), color=TXT_DARK)

    # ── Col C: Speedup & FEA ──
    tb(sl, cx[2], ry, cw3, 0.20, "Speedup vs. SIMP", sz=Pt(11), bold=True, color=TXT_DARK)

    sf = pfig("fig12_speedup") or img_path("Image 14")
    if sf: img_safe(sl, sf, cx[2], ry+0.24, cw3, 2.00)
    cap(sl, cx[2], ry+2.28, cw3, "Fig. 12. SIMP 64³: 94 s vs SASTO 128³: 50 s → 23–92× faster")

    tb(sl, cx[2], ry+2.50, cw3, 0.50,
       "SASTO at 128³ (median 50s, consumer GPU) vs. SIMP at 64³ (94s) and 128³ (77 min). SIMP requires hundreds of FEA solves; SASTO uses only surrogate forward+backward passes.",
       sz=Pt(8), color=TXT_DARK)

    tb(sl, cx[2], ry+3.04, cw3, 0.20, "Independent FEA Re-analysis (n=1,114)", sz=Pt(11), bold=True, color=TXT_DARK)

    ff = pfig("fig13_fea_compliance") or img_path("Image 15")
    if ff: img_safe(sl, ff, cx[2], ry+3.28, cw3, 2.30)
    badge(sl, cx[2]+0.06, ry+3.38, 1.50, 0.20, "0/1,114 violations", TEAL)
    badge(sl, cx[2]+0.06, ry+3.62, 1.60, 0.20, "P(violation) ≤ 0.09%", TEAL)
    cap(sl, cx[2], ry+5.62, cw3, "Fig. 13. All C_opt/C_base ≤ 1.15. Max observed: 1.004")

    tb(sl, cx[2], ry+5.84, cw3, 0.50,
       "Same-method FEA reanalysis on all 1,114 optimized designs: zero violations of stress, displacement, or stiffness constraints. Max compliance ratio = 1.004 (vs. 1.15 limit). Confirms surrogate conservatism.",
       sz=Pt(8), color=TXT_DARK)

    # ── Bottom Stats Banner ──
    bh2 = 1.20
    by = y + c2_h - bh2 - CP
    rect(sl, CX+CP, by, CW-2*CP, bh2, fill=SEC_BAR, r=30000)
    stats = [
        ("23.5%", "Mean material\nreduction"),
        ("23–92×", "Speedup\nvs. SIMP"),
        ("0 / 1,114", "FEA constraint\nviolations"),
        ("50 sec", "Median\nruntime"),
    ]
    stw = (CW - 2*CP) / 4
    for i, (num, lab) in enumerate(stats):
        sx = CX + CP + i * stw
        tb(sl, sx, by+0.08, stw, 0.55, num, font="Arial Black", sz=Pt(28), bold=True, color=GOLD, align=PP_ALIGN.CENTER)
        tb(sl, sx, by+0.72, stw, 0.40, lab, sz=Pt(9), color=TXT_WHITE, align=PP_ALIGN.CENTER)
        if i > 0:
            sep = sl.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(sx), Inches(by+0.18), Inches(0.02), Inches(bh2-0.36))
            sep.fill.solid(); sep.fill.fore_color.rgb = TXT_WHITE; sep.line.fill.background()

    print("  C2 Results")

    # ═══════════════════════════════════════════════════════════
    # RIGHT PANEL
    # ═══════════════════════════════════════════════════════════
    y = CT

    # ─── STATISTICAL ANALYSIS ─
    y = hdr(sl, RX, y, RW, "STATISTICAL ANALYSIS")
    r1_h = 13.80
    card(sl, RX, y, RW, r1_h)

    ry2 = y + CP
    # Surrogate metrics table
    tb(sl, RX+CP, ry2, RW-2*CP, 0.20, "Surrogate Model Performance", sz=Pt(12), bold=True, color=TXT_DARK)
    ry2 += 0.22
    sur_rows = [
        ["Target", "Spearman ρ", "R²_log", "MAPE%"],
        ["Von Mises stress", "0.737", "0.419", "37.4"],
        ["Displacement", "0.970 ★", "0.842", "10.9"],
        ["Compliance", "0.948 ★", "0.814", "18.5"],
    ]
    _, sty = tbl(sl, RX+CP, ry2, RW-2*CP, sur_rows, [0.35, 0.22, 0.22, 0.21])

    # Explanation
    rect(sl, RX+CP, sty+0.04, RW-2*CP, 0.34, fill=EQ_BG, r=20000)
    tb(sl, RX+CP+0.06, sty+0.07, RW-2*CP-0.12, 0.28,
       "Ranking accuracy (Spearman ρ), not pointwise prediction, drives optimization safety. Compliance ρ=0.948 → correct voxel ordering for removal. Stress ρ=0.737 is adequate because the penalty acts as a soft gate, not a hard classifier.",
       sz=Pt(8), color=TXT_DARK, align=PP_ALIGN.CENTER)
    ry2 = sty + 0.42

    # R1-B: Convergence
    tb(sl, RX+CP, ry2, RW-2*CP, 0.20, "Optimization Convergence", sz=Pt(12), bold=True, color=TXT_DARK)
    ry2 += 0.22
    cf = pfig("fig14_convergence") or img_path("Image 16")
    if cf: img_safe(sl, cf, RX+CP, ry2, RW-2*CP, 2.30)
    ry2 += 2.34
    cap(sl, RX+CP, ry2, RW-2*CP,
        "Fig. 14. SASTO-PA (teal) vs. SASTO-U (gold). Part-aware enables deeper removal.")
    ry2 += 0.22
    tb(sl, RX+CP, ry2, RW-2*CP, 0.40,
       "SASTO-PA converges to 45.0% reduction vs. SASTO-U at 34.3% on the reference geometry. Phase 1 removes >99% of material in large batches; Phases 2–3 squeeze the remaining feasible voxels near the constraint boundary.",
       sz=Pt(8), color=TXT_DARK)
    ry2 += 0.44

    # R1-C: k-Factor
    tb(sl, RX+CP, ry2, RW-2*CP, 0.20, "k-Factor Sensitivity (Pareto Frontier)", sz=Pt(12), bold=True, color=TXT_DARK)
    ry2 += 0.22
    kf = pfig("fig15_k_factor") or img_path("Image 17")
    if kf: img_safe(sl, kf, RX+CP, ry2, RW-2*CP, 2.00)
    badge(sl, RX+RW/2-0.80, ry2+1.20, 1.60, 0.20, "k=1.0 Operating Point", GOLD)
    ry2 += 2.04
    cap(sl, RX+CP, ry2, RW-2*CP,
        "Fig. 15. Non-monotonic Pareto: both gate stringency and budget depend on k.")
    ry2 += 0.22
    tb(sl, RX+CP, ry2, RW-2*CP, 0.40,
       "k<0.5: under-conservative → risks violations. k>2.0: over-conservative → less removal. k=1.0 achieves Pareto-optimal balance: maximum reduction with zero violations. Conformal calibration: k_conformal = 1.90 for 84.1% coverage (heavier-tailed than Gaussian).",
       sz=Pt(8), color=TXT_DARK)
    ry2 += 0.44

    # R1-D: Conformal Prediction
    tb(sl, RX+CP, ry2, RW-2*CP, 0.20, "Conformal Prediction & UQ", sz=Pt(12), bold=True, color=TXT_DARK)
    ry2 += 0.22
    uf = pfig("fig16_uncertainty") or img_path("Image 18")
    if uf: img_safe(sl, uf, RX+CP, ry2, RW-2*CP, 1.80)
    ry2 += 1.84
    cap(sl, RX+CP, ry2, RW-2*CP,
        "Fig. 16. Uncertainty bands widen during optimization. Γ_D ≈ 0.184 (sub-linear).")
    ry2 += 0.22
    tb(sl, RX+CP, ry2, RW-2*CP, 0.60,
       "Ensemble disagreement D(ϕ) = (1/T)·Σⱼ σⱼ/μⱼ tracks distribution shift as material is removed. Disagreement divergence rate Γ_D ≈ 0.184 signals sub-linear uncertainty growth. Distribution-free conformal certification on n=1,114 FEA-validated designs: P(violation) ≤ 0.09% (one-sided Clopper-Pearson at α=0.05). Ensemble residuals are heavier-tailed than Gaussian.",
       sz=Pt(8), color=TXT_DARK)

    y += r1_h + SG
    print("  R1 Statistics")

    # ─── CONCLUSIONS ──────────
    y = hdr(sl, RX, y, RW, "CONCLUSIONS")
    r2_h = 5.50
    card(sl, RX, y, RW, r2_h)

    conclusions = [
        ("1", "SASTO achieves 23.5%±7.8% mean material reduction across 1,114 held-out geometries, up to 45.0% on individual designs — the first surrogate-accelerated, topology-preserving optimizer at building scale."),
        ("2", "Deep ensemble surrogate (5×8.76M params) provides 23–92× speedup vs. SIMP: median 50 s on consumer GPU vs. 19–77 min for SIMP. Predicts global scalars enabling fast backprop sensitivity."),
        ("3", "6-connectivity criterion (new in building-scale TO) eliminates thousands of floating mesh fragments produced by standard 26-connectivity, guaranteeing watertight single-component STLs compatible with AM toolpath generation."),
        ("4", "Part-aware thickness yields 10.7pp more reduction than uniform baseline by permitting 1-voxel (78 mm) interior walls while protecting 2-voxel (156 mm) load-bearing exterior members."),
        ("5", "Independent FEA re-analysis: zero violations across all 1,114 designs (max C_opt/C_base = 1.004 vs. 1.15 limit). Conformal certification: P(violation) ≤ 0.09% at 95% confidence."),
    ]
    coy = y + CP
    for num, text in conclusions:
        rect(sl, RX+CP, coy, 0.22, 0.22, fill=TEAL, r=50000)
        tb(sl, RX+CP, coy, 0.22, 0.22, num, sz=Pt(10), bold=True, color=TXT_WHITE, align=PP_ALIGN.CENTER)
        tb(sl, RX+CP+0.28, coy, RW-2*CP-0.28, 0.80, text, sz=Pt(9), color=TXT_DARK)
        coy += 0.88

    # Impact strip
    imy = y + r2_h - 0.38
    rect(sl, RX+CP, imy, RW-2*CP, 0.32, fill=EQ_BG, border=GOLD, bw=Pt(1), r=20000)
    tb(sl, RX+CP+0.06, imy+0.03, RW-2*CP-0.12, 0.26,
       "Impact: 8% of global CO₂ = cement. 23.5% less concrete per house → millions of tons saved at scale if adopted by AM construction.",
       sz=Pt(8), bold=True, color=TXT_DARK, align=PP_ALIGN.CENTER)

    y += r2_h + SG
    print("  R2 Conclusions")

    # ─── FUTURE WORK ──────────
    y = hdr(sl, RX, y, RW, "FUTURE WORK")
    r3_h = 3.80
    card(sl, RX, y, RW, r3_h)

    future = [
        ("1", "FEA-IN-THE-LOOP ACTIVE LEARNING",
         "When ensemble Γ_D > threshold τ, trigger ground-truth FEA mid-optimization → self-correcting safety net. Extends to multi-story and seismic load cases where displacement may become active."),
        ("2", "NONLINEAR FEA SPOT CHECKS",
         "Concrete damaged plasticity (CDP) on 5 representative designs to assess tension cracking, compression softening, and buckling in 78 mm interior partitions."),
        ("3", "PHYSICAL PRINT VALIDATION",
         "1:10 scale concrete print of optimized house. Compression testing + digital image correlation (DIC) for full-field strain mapping. Compare measured vs. predicted structural response."),
    ]
    fy = y + CP
    for num, title, desc in future:
        rect(sl, RX+CP, fy, 0.22, 0.22, fill=TEAL, r=50000)
        tb(sl, RX+CP, fy, 0.22, 0.22, num, sz=Pt(10), bold=True, color=TXT_WHITE, align=PP_ALIGN.CENTER)
        tb(sl, RX+CP+0.28, fy, RW-2*CP-0.28, 0.20, title, sz=Pt(10), bold=True, color=SEC_BAR)
        tb(sl, RX+CP+0.28, fy+0.22, RW-2*CP-0.28, 0.64, desc, sz=Pt(9), color=TXT_DARK)
        fy += 0.92

    # Protocol pipeline
    pp_y = fy + 0.04
    ppl = ["Optimized STL", "3D Print 1:10", "Compression + DIC"]
    ppw = 3.30; ppg = 0.30
    ppt = len(ppl)*ppw + (len(ppl)-1)*ppg
    pps = RX + (RW - ppt) / 2
    for i, lab in enumerate(ppl):
        px = pps + i*(ppw+ppg)
        rect(sl, px, pp_y, ppw, 0.34, fill=TEAL, r=20000)
        tb(sl, px, pp_y+0.04, ppw, 0.26, lab, sz=Pt(9), bold=True, color=TXT_WHITE, align=PP_ALIGN.CENTER)
        if i < len(ppl)-1:
            tb(sl, px+ppw+0.02, pp_y+0.02, ppg-0.04, 0.30, "►", sz=Pt(13), bold=True, color=GOLD, align=PP_ALIGN.CENTER)

    y += r3_h + SG
    print("  R3 Future Work")

    # ─── KEY REFERENCES ───────
    y = hdr(sl, RX, y, RW, "KEY REFERENCES")
    r4_h = CB - y
    card(sl, RX, y, RW, r4_h)

    refs = [
        "1. Bendsøe & Sigmund (2003). Topology Optimization: Theory, Methods and Applications. Springer.",
        "2. Lakshminarayanan et al. (2017). Simple and scalable predictive uncertainty estimation. NeurIPS.",
        "3. Kong & Rosenfeld (1989). Digital topology: Introduction and survey. CVGIP 48(3).",
        "4. ASCE (2022). ASCE/SEI 7-22 Minimum Design Loads and Associated Criteria.",
        "5. Buswell et al. (2018). 3D printing using concrete extrusion. Cement & Concrete Research 112.",
        "6. Sigmund & Maute (2013). Topology optimization approaches. Struct Multidisc Optim 48.",
        "7. Lin et al. (2024). 3DWire: large-scale 3D wireframe dataset. KAUST.",
        "8. IEA (2021). Global Status Report for Buildings and Construction.",
        "9. Lorensen & Cline (1987). Marching Cubes. SIGGRAPH.",
        "10. Vovk et al. (2005). Algorithmic Learning in a Random World. Springer.",
    ]
    ref_paras = []
    for r in refs:
        ref_paras.append({"text": r, "sz": Pt(9), "color": TXT_DARK, "sb": 0, "sa": 1})
    mtb(sl, RX+CP, y+CP, RW-2*CP, r4_h-2*CP, ref_paras)

    print("  R4 References")

    # ── SAVE ──
    prs.save(OUTPUT)
    print(f"\n✓ Saved: {OUTPUT}")
    print(f"  {BW}×{BH} in, ~{len(sl.shapes)} shapes")

if __name__ == "__main__":
    build()
