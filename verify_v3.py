"""Quick verification of the rebuilt poster PPTX v3."""
import sys
sys.stdout.reconfigure(encoding='utf-8')
from pptx import Presentation
from pptx.util import Inches

prs = Presentation("SASTO_ISEF_Poster_v3.pptx")
slide = prs.slides[0]
w = prs.slide_width / 914400
h = prs.slide_height / 914400
print(f"Slide: {w:.1f} × {h:.1f} inches")
print(f"Total shapes: {len(slide.shapes)}")

# Count by type
from collections import Counter
types = Counter()
images = 0
texts = 0
for s in slide.shapes:
    types[s.shape_type] += 1
    if s.shape_type == 13:  # picture
        images += 1
    if s.has_text_frame:
        texts += 1

print(f"Shape types: {dict(types)}")
print(f"Images: {images}")
print(f"Text frames: {texts}")

# Check key sections by scanning text
section_headers = []
for s in slide.shapes:
    if s.has_text_frame:
        txt = s.text_frame.text.strip()
        if txt.isupper() and len(txt) > 5 and len(txt) < 50:
            left = s.left / 914400
            top = s.top / 914400
            section_headers.append((top, left, txt))

section_headers.sort()
print("\nSection headers found:")
for top, left, txt in section_headers:
    panel = "LEFT" if left < 12 else ("CENTER" if left < 36 else "RIGHT")
    print(f"  y={top:.2f}  {panel:8s}  {txt}")

# Check images
print("\nImages:")
for s in slide.shapes:
    if s.shape_type == 13:
        left = s.left / 914400
        top = s.top / 914400
        w = s.width / 914400
        h = s.height / 914400
        print(f"  [{left:.1f}, {top:.1f}] {w:.2f}×{h:.2f} in")
