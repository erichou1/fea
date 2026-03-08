#!/usr/bin/env python3
"""
ISEF Poster PPTX Rebuild — Complete rewrite matching reference poster style.

Fixes ALL 26 issues identified in the poster plan:
  1. Title band reduced to 3.25"
  2. Section headers CENTERED with Arial Black weight
  3. Card corner radius 6px on all cards
  4. Image corner radius 5px on all image frames
  5. Proper font sizes: title 44pt Black, subtitle 26pt Bold Italic, author 22pt Bold
  6. Right credit block 13pt (was 10.5pt)
  7. Section gaps tightened (0.10")
  8. Gradient background (radial approximation)
  9. Figures replaced with poster-styled versions where available
 10. Consistent card borders (#B7C5E3)
 11. No distorted images — aspect-ratio-locked placement
 12. Research objectives: text rows + horizontal pipeline boxes
 13. Bottom banner redesigned or made inline
 14. Left element: single house render, transparent bg
 15. Math formulas with proper Unicode symbols
 16. Image-beside-text layouts where appropriate
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
import os
import copy

# ═══════════════════════════════════════════════════════════════
# COLOR PALETTE
# ═══════════════════════════════════════════════════════════════
BG_NAVY = RGBColor(0x06, 0x2B, 0x7A)
BG_LIGHTER = RGBColor(0x0A, 0x3D, 0x9A)  # for gradient top
TITLE_BAND = RGBColor(0x03, 0x20, 0x61)
SECTION_BAR = RGBColor(0x0A, 0x3D, 0x9A)
CARD_FILL = RGBColor(0xF7, 0xF9, 0xFC)
CARD_BORDER = RGBColor(0xB7, 0xC5, 0xE3)
TEXT_DARK = RGBColor(0x0B, 0x17, 0x36)
TEXT_WHITE = RGBColor(0xFF, 0xFF, 0xFF)
ACCENT_RED = RGBColor(0xD7, 0x26, 0x3D)
ACCENT_TEAL = RGBColor(0x00, 0x8C, 0x9E)
ACCENT_GOLD = RGBColor(0xCF, 0xA5, 0x35)
EQ_BG = RGBColor(0xE8, 0xEE, 0xF2)
ALT_ROW = RGBColor(0xEE, 0xF2, 0xFA)
IMPACT_BG = RGBColor(0xE8, 0xEE, 0xF2)

# ═══════════════════════════════════════════════════════════════
# DIMENSIONS (inches)
# ═══════════════════════════════════════════════════════════════
BOARD_W = 48.0
BOARD_H = 36.0
TITLE_H = 3.25
GOLD_LINE_Y = TITLE_H - 0.02
CONTENT_TOP = TITLE_H + 0.18  # y where content starts
CONTENT_BOT = 35.65
MARGIN = 0.35
SECTION_GAP = 0.10  # tightened from 0.18
HEADER_H = 0.60     # section header bar height
CARD_PAD = 0.18     # inner padding of cards

# Panel boundaries
LEFT_X = MARGIN
LEFT_W = 12.0 - 2 * MARGIN  # 11.30
CENTER_X = 12.0 + MARGIN
CENTER_W = 24.0 - 2 * MARGIN  # 23.30
RIGHT_X = 36.0 + MARGIN
RIGHT_W = 12.0 - 2 * MARGIN  # 11.30

# Corner radius
CARD_RADIUS = 54000  # ~6px at 300dpi (in EMU, 1/914400 inch)
IMG_RADIUS = 45000   # ~5px

# ═══════════════════════════════════════════════════════════════
# PATHS
# ═══════════════════════════════════════════════════════════════
BASE = os.path.dirname(os.path.abspath(__file__))
EXTRACTED = os.path.join(BASE, "poster_images_extracted")
POSTER_FIGS = os.path.join(BASE, "poster_figures")
OUTPUT = os.path.join(BASE, "SASTO_ISEF_Poster_v3.pptx")


def img_path(name):
    """Get path for an extracted image."""
    for ext in (".jpg", ".png"):
        p = os.path.join(EXTRACTED, name + ext)
        if os.path.isfile(p):
            return p
    return None

def poster_fig(name):
    """Get path for a regenerated poster figure."""
    p = os.path.join(POSTER_FIGS, name + ".png")
    return p if os.path.isfile(p) else None


# ═══════════════════════════════════════════════════════════════
# HELPER FUNCTIONS
# ═══════════════════════════════════════════════════════════════

def set_corner_radius(shape, radius_emu=CARD_RADIUS):
    """Set corner radius on a shape via XML manipulation."""
    sp = shape._element
    prstGeom = sp.find(qn("a:prstGeom"), sp.nsmap) if hasattr(sp, 'nsmap') else None
    if prstGeom is None:
        # For auto shapes, try spPr
        spPr = sp.find(qn("p:spPr"))
        if spPr is None:
            spPr = sp.find(qn("xdr:spPr"))
        if spPr is not None:
            prstGeom = spPr.find(qn("a:prstGeom"))
    # Use roundRect preset
    spPr = shape._element.spPr
    prstGeom = spPr.find(qn("a:prstGeom"))
    if prstGeom is not None:
        prstGeom.set("prst", "roundRect")
        # Set corner size
        avLst = prstGeom.find(qn("a:avLst"))
        if avLst is None:
            avLst = spPr.makeelement(qn("a:avLst"), {})
            prstGeom.append(avLst)
        else:
            for child in list(avLst):
                avLst.remove(child)
        gd = spPr.makeelement(qn("a:gd"), {"name": "adj", "fmla": f"val {radius_emu}"})
        avLst.append(gd)
    return shape


def add_rounded_rect(slide, left, top, width, height, fill_color=None,
                     border_color=None, border_width=Pt(1), radius=CARD_RADIUS):
    """Add a rounded rectangle shape."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top),
        Inches(width), Inches(height)
    )
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = border_width
    else:
        shape.line.fill.background()

    # Set corner radius
    spPr = shape._element.spPr
    prstGeom = spPr.find(qn("a:prstGeom"))
    if prstGeom is not None:
        avLst = prstGeom.find(qn("a:avLst"))
        if avLst is None:
            avLst = spPr.makeelement(qn("a:avLst"), {})
            prstGeom.append(avLst)
        for child in list(avLst):
            avLst.remove(child)
        gd = spPr.makeelement(qn("a:gd"), {"name": "adj", "fmla": f"val {radius}"})
        avLst.append(gd)

    return shape


def add_text_box(slide, left, top, width, height, text, font_name="Arial",
                 font_size=Pt(13), bold=False, italic=False, color=TEXT_DARK,
                 align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP):
    """Add a text box with consistent styling."""
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top),
                                      Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.auto_size = None
    p = tf.paragraphs[0]
    p.text = text
    p.font.name = font_name
    p.font.size = font_size
    p.font.bold = bold
    p.font.italic = italic
    if color:
        p.font.color.rgb = color
    p.alignment = align

    # Try to set vertical anchor
    try:
        tf.paragraphs[0].space_before = Pt(0)
        tf.paragraphs[0].space_after = Pt(0)
    except:
        pass

    return txBox


def add_section_header(slide, x, y, w, text):
    """Add a section header bar with centered text."""
    # Background bar
    bar = add_rounded_rect(slide, x, y, w, HEADER_H,
                          fill_color=SECTION_BAR, border_color=None,
                          radius=20000)

    # Centered text
    tx = add_text_box(slide, x, y, w, HEADER_H, text,
                      font_name="Arial Black", font_size=Pt(22),
                      bold=True, color=TEXT_WHITE, align=PP_ALIGN.CENTER)
    return y + HEADER_H


def add_card(slide, x, y, w, h):
    """Add a white content card with border and corner radius."""
    return add_rounded_rect(slide, x, y, w, h,
                           fill_color=CARD_FILL, border_color=CARD_BORDER,
                           border_width=Pt(1), radius=CARD_RADIUS)


def add_image_safe(slide, img_file, left, top, width, height):
    """Add an image maintaining aspect ratio within the given bounds."""
    if not os.path.isfile(img_file):
        print(f"  WARNING: Image not found: {img_file}")
        return None

    from PIL import Image as PILImage
    with PILImage.open(img_file) as im:
        iw, ih = im.size

    aspect = iw / ih
    box_aspect = width / height

    if aspect > box_aspect:
        # Width-limited
        final_w = width
        final_h = width / aspect
        final_left = left
        final_top = top + (height - final_h) / 2
    else:
        # Height-limited
        final_h = height
        final_w = height * aspect
        final_left = left + (width - final_w) / 2
        final_top = top

    pic = slide.shapes.add_picture(img_file, Inches(final_left), Inches(final_top),
                                    Inches(final_w), Inches(final_h))
    return pic


def add_multi_para_text(slide, left, top, width, height, paragraphs_data):
    """Add a text box with multiple paragraphs.
    paragraphs_data: list of dicts with keys: text, font_size, bold, italic, color, align
    """
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top),
                                      Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.auto_size = None

    for i, pd in enumerate(paragraphs_data):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = pd.get("text", "")
        p.font.name = pd.get("font", "Arial")
        p.font.size = pd.get("font_size", Pt(13))
        p.font.bold = pd.get("bold", False)
        p.font.italic = pd.get("italic", False)
        if pd.get("color"):
            p.font.color.rgb = pd["color"]
        p.alignment = pd.get("align", PP_ALIGN.LEFT)
        p.space_before = Pt(pd.get("space_before", 2))
        p.space_after = Pt(pd.get("space_after", 2))

    return txBox


def add_table(slide, x, y, w, rows_data, col_widths, header_fill=SECTION_BAR):
    """Add a styled table.
    rows_data: list of lists of strings (first row is header)
    col_widths: list of floats (inches)
    Returns (table_shape, bottom_y).
    """
    n_rows = len(rows_data)
    n_cols = len(rows_data[0])
    row_h = 0.24
    table_h = n_rows * row_h

    table_shape = slide.shapes.add_table(n_rows, n_cols,
                                          Inches(x), Inches(y),
                                          Inches(w), Inches(table_h))
    table = table_shape.table

    # Set column widths
    for ci, cw in enumerate(col_widths):
        table.columns[ci].width = Inches(cw)

    for ri, row in enumerate(rows_data):
        for ci, cell_text in enumerate(row):
            cell = table.cell(ri, ci)
            cell.text = cell_text

            # Style
            for p in cell.text_frame.paragraphs:
                p.font.name = "Arial"
                p.font.size = Pt(11)
                p.alignment = PP_ALIGN.LEFT
                if ri == 0:
                    p.font.bold = True
                    p.font.color.rgb = TEXT_WHITE
                    p.font.size = Pt(11)
                else:
                    p.font.color.rgb = TEXT_DARK

            # Cell fill
            if ri == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = header_fill
            elif ri % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = ALT_ROW
            else:
                cell.fill.solid()
                cell.fill.fore_color.rgb = CARD_FILL

    return table_shape, y + table_h


# ═══════════════════════════════════════════════════════════════
# BUILD THE POSTER
# ═══════════════════════════════════════════════════════════════

def build_poster():
    prs = Presentation()
    prs.slide_width = Inches(BOARD_W)
    prs.slide_height = Inches(BOARD_H)

    # Add blank slide
    layout = prs.slide_layouts[6]  # blank
    slide = prs.slides.add_slide(layout)

    print("Building poster...")

    # ─── BACKGROUND ─────────────────────────────────────────────
    # Main background (navy)
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                 Inches(0), Inches(0),
                                 Inches(BOARD_W), Inches(BOARD_H))
    bg.fill.solid()
    bg.fill.fore_color.rgb = BG_NAVY
    bg.line.fill.background()

    # Gradient overlay — lighter at top-center (simulated with semi-transparent rect)
    grad_rect = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                        Inches(8), Inches(0),
                                        Inches(32), Inches(18))
    grad_fill = grad_rect.fill
    grad_fill.gradient()
    grad_fill.gradient_stops[0].color.rgb = BG_LIGHTER
    grad_fill.gradient_stops[0].position = 0.0
    try:
        grad_fill.gradient_stops[1].color.rgb = BG_NAVY
        grad_fill.gradient_stops[1].position = 1.0
    except:
        pass
    # Make it subtle by setting transparency
    try:
        for stop in grad_fill.gradient_stops:
            # Set alpha via XML
            srgb = stop._element.find(qn("a:srgbClr"))
            if srgb is not None:
                alpha = srgb.makeelement(qn("a:alpha"), {"val": "25000"})  # 25% opacity
                srgb.append(alpha)
    except Exception as e:
        print(f"  Note: gradient transparency not applied: {e}")
    grad_rect.line.fill.background()
    grad_rect.rotation = 270  # top to bottom

    # ─── TITLE BAND ─────────────────────────────────────────────
    title_bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                       Inches(0), Inches(0),
                                       Inches(BOARD_W), Inches(TITLE_H))
    title_bg.fill.solid()
    title_bg.fill.fore_color.rgb = TITLE_BAND
    title_bg.line.fill.background()

    # Gold accent line at bottom of title band
    gold_line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                        Inches(0), Inches(GOLD_LINE_Y),
                                        Inches(BOARD_W), Inches(0.03))
    gold_line.fill.solid()
    gold_line.fill.fore_color.rgb = ACCENT_GOLD
    gold_line.line.fill.background()

    # ─── LEFT TITLE ELEMENT: Single optimized house ────────────
    house_img = img_path("Image 0")
    if house_img:
        # Border/frame for the house image
        frame = add_rounded_rect(slide, 0.34, 0.14, 2.98, 2.76,
                                fill_color=None, border_color=TEXT_WHITE,
                                border_width=Pt(1.2), radius=IMG_RADIUS)
        add_image_safe(slide, house_img, 0.38, 0.18, 2.90, 2.68)
    # Label below
    add_text_box(slide, 0.34, 2.96, 2.98, 0.22, "SASTO Optimization",
                 font_size=Pt(11), italic=True, color=TEXT_WHITE,
                 align=PP_ALIGN.CENTER)

    # ─── TITLE TEXT (centered in middle area) ──────────────────
    # Row 1: Main title
    add_text_box(slide, 3.60, 0.10, 34.0, 1.05,
                 "SURROGATE-ACCELERATED STRUCTURAL OPTIMIZATION",
                 font_name="Arial Black", font_size=Pt(44), bold=True,
                 color=TEXT_WHITE, align=PP_ALIGN.CENTER)

    # Row 2: Subtitle
    add_text_box(slide, 3.60, 1.20, 34.0, 0.72,
                 "Additive Manufacturing: Harnessing FEA to Optimize Material Efficiency",
                 font_name="Arial", font_size=Pt(26), bold=True, italic=True,
                 color=TEXT_WHITE, align=PP_ALIGN.CENTER)

    # Row 3: Author
    add_text_box(slide, 3.60, 2.00, 34.0, 0.55,
                 "Eric Hou",
                 font_name="Arial", font_size=Pt(22), bold=True,
                 color=TEXT_WHITE, align=PP_ALIGN.CENTER)

    # ─── RIGHT CREDIT BLOCK ───────────────────────────────────
    credit_paras = [
        {"text": "Intel ISEF 2026 — Engineering Mechanics", "font_size": Pt(14), "bold": True, "color": TEXT_WHITE},
        {"text": "", "font_size": Pt(6), "bold": False, "color": TEXT_WHITE},
        {"text": "* Images from public domain or adapted from publicly available sources.",
         "font_size": Pt(12), "bold": False, "color": TEXT_WHITE},
        {"text": "All other graphics & figures created by Eric Hou, 2026.",
         "font_size": Pt(12), "bold": False, "color": TEXT_WHITE},
    ]
    add_multi_para_text(slide, 38.50, 0.18, 9.10, 2.88, credit_paras)

    print("  Title band complete")

    # ═══════════════════════════════════════════════════════════
    # LEFT PANEL
    # ═══════════════════════════════════════════════════════════
    lx = LEFT_X
    lw = LEFT_W
    ly = CONTENT_TOP

    # ─── L1: VISUAL ABSTRACT ──────────────────────────────────
    ly = add_section_header(slide, lx, ly, lw, "VISUAL ABSTRACT")
    va_h = 5.00
    add_card(slide, lx, ly, lw, va_h)

    # Visual abstract image (pipeline diagram)
    va_img = img_path("Image 1")
    if va_img:
        add_image_safe(slide, va_img, lx + CARD_PAD, ly + CARD_PAD,
                       lw - 2*CARD_PAD, va_h - 0.45)
    add_text_box(slide, lx + CARD_PAD, ly + va_h - 0.30, lw - 2*CARD_PAD, 0.25,
                 "Fig. 1. SASTO: 3D wireframe → volumetric parts → FEA dataset → surrogate → optimization → watertight STL",
                 font_size=Pt(9), italic=True, color=TEXT_DARK, align=PP_ALIGN.CENTER)
    ly += va_h + SECTION_GAP
    print("  L1 Visual Abstract complete")

    # ─── L2: INTRODUCTION ─────────────────────────────────────
    ly = add_section_header(slide, lx, ly, lw, "INTRODUCTION")
    intro_h = 5.80
    add_card(slide, lx, ly, lw, intro_h)

    # Text on left, images on right (Option A layout)
    txt_w = 6.20  # wider text column for readability
    img_w = lw - txt_w - 3*CARD_PAD

    intro_paras = [
        {"text": "▶ Concrete & Construction", "font_size": Pt(13), "bold": True, "color": ACCENT_TEAL, "space_after": 1},
        {"text": "Concrete = ~8% of global CO₂ [IEA 2021]. Uniform-thickness walls are determined by formwork, not structural need — wasting material.",
         "font_size": Pt(12), "bold": False, "color": TEXT_DARK, "space_after": 4},
        {"text": "▶ Additive Manufacturing", "font_size": Pt(13), "bold": True, "color": ACCENT_TEAL, "space_after": 1},
        {"text": "Large-scale 3D printing (ICON, COBOD, Apis Cor) enables topology-optimized structures with arbitrary wall profiles at no extra tooling cost.",
         "font_size": Pt(12), "bold": False, "color": TEXT_DARK, "space_after": 4},
        {"text": "▶ Computational Bottleneck", "font_size": Pt(13), "bold": True, "color": ACCENT_TEAL, "space_after": 1},
        {"text": "Classical SIMP: hundreds–thousands FEA solves, each minutes–hours at building scale. Produces disconnected mesh fragments.",
         "font_size": Pt(12), "bold": False, "color": TEXT_DARK, "space_after": 4},
    ]
    add_multi_para_text(slide, lx + CARD_PAD, ly + CARD_PAD,
                        txt_w, intro_h - 1.20, intro_paras)

    # Images on right side
    img2 = img_path("Image 2")
    img3 = img_path("Image 3")
    img_x = lx + txt_w + 2*CARD_PAD
    if img2:
        add_image_safe(slide, img2, img_x, ly + CARD_PAD, img_w, 2.00)
    if img3:
        # Image 3 is portrait (ratio 0.70) — place in tighter box
        add_image_safe(slide, img3, img_x, ly + CARD_PAD + 2.10, img_w, 2.00)

    # Key insight box at bottom
    insight_y = ly + intro_h - 1.10
    insight = add_rounded_rect(slide, lx + CARD_PAD, insight_y,
                               lw - 2*CARD_PAD, 0.80,
                               fill_color=EQ_BG, border_color=ACCENT_TEAL,
                               border_width=Pt(1.5), radius=35000)
    add_text_box(slide, lx + 2*CARD_PAD, insight_y + 0.08,
                 lw - 4*CARD_PAD, 0.64,
                 "SASTO replaces FEA with a deep ensemble surrogate achieving 23–92× speedup and zero constraint violations across 1,114 designs.",
                 font_size=Pt(11), bold=True, color=TEXT_DARK, align=PP_ALIGN.CENTER)

    ly += intro_h + SECTION_GAP
    print("  L2 Introduction complete")

    # ─── L3: RESEARCH OBJECTIVES ──────────────────────────────
    ly = add_section_header(slide, lx, ly, lw, "RESEARCH OBJECTIVES")
    obj_h = 3.60
    add_card(slide, lx, ly, lw, obj_h)

    # Objective rows
    objectives = [
        ("1", "SPEED", "23–92× speedup via deep ensemble surrogate", ACCENT_TEAL),
        ("2", "PRINTABILITY", "6-connected single-mesh watertight guarantee", ACCENT_TEAL),
        ("3", "EFFICIENCY", "23.5% mean material reduction (up to 45%)", ACCENT_GOLD),
        ("4", "SAFETY", "0/1,114 FEA violations; P(violation) ≤ 0.09%", ACCENT_RED),
    ]

    oy = ly + CARD_PAD
    for num, title, desc, badge_color in objectives:
        # Number badge
        badge = add_rounded_rect(slide, lx + CARD_PAD, oy, 0.28, 0.28,
                                fill_color=badge_color, border_color=None, radius=50000)
        add_text_box(slide, lx + CARD_PAD, oy, 0.28, 0.28, num,
                     font_size=Pt(12), bold=True, color=TEXT_WHITE, align=PP_ALIGN.CENTER)
        # Title + description
        add_text_box(slide, lx + CARD_PAD + 0.38, oy, 1.60, 0.26, title,
                     font_size=Pt(12), bold=True, color=badge_color, align=PP_ALIGN.LEFT)
        add_text_box(slide, lx + CARD_PAD + 2.00, oy, lw - 2*CARD_PAD - 2.00, 0.26, desc,
                     font_size=Pt(11), bold=False, color=TEXT_DARK, align=PP_ALIGN.LEFT)
        oy += 0.36

    # Pipeline boxes at bottom
    pipe_y = oy + 0.15
    pipe_labels = ["1. Data\nGeneration", "2. Surrogate\nTraining", "3. SASTO\nOptimize", "4. Validate\n& Certify"]
    pipe_w = 2.40
    pipe_gap = 0.30
    pipe_total = len(pipe_labels) * pipe_w + (len(pipe_labels) - 1) * pipe_gap
    pipe_start = lx + (lw - pipe_total) / 2

    for i, label in enumerate(pipe_labels):
        px = pipe_start + i * (pipe_w + pipe_gap)
        box = add_rounded_rect(slide, px, pipe_y, pipe_w, 0.70,
                              fill_color=ACCENT_TEAL, border_color=None, radius=35000)
        add_text_box(slide, px, pipe_y + 0.05, pipe_w, 0.60, label,
                     font_size=Pt(10), bold=True, color=TEXT_WHITE, align=PP_ALIGN.CENTER)

        # Arrow between boxes
        if i < len(pipe_labels) - 1:
            ax_x = px + pipe_w + 0.02
            add_text_box(slide, ax_x, pipe_y + 0.18, pipe_gap - 0.04, 0.30, "►",
                         font_size=Pt(14), bold=True, color=ACCENT_GOLD, align=PP_ALIGN.CENTER)

    ly += obj_h + SECTION_GAP
    print("  L3 Research Objectives complete")

    # ─── L4: ENGINEERING DESIGN CRITERIA ──────────────────────
    ly = add_section_header(slide, lx, ly, lw, "ENGINEERING DESIGN CRITERIA")
    crit_h = 4.20
    add_card(slide, lx, ly, lw, crit_h)

    # Table
    crit_data = [
        ["Constraint", "Limit", "Basis"],
        ["Von Mises stress", "σ_VM ≤ 5.0 MPa", "f'c/(γ_m×γ_f)"],
        ["Compliance ratio", "C_opt/C_base ≤ 1.15", "15% stiffness limit"],
        ["Displacement", "u_max ≤ L/360 ≈ 28 mm", "ASCE 7-22"],
        ["Wall t (exterior)", "2Δx ≈ 156 mm", "Load path"],
        ["Wall t (interior)", "1Δx ≈ 78 mm", "Non-structural"],
        ["Mesh integrity", "1 component", "Printability"],
    ]
    col_ws = [3.80, 3.50, 3.50]
    _, tbl_bot = add_table(slide, lx + CARD_PAD, ly + CARD_PAD,
                           lw - 2*CARD_PAD, crit_data, col_ws)

    # Equation pill below table
    eq_y = tbl_bot + 0.15
    eq = add_rounded_rect(slide, lx + CARD_PAD, eq_y,
                          lw - 2*CARD_PAD, 0.50,
                          fill_color=EQ_BG, border_color=None, radius=30000)
    add_text_box(slide, lx + 2*CARD_PAD, eq_y + 0.04,
                 lw - 4*CARD_PAD, 0.42,
                 "σ_VM,allow = f'c / (γ_m × γ_f) = 30 MPa / (3.0 × 2.0) = 5.0 MPa\nConservative: μ_σ + k·σ_σ,  k = 1.0  |  P(violation) ≤ 0.09%",
                 font_size=Pt(10), bold=False, color=TEXT_DARK, align=PP_ALIGN.CENTER)

    ly += crit_h + SECTION_GAP
    print("  L4 Design Criteria complete")

    # ─── L5: PROBLEM FRAMING ──────────────────────────────────
    ly = add_section_header(slide, lx, ly, lw, "PROBLEM FRAMING")
    pf_h = CONTENT_BOT - ly  # fill remaining space
    add_card(slide, lx, ly, lw, pf_h)

    # Optimization objective
    add_text_box(slide, lx + CARD_PAD, ly + CARD_PAD, lw - 2*CARD_PAD, 0.24,
                 "Optimization Objective", font_size=Pt(14), bold=True, color=TEXT_DARK)

    eq1 = add_rounded_rect(slide, lx + CARD_PAD, ly + 0.50,
                           lw - 2*CARD_PAD, 0.50,
                           fill_color=EQ_BG, border_color=None, radius=30000)
    add_text_box(slide, lx + 2*CARD_PAD, ly + 0.54,
                 lw - 4*CARD_PAD, 0.42,
                 "min J(ρ) = w_V·(V/V₀) + w_S·(S/V₀) + P_constraint(ρ)",
                 font_size=Pt(12), bold=False, color=TEXT_DARK, align=PP_ALIGN.CENTER)

    # Sensitivity formula
    add_text_box(slide, lx + CARD_PAD, ly + 1.10, lw - 2*CARD_PAD, 0.24,
                 "Sensitivity Formula", font_size=Pt(14), bold=True, color=TEXT_DARK)

    eq2 = add_rounded_rect(slide, lx + CARD_PAD, ly + 1.40,
                           lw - 2*CARD_PAD, 0.44,
                           fill_color=EQ_BG, border_color=None, radius=30000)
    add_text_box(slide, lx + 2*CARD_PAD, ly + 1.44,
                 lw - 4*CARD_PAD, 0.36,
                 "sᵢ = (1/5)·Σₘ ∂/∂ρᵢ[fₘ(C) + 0.3·fₘ(σ)]   |   sᵢ > 0: safe to remove   |   sᵢ < 0: essential",
                 font_size=Pt(10), bold=False, color=TEXT_DARK, align=PP_ALIGN.CENTER)

    # Part-aware thickness figure
    add_text_box(slide, lx + CARD_PAD, ly + 1.96, lw - 2*CARD_PAD, 0.24,
                 "Part-Aware Thickness Schematic", font_size=Pt(14), bold=True, color=TEXT_DARK)

    fig5_img = img_path("Image 5")
    if fig5_img:
        add_image_safe(slide, fig5_img, lx + CARD_PAD, ly + 2.26,
                       lw - 2*CARD_PAD, pf_h - 3.10)

    # Equation at bottom
    eq3_y = ly + pf_h - 0.65
    eq3 = add_rounded_rect(slide, lx + CARD_PAD, eq3_y,
                           lw - 2*CARD_PAD, 0.45,
                           fill_color=EQ_BG, border_color=None, radius=30000)
    add_text_box(slide, lx + 2*CARD_PAD, eq3_y + 0.04,
                 lw - 4*CARD_PAD, 0.37,
                 "t_min(p) = 2Δx = 156 mm  if p ∈ {ext, roof, floor}  |  1Δx = 78 mm  if p = interior wall",
                 font_size=Pt(10), bold=False, color=TEXT_DARK, align=PP_ALIGN.CENTER)

    print("  L5 Problem Framing complete")

    # ═══════════════════════════════════════════════════════════
    # CENTER PANEL
    # ═══════════════════════════════════════════════════════════
    cx = CENTER_X
    cw = CENTER_W
    cy = CONTENT_TOP

    # ─── C1: ENGINEERING METHODOLOGY ──────────────────────────
    cy = add_section_header(slide, cx, cy, cw, "ENGINEERING METHODOLOGY")
    c1_h = 15.80
    add_card(slide, cx, cy, cw, c1_h)

    # 2x2 sub-panel grid
    gutter = 0.20
    sub_w = (cw - 2*CARD_PAD - gutter) / 2  # ~11.27
    top_h = 7.00
    bot_h = c1_h - CARD_PAD - top_h - gutter - CARD_PAD

    # Sub-panel positions
    sp_x1 = cx + CARD_PAD
    sp_x2 = sp_x1 + sub_w + gutter
    sp_y1 = cy + CARD_PAD
    sp_y2 = sp_y1 + top_h + gutter

    # C1-A: Dataset Generation (top-left)
    add_rounded_rect(slide, sp_x1, sp_y1, sub_w, top_h,
                    fill_color=CARD_FILL, border_color=CARD_BORDER,
                    border_width=Pt(0.5), radius=40000)
    add_text_box(slide, sp_x1 + 0.10, sp_y1 + 0.08, sub_w - 0.20, 0.24,
                 "Dataset Generation Pipeline", font_size=Pt(13), bold=True, color=TEXT_DARK)

    ds_img = img_path("Image 6")
    if ds_img:
        add_image_safe(slide, ds_img, sp_x1 + 0.10, sp_y1 + 0.36,
                       sub_w - 0.20, 3.50)

    # Dataset table
    ds_data = [
        ["Split", "n", "Targets"],
        ["Train", "8,943", "σ_VM, u_max, C"],
        ["Validation", "1,121", "—"],
        ["Test", "1,114", "—"],
    ]
    ds_col_ws = [2.80, 2.50, 5.40]
    _, ds_bot = add_table(slide, sp_x1 + 0.10, sp_y1 + 3.96,
                          sub_w - 0.20, ds_data, ds_col_ws)

    add_text_box(slide, sp_x1 + 0.10, ds_bot + 0.05, sub_w - 0.20, 0.24,
                 "Fig. 6. 14,293 wireframes → 11,178 FEA sims → 128³ voxel grids",
                 font_size=Pt(9), italic=True, color=TEXT_DARK, align=PP_ALIGN.CENTER)

    # C1-B: Deep Ensemble (top-right) — Image 7 is portrait (ratio 0.62)
    add_rounded_rect(slide, sp_x2, sp_y1, sub_w, top_h,
                    fill_color=CARD_FILL, border_color=CARD_BORDER,
                    border_width=Pt(0.5), radius=40000)
    add_text_box(slide, sp_x2 + 0.10, sp_y1 + 0.08, sub_w - 0.20, 0.24,
                 "Deep Ensemble Surrogate (5×8.76M params)", font_size=Pt(13), bold=True, color=TEXT_DARK)

    # Ensemble badge
    badge = add_rounded_rect(slide, sp_x2 + sub_w - 1.80, sp_y1 + 0.06, 1.60, 0.26,
                            fill_color=ACCENT_RED, border_color=None, radius=30000)
    add_text_box(slide, sp_x2 + sub_w - 1.80, sp_y1 + 0.06, 1.60, 0.26,
                 "×5 Ensemble", font_size=Pt(10), bold=True, color=TEXT_WHITE, align=PP_ALIGN.CENTER)

    # Layout: image on left (portrait), text+table on right
    arch_img = img_path("Image 7")
    img7_w = 4.20  # portrait image gets narrower column
    txt7_x = sp_x2 + 0.10 + img7_w + 0.15
    txt7_w = sub_w - 0.20 - img7_w - 0.15
    if arch_img:
        add_image_safe(slide, arch_img, sp_x2 + 0.10, sp_y1 + 0.40,
                       img7_w, 5.40)

    # Architecture description on the right side
    arch_desc = [
        {"text": "Architecture:", "font_size": Pt(11), "bold": True, "color": TEXT_DARK, "space_after": 2},
        {"text": "• 3D CNN → Avg+Max dual-pool → 512-d embedding", "font_size": Pt(10), "bold": False, "color": TEXT_DARK, "space_after": 1},
        {"text": "• Load features concatenated after pool", "font_size": Pt(10), "bold": False, "color": TEXT_DARK, "space_after": 1},
        {"text": "• 3 regression heads: σ_VM, u_max, C", "font_size": Pt(10), "bold": False, "color": TEXT_DARK, "space_after": 4},
        {"text": "Training:", "font_size": Pt(11), "bold": True, "color": TEXT_DARK, "space_after": 2},
        {"text": "• Huber loss (SmoothL1)", "font_size": Pt(10), "bold": False, "color": TEXT_DARK, "space_after": 1},
        {"text": "• AdamW, lr = 5×10⁻⁴", "font_size": Pt(10), "bold": False, "color": TEXT_DARK, "space_after": 1},
        {"text": "• 90° rotations + flips + noise σ=0.02", "font_size": Pt(10), "bold": False, "color": TEXT_DARK, "space_after": 1},
        {"text": "• Input: log(1+|y|) → z-score", "font_size": Pt(10), "bold": False, "color": TEXT_DARK, "space_after": 4},
    ]
    add_multi_para_text(slide, txt7_x, sp_y1 + 0.40, txt7_w, 3.50, arch_desc)

    # Hyperparameter table (full width below image)
    hp_data = [
        ["Hyperparameter", "Value"],
        ["Loss", "Huber (SmoothL1)"],
        ["Optimizer", "AdamW, lr = 5e-4"],
        ["Augmentation", "90° rotations, flips, noise σ = 0.02"],
        ["Input", "log(1+|y|) → z-score"],
    ]
    hp_col_ws = [4.50, 6.20]
    _, hp_bot = add_table(slide, sp_x2 + 0.10, sp_y1 + 5.90,
                          sub_w - 0.20, hp_data, hp_col_ws)

    add_text_box(slide, sp_x2 + 0.10, hp_bot + 0.05, sub_w - 0.20, 0.24,
                 "Fig. 7. Dual pooling → 512-d embed fused with load features → 3 outputs",
                 font_size=Pt(9), italic=True, color=TEXT_DARK, align=PP_ALIGN.CENTER)

    # C1-C: SASTO Algorithm (bottom-left)
    add_rounded_rect(slide, sp_x1, sp_y2, sub_w, bot_h,
                    fill_color=CARD_FILL, border_color=CARD_BORDER,
                    border_width=Pt(0.5), radius=40000)
    add_text_box(slide, sp_x1 + 0.10, sp_y2 + 0.08, sub_w - 0.20, 0.24,
                 "SASTO Algorithm — Sensitivity-Guided Erosion", font_size=Pt(13), bold=True, color=TEXT_DARK)

    algo_img = img_path("Image 8")
    if algo_img:
        add_image_safe(slide, algo_img, sp_x1 + 0.10, sp_y2 + 0.40,
                       sub_w - 0.20, bot_h - 1.80)

    # Phase labels below the figure
    phase_y = sp_y2 + bot_h - 1.30
    phases = [
        ("Phase 1: Sensitivity-Guided Erosion (>99% of removal)", ACCENT_TEAL),
        ("Phase 2: Endgame (B=5, then 1)", ACCENT_GOLD),
        ("Phase 3: Swap Moves (interior ↔ removed neighbor)", ACCENT_RED),
    ]
    for i, (plabel, pcolor) in enumerate(phases):
        py = phase_y + i * 0.28
        pbar = add_rounded_rect(slide, sp_x1 + 0.10, py, sub_w - 0.20, 0.24,
                               fill_color=pcolor, border_color=None, radius=20000)
        add_text_box(slide, sp_x1 + 0.20, py, sub_w - 0.40, 0.24, plabel,
                     font_size=Pt(9), bold=True, color=TEXT_WHITE, align=PP_ALIGN.LEFT)

    # Trust region note
    tr_y = phase_y + 0.88
    tr = add_rounded_rect(slide, sp_x1 + 0.10, tr_y, sub_w - 0.20, 0.38,
                          fill_color=EQ_BG, border_color=None, radius=25000)
    add_text_box(slide, sp_x1 + 0.20, tr_y + 0.04, sub_w - 0.40, 0.30,
                 "Trust-region: constraint violated → undo, halve B → max(B/2, 10)",
                 font_size=Pt(9), bold=False, color=TEXT_DARK, align=PP_ALIGN.CENTER)

    add_text_box(slide, sp_x1 + 0.10, sp_y2 + bot_h - 0.30, sub_w - 0.20, 0.24,
                 "Fig. 8. Phase 1 provides >99% of removal. Phases 2–3 squeeze remaining feasible voxels.",
                 font_size=Pt(9), italic=True, color=TEXT_DARK, align=PP_ALIGN.CENTER)

    # C1-D: 6-Connectivity (bottom-right)
    add_rounded_rect(slide, sp_x2, sp_y2, sub_w, bot_h,
                    fill_color=CARD_FILL, border_color=CARD_BORDER,
                    border_width=Pt(0.5), radius=40000)
    add_text_box(slide, sp_x2 + 0.10, sp_y2 + 0.08, sub_w - 0.20, 0.24,
                 "Topology: 6-Connectivity Guarantee", font_size=Pt(13), bold=True, color=TEXT_DARK)

    # Side-by-side comparison
    half_w = (sub_w - 0.40) / 2 - 0.05
    comp_y = sp_y2 + 0.40

    # Left: 26-conn fails
    fail_frame = add_rounded_rect(slide, sp_x2 + 0.10, comp_y, half_w, 3.20,
                                 fill_color=CARD_FILL, border_color=ACCENT_RED,
                                 border_width=Pt(1.5), radius=30000)
    add_text_box(slide, sp_x2 + 0.10, comp_y + 0.02, half_w, 0.22,
                 "26-CONN — FAILS", font_size=Pt(10), bold=True, color=ACCENT_RED, align=PP_ALIGN.CENTER)
    fail_img = img_path("Image 9")
    if fail_img:
        add_image_safe(slide, fail_img, sp_x2 + 0.20, comp_y + 0.28, half_w - 0.20, 2.10)
    add_text_box(slide, sp_x2 + 0.10, comp_y + 2.44, half_w, 0.70,
                 "✗ Thousands of floating\nfragments — Unprintable",
                 font_size=Pt(10), bold=True, color=ACCENT_RED, align=PP_ALIGN.CENTER)

    # vs. divider
    add_text_box(slide, sp_x2 + 0.10 + half_w, comp_y + 1.20, 0.30, 0.30,
                 "vs.", font_size=Pt(12), bold=True, color=ACCENT_RED, align=PP_ALIGN.CENTER)

    # Right: 6-conn works
    pass_frame = add_rounded_rect(slide, sp_x2 + 0.10 + half_w + 0.10, comp_y,
                                 half_w, 3.20,
                                 fill_color=CARD_FILL, border_color=ACCENT_TEAL,
                                 border_width=Pt(1.5), radius=30000)
    add_text_box(slide, sp_x2 + 0.10 + half_w + 0.10, comp_y + 0.02, half_w, 0.22,
                 "6-CONN — WORKS", font_size=Pt(10), bold=True, color=ACCENT_TEAL, align=PP_ALIGN.CENTER)
    pass_img = img_path("Image 10")
    if pass_img:
        add_image_safe(slide, pass_img, sp_x2 + 0.20 + half_w + 0.10, comp_y + 0.28,
                       half_w - 0.20, 2.10)
    add_text_box(slide, sp_x2 + 0.10 + half_w + 0.10, comp_y + 2.44, half_w, 0.70,
                 "✓ 1 connected component\nWatertight STL confirmed",
                 font_size=Pt(10), bold=True, color=ACCENT_TEAL, align=PP_ALIGN.CENTER)

    # Proposition box
    prop_y = comp_y + 3.30
    prop = add_rounded_rect(slide, sp_x2 + 0.10, prop_y, sub_w - 0.20, 0.50,
                           fill_color=EQ_BG, border_color=ACCENT_TEAL,
                           border_width=Pt(1), radius=25000)
    add_text_box(slide, sp_x2 + 0.20, prop_y + 0.06, sub_w - 0.40, 0.38,
                 "Proposition: A binary voxel field with exactly one 6-connected foreground component yields a single-component marching-cubes surface mesh.",
                 font_size=Pt(9), bold=True, color=TEXT_DARK, align=PP_ALIGN.CENTER)

    # Face vs corner diagram
    diag_y = prop_y + 0.58
    face_box = add_rounded_rect(slide, sp_x2 + 0.10, diag_y, half_w, 0.40,
                               fill_color=CARD_FILL, border_color=ACCENT_TEAL,
                               border_width=Pt(1), radius=20000)
    add_text_box(slide, sp_x2 + 0.10, diag_y + 0.06, half_w, 0.28,
                 "FACE-share (6-adj) = printable ✓",
                 font_size=Pt(9), bold=True, color=ACCENT_TEAL, align=PP_ALIGN.CENTER)

    corner_box = add_rounded_rect(slide, sp_x2 + 0.10 + half_w + 0.10, diag_y,
                                 half_w, 0.40,
                                 fill_color=CARD_FILL, border_color=ACCENT_RED,
                                 border_width=Pt(1), radius=20000)
    add_text_box(slide, sp_x2 + 0.10 + half_w + 0.10, diag_y + 0.06, half_w, 0.28,
                 "CORNER-share only = fragment ✗",
                 font_size=Pt(9), bold=True, color=ACCENT_RED, align=PP_ALIGN.CENTER)

    add_text_box(slide, sp_x2 + 0.10, diag_y + 0.46, sub_w - 0.20, 0.24,
                 "Fig. 9. 6-connectivity eliminates floating fragments incompatible with AM toolpath generation.",
                 font_size=Pt(9), italic=True, color=TEXT_DARK, align=PP_ALIGN.CENTER)

    cy += c1_h + SECTION_GAP
    print("  C1 Methodology complete")

    # ─── C2: RESULTS & IN-SILICO VALIDATION ──────────────────
    cy = add_section_header(slide, cx, cy, cw, "RESULTS & IN-SILICO VALIDATION")
    c2_h = CONTENT_BOT - cy
    add_card(slide, cx, cy, cw, c2_h)

    # 3-column layout
    col_gap = 0.20
    col_w = (cw - 2*CARD_PAD - 2*col_gap) / 3

    col_x = [
        cx + CARD_PAD,
        cx + CARD_PAD + col_w + col_gap,
        cx + CARD_PAD + 2*(col_w + col_gap)
    ]
    content_y = cy + CARD_PAD

    # ── Col A: Reference Case ──
    add_text_box(slide, col_x[0], content_y, col_w, 0.24,
                 "Reference Case (Sample 00472)", font_size=Pt(12), bold=True, color=TEXT_DARK)

    ref_img = img_path("Image 11")
    if ref_img:
        add_image_safe(slide, ref_img, col_x[0], content_y + 0.30, col_w, 3.60)

    # "-45% material" badge
    badge_y = content_y + 2.20
    badge = add_rounded_rect(slide, col_x[0] + col_w/2 - 0.80, badge_y,
                            1.60, 0.26,
                            fill_color=ACCENT_RED, border_color=None, radius=30000)
    add_text_box(slide, col_x[0] + col_w/2 - 0.80, badge_y, 1.60, 0.26,
                 "−45.0% material", font_size=Pt(10), bold=True, color=TEXT_WHITE, align=PP_ALIGN.CENTER)

    # Results table
    ref_data = [
        ["Metric", "Baseline", "SASTO-PA"],
        ["Vol. reduction", "—", "45.0% ★"],
        ["VM stress (Pa)", "3.08e6", "3.08e6 ✓"],
        ["Compliance ratio", "1.00", "1.004 ✓"],
        ["Mesh components", "1", "1 ✓"],
        ["Runtime", "—", "160 s"],
        ["EI Index", "—", "0.358 ★"],
    ]
    ref_col_ws = [col_w * 0.40, col_w * 0.28, col_w * 0.32]
    _, ref_bot = add_table(slide, col_x[0], content_y + 4.00,
                           col_w, ref_data, ref_col_ws)

    add_text_box(slide, col_x[0], ref_bot + 0.05, col_w, 0.24,
                 "Table 1. SASTO-PA: 10.7pp more than SASTO-U",
                 font_size=Pt(8), italic=True, color=TEXT_DARK, align=PP_ALIGN.CENTER)

    # ── Col B: Multi-Geometry ──
    add_text_box(slide, col_x[1], content_y, col_w, 0.24,
                 "1,114-Geometry Generalization", font_size=Pt(12), bold=True, color=TEXT_DARK)

    # Histogram (use regenerated figure)
    hist_fig = poster_fig("fig10_histogram")
    hist_img_old = img_path("Image 12")
    hist_img = hist_fig if hist_fig else hist_img_old
    if hist_img:
        add_image_safe(slide, hist_img, col_x[1], content_y + 0.30, col_w, 3.20)

    add_text_box(slide, col_x[1], content_y + 3.52, col_w, 0.24,
                 "Fig. 10. Volume reduction. n=1,114 | Mean: 23.5%±7.8%",
                 font_size=Pt(8), italic=True, color=TEXT_DARK, align=PP_ALIGN.CENTER)

    # Per-part retention
    add_text_box(slide, col_x[1], content_y + 3.85, col_w, 0.24,
                 "Per-Part Material Retention", font_size=Pt(12), bold=True, color=TEXT_DARK)

    part_fig = poster_fig("fig11_per_part")
    part_img_old = img_path("Image 13")
    part_img = part_fig if part_fig else part_img_old
    if part_img:
        add_image_safe(slide, part_img, col_x[1], content_y + 4.15, col_w, 2.80)

    add_text_box(slide, col_x[1], content_y + 6.99, col_w, 0.24,
                 "Fig. 11. Load-bearing >91% retained. Interior walls: primary target.",
                 font_size=Pt(8), italic=True, color=TEXT_DARK, align=PP_ALIGN.CENTER)

    # ── Col C: Speedup & Validation ──
    add_text_box(slide, col_x[2], content_y, col_w, 0.24,
                 "Speedup vs. SIMP", font_size=Pt(12), bold=True, color=TEXT_DARK)

    speed_fig = poster_fig("fig12_speedup")
    speed_img_old = img_path("Image 14")
    speed_img = speed_fig if speed_fig else speed_img_old
    if speed_img:
        add_image_safe(slide, speed_img, col_x[2], content_y + 0.30, col_w, 2.50)

    add_text_box(slide, col_x[2], content_y + 2.84, col_w, 0.24,
                 "Fig. 12. SIMP 64³: 94s; SASTO 128³: 50s. 23–92× faster.",
                 font_size=Pt(8), italic=True, color=TEXT_DARK, align=PP_ALIGN.CENTER)

    # FEA compliance
    add_text_box(slide, col_x[2], content_y + 3.22, col_w, 0.24,
                 "Independent FEA Re-analysis (n=1,114)", font_size=Pt(12), bold=True, color=TEXT_DARK)

    fea_fig = poster_fig("fig13_fea_compliance")
    fea_img_old = img_path("Image 15")
    fea_img = fea_fig if fea_fig else fea_img_old
    if fea_img:
        add_image_safe(slide, fea_img, col_x[2], content_y + 3.50, col_w, 3.10)

    # Badges on FEA chart
    badge1 = add_rounded_rect(slide, col_x[2] + 0.10, content_y + 3.60,
                              1.60, 0.26,
                              fill_color=ACCENT_TEAL, border_color=None, radius=30000)
    add_text_box(slide, col_x[2] + 0.10, content_y + 3.60, 1.60, 0.26,
                 "0/1,114 violations", font_size=Pt(9), bold=True, color=TEXT_WHITE, align=PP_ALIGN.CENTER)

    badge2 = add_rounded_rect(slide, col_x[2] + 0.10, content_y + 3.90,
                              1.80, 0.24,
                              fill_color=ACCENT_TEAL, border_color=None, radius=30000)
    add_text_box(slide, col_x[2] + 0.10, content_y + 3.90, 1.80, 0.24,
                 "P(violation) ≤ 0.09%", font_size=Pt(9), bold=True, color=TEXT_WHITE, align=PP_ALIGN.CENTER)

    add_text_box(slide, col_x[2], content_y + 6.64, col_w, 0.24,
                 "Fig. 13. All C_opt/C_base ≤ 1.15. Max: 1.004.",
                 font_size=Pt(8), italic=True, color=TEXT_DARK, align=PP_ALIGN.CENTER)

    # ── Bottom Stats Banner (Option A with icons) ──
    banner_h = 1.50
    banner_y = cy + c2_h - banner_h - CARD_PAD
    ban = add_rounded_rect(slide, cx + CARD_PAD, banner_y,
                          cw - 2*CARD_PAD, banner_h,
                          fill_color=SECTION_BAR, border_color=None, radius=40000)

    stats = [
        ("23.5%", "Mean material\nreduction"),
        ("23–92×", "Speedup\nvs. SIMP"),
        ("0/1,114", "FEA constraint\nviolations"),
        ("50 sec", "Median\noptimization time"),
    ]
    stat_w = (cw - 2*CARD_PAD) / 4
    for i, (number, label) in enumerate(stats):
        sx = cx + CARD_PAD + i * stat_w
        add_text_box(slide, sx, banner_y + 0.15, stat_w, 0.70, number,
                     font_name="Arial Black", font_size=Pt(32), bold=True,
                     color=ACCENT_GOLD, align=PP_ALIGN.CENTER)
        add_text_box(slide, sx, banner_y + 0.90, stat_w, 0.50, label,
                     font_size=Pt(10), bold=False, color=TEXT_WHITE, align=PP_ALIGN.CENTER)

        # Separator lines between cells
        if i > 0:
            sep = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                          Inches(sx), Inches(banner_y + 0.25),
                                          Inches(0.02), Inches(banner_h - 0.50))
            sep.fill.solid()
            sep.fill.fore_color.rgb = TEXT_WHITE
            sep.line.fill.background()

    print("  C2 Results complete")

    # ═══════════════════════════════════════════════════════════
    # RIGHT PANEL
    # ═══════════════════════════════════════════════════════════
    rx = RIGHT_X
    rw = RIGHT_W
    ry = CONTENT_TOP

    # ─── R1: STATISTICAL ANALYSIS ────────────────────────────
    ry = add_section_header(slide, rx, ry, rw, "STATISTICAL ANALYSIS")
    r1_h = 13.80
    add_card(slide, rx, ry, rw, r1_h)

    # R1-A: Surrogate metrics table
    add_text_box(slide, rx + CARD_PAD, ry + CARD_PAD, rw - 2*CARD_PAD, 0.24,
                 "Surrogate Model Performance", font_size=Pt(13), bold=True, color=TEXT_DARK)

    sur_data = [
        ["Target", "Spearman ρ", "R²_log", "MAPE%"],
        ["Von Mises stress", "0.737", "0.419", "37.4"],
        ["Displacement", "0.970 ★", "0.842", "10.9"],
        ["Compliance", "0.948 ★", "0.814", "18.5"],
    ]
    sur_col_ws = [3.50, 2.20, 2.10, 2.20]
    _, sur_bot = add_table(slide, rx + CARD_PAD, ry + 0.48,
                           rw - 2*CARD_PAD, sur_data, sur_col_ws)

    # Callout
    callout = add_rounded_rect(slide, rx + CARD_PAD, sur_bot + 0.08,
                               rw - 2*CARD_PAD, 0.36,
                               fill_color=EQ_BG, border_color=None, radius=25000)
    add_text_box(slide, rx + 2*CARD_PAD, sur_bot + 0.12,
                 rw - 4*CARD_PAD, 0.28,
                 "Ranking accuracy (Spearman ρ), not pointwise, drives optimization safety. Compliance ρ = 0.948.",
                 font_size=Pt(9), bold=False, color=TEXT_DARK, align=PP_ALIGN.CENTER)

    # R1-B: Convergence
    conv_y = sur_bot + 0.52
    add_text_box(slide, rx + CARD_PAD, conv_y, rw - 2*CARD_PAD, 0.24,
                 "Optimization Convergence", font_size=Pt(13), bold=True, color=TEXT_DARK)

    conv_fig = poster_fig("fig14_convergence")
    conv_img_old = img_path("Image 16")
    conv_img = conv_fig if conv_fig else conv_img_old
    if conv_img:
        add_image_safe(slide, conv_img, rx + CARD_PAD, conv_y + 0.28,
                       rw - 2*CARD_PAD, 3.00)
    add_text_box(slide, rx + CARD_PAD, conv_y + 3.32, rw - 2*CARD_PAD, 0.24,
                 "Fig. 14. SASTO-PA (teal) vs. SASTO-U (gold). Part-aware thinning enables deeper removal.",
                 font_size=Pt(8), italic=True, color=TEXT_DARK, align=PP_ALIGN.CENTER)

    # R1-C: k-Factor
    kf_y = conv_y + 3.62
    add_text_box(slide, rx + CARD_PAD, kf_y, rw - 2*CARD_PAD, 0.24,
                 "k-Factor Sensitivity (Pareto Frontier)", font_size=Pt(13), bold=True, color=TEXT_DARK)

    kf_fig = poster_fig("fig15_k_factor")
    kf_img_old = img_path("Image 17")
    kf_img = kf_fig if kf_fig else kf_img_old
    if kf_img:
        add_image_safe(slide, kf_img, rx + CARD_PAD, kf_y + 0.28,
                       rw - 2*CARD_PAD, 2.70)

    # k=1.0 badge
    kbadge = add_rounded_rect(slide, rx + rw/2 - 0.40, kf_y + 1.60, 1.80, 0.24,
                             fill_color=ACCENT_GOLD, border_color=None, radius=20000)
    add_text_box(slide, rx + rw/2 - 0.40, kf_y + 1.60, 1.80, 0.24,
                 "k=1.0 Operating Point", font_size=Pt(9), bold=True, color=TEXT_WHITE, align=PP_ALIGN.CENTER)

    add_text_box(slide, rx + CARD_PAD, kf_y + 3.02, rw - 2*CARD_PAD, 0.24,
                 "Fig. 15. Non-monotonic: both gate and budget depend on k. k=1.0 is Pareto-optimal.",
                 font_size=Pt(8), italic=True, color=TEXT_DARK, align=PP_ALIGN.CENTER)

    # R1-D: Conformal Prediction
    cp_y = kf_y + 3.32
    add_text_box(slide, rx + CARD_PAD, cp_y, rw - 2*CARD_PAD, 0.24,
                 "Conformal Prediction Calibration", font_size=Pt(13), bold=True, color=TEXT_DARK)

    unc_fig = poster_fig("fig16_uncertainty")
    unc_img_old = img_path("Image 18")
    unc_img = unc_fig if unc_fig else unc_img_old
    if unc_img:
        add_image_safe(slide, unc_img, rx + CARD_PAD, cp_y + 0.28,
                       rw - 2*CARD_PAD, 2.40)

    add_text_box(slide, rx + CARD_PAD, cp_y + 2.72, rw - 2*CARD_PAD, 0.24,
                 "Fig. 16. Uncertainty bands widen. Γ_D ≈ 0.184 (ref). P(violation) ≤ 0.09%.",
                 font_size=Pt(8), italic=True, color=TEXT_DARK, align=PP_ALIGN.CENTER)

    ry += r1_h + SECTION_GAP
    print("  R1 Statistical Analysis complete")

    # ─── R2: CONCLUSIONS ──────────────────────────────────────
    ry = add_section_header(slide, rx, ry, rw, "CONCLUSIONS")
    r2_h = 5.40
    add_card(slide, rx, ry, rw, r2_h)

    conclusions = [
        "SASTO achieves 23.5%±7.8% mean material reduction across 1,114 held-out geometries, up to 45.0% on individual designs.",
        "Deep ensemble surrogate provides 23–92× speedup vs. SIMP: median 50 s on consumer GPU vs. 19–77 min for SIMP.",
        "6-connectivity eliminates thousands of floating mesh fragments (26-conn), guaranteeing watertight single-component STLs.",
        "Part-aware thickness yields 10.7pp more reduction than uniform by permitting 1-voxel interior walls.",
        "Independent FEA re-analysis of all 1,114 designs: zero violations (max 1.004). Conformal: P(violation) ≤ 0.09%.",
    ]

    coy = ry + CARD_PAD
    for i, text in enumerate(conclusions):
        num = str(i + 1)
        # Number badge
        add_rounded_rect(slide, rx + CARD_PAD, coy, 0.26, 0.26,
                        fill_color=ACCENT_TEAL, border_color=None, radius=50000)
        add_text_box(slide, rx + CARD_PAD, coy, 0.26, 0.26, num,
                     font_size=Pt(11), bold=True, color=TEXT_WHITE, align=PP_ALIGN.CENTER)
        add_text_box(slide, rx + CARD_PAD + 0.34, coy, rw - 2*CARD_PAD - 0.34, 0.76, text,
                     font_size=Pt(11), bold=False, color=TEXT_DARK)
        coy += 0.86

    # Impact strip
    impact_y = ry + r2_h - 0.42
    imp = add_rounded_rect(slide, rx + CARD_PAD, impact_y,
                          rw - 2*CARD_PAD, 0.36,
                          fill_color=IMPACT_BG, border_color=ACCENT_GOLD,
                          border_width=Pt(1), radius=25000)
    add_text_box(slide, rx + 2*CARD_PAD, impact_y + 0.04,
                 rw - 4*CARD_PAD, 0.28,
                 "🌍 8% of global CO₂ = cement. 23.5% less concrete per house → millions of tons saved at scale.",
                 font_size=Pt(9), bold=True, color=TEXT_DARK, align=PP_ALIGN.CENTER)

    ry += r2_h + SECTION_GAP
    print("  R2 Conclusions complete")

    # ─── R3: FUTURE WORK ──────────────────────────────────────
    ry = add_section_header(slide, rx, ry, rw, "FUTURE WORK")
    r3_h = 4.20
    add_card(slide, rx, ry, rw, r3_h)

    future_items = [
        ("1", "FEA-IN-THE-LOOP ACTIVE LEARNING",
         "When ensemble Γ_D > threshold τ, trigger ground-truth FEA mid-optimization. Self-correcting safety net."),
        ("2", "NONLINEAR FEA SPOT CHECKS",
         "Concrete damaged plasticity (CDP) on 5 designs to assess cracking/buckling in 78mm partitions."),
        ("3", "PHYSICAL PRINT VALIDATION",
         "1:10 scale print. Compression + DIC full-field strain. Compare to model predictions."),
    ]

    fwy = ry + CARD_PAD
    for num, title, desc in future_items:
        add_rounded_rect(slide, rx + CARD_PAD, fwy, 0.26, 0.26,
                        fill_color=ACCENT_TEAL, border_color=None, radius=50000)
        add_text_box(slide, rx + CARD_PAD, fwy, 0.26, 0.26, num,
                     font_size=Pt(11), bold=True, color=TEXT_WHITE, align=PP_ALIGN.CENTER)
        add_text_box(slide, rx + CARD_PAD + 0.34, fwy, rw - 2*CARD_PAD - 0.34, 0.24, title,
                     font_size=Pt(11), bold=True, color=SECTION_BAR)
        add_text_box(slide, rx + CARD_PAD + 0.34, fwy + 0.28, rw - 2*CARD_PAD - 0.34, 0.70, desc,
                     font_size=Pt(10), bold=False, color=TEXT_DARK)
        fwy += 1.10

    # Print protocol pipeline
    pipe_y2 = fwy + 0.05
    pipe_labels2 = ["Optimized STL", "3D Print 1:10", "Compression + DIC"]
    pw2 = 3.20
    pg2 = 0.35
    pt2 = len(pipe_labels2) * pw2 + (len(pipe_labels2) - 1) * pg2
    ps2 = rx + (rw - pt2) / 2

    for i, label in enumerate(pipe_labels2):
        px2 = ps2 + i * (pw2 + pg2)
        add_rounded_rect(slide, px2, pipe_y2, pw2, 0.40,
                        fill_color=ACCENT_TEAL, border_color=None, radius=25000)
        add_text_box(slide, px2, pipe_y2 + 0.05, pw2, 0.30, label,
                     font_size=Pt(9), bold=True, color=TEXT_WHITE, align=PP_ALIGN.CENTER)
        if i < len(pipe_labels2) - 1:
            add_text_box(slide, px2 + pw2 + 0.02, pipe_y2 + 0.05, pg2 - 0.04, 0.30, "►",
                         font_size=Pt(14), bold=True, color=ACCENT_GOLD, align=PP_ALIGN.CENTER)

    ry += r3_h + SECTION_GAP
    print("  R3 Future Work complete")

    # ─── R4: KEY REFERENCES ───────────────────────────────────
    ry = add_section_header(slide, rx, ry, rw, "KEY REFERENCES")
    r4_h = CONTENT_BOT - ry
    add_card(slide, rx, ry, rw, r4_h)

    refs = [
        "1. Bendsøe & Sigmund (2003). Topology Optimization. Springer.",
        "2. Lakshminarayanan et al. (2017). Deep ensembles. NeurIPS.",
        "3. Kong & Rosenfeld (1989). Digital topology. CVGIP 48(3).",
        "4. ASCE (2022). ASCE/SEI 7-22 Minimum Design Loads.",
        "5. Buswell et al. (2018). 3D printing concrete. Cem. Concr. Res. 112.",
        "6. Sigmund & Maute (2013). Topology optimization. S.M.O. 48.",
        "7. Lin et al. (2024). 3DWire: 3D Wireframe Dataset. KAUST.",
        "8. IEA (2021). Global Status Report — Buildings.",
        "9. Lorensen & Cline (1987). Marching Cubes. SIGGRAPH.",
        "10. Vovk et al. (2005). Algorithmic Learning. Springer.",
    ]

    ref_paras = []
    for ref in refs:
        ref_paras.append({"text": ref, "font_size": Pt(10), "bold": False,
                         "color": TEXT_DARK, "space_before": 1, "space_after": 1})

    add_multi_para_text(slide, rx + CARD_PAD, ry + CARD_PAD,
                        rw - 2*CARD_PAD, r4_h - 2*CARD_PAD, ref_paras)

    print("  R4 References complete")

    # ─── SAVE ─────────────────────────────────────────────────
    prs.save(OUTPUT)
    print(f"\n✓ Poster saved to: {OUTPUT}")
    print(f"  Size: {BOARD_W}×{BOARD_H} inches")
    print(f"  Shapes: ~{len(slide.shapes)}")


if __name__ == "__main__":
    build_poster()
