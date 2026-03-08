"""Detailed inspection of the user-modified v3(1) poster."""
import sys
sys.stdout.reconfigure(encoding='utf-8')
from pptx import Presentation
from pptx.util import Inches, Pt, Emu

prs = Presentation("SASTO_ISEF_Poster_v3 (1).pptx")
slide = prs.slides[0]
w = prs.slide_width / 914400
h = prs.slide_height / 914400
print(f"Slide: {w:.1f} x {h:.1f} inches")
print(f"Total shapes: {len(slide.shapes)}")

# Dump every shape with position, size, type, and text preview
for i, s in enumerate(slide.shapes):
    left = s.left / 914400
    top = s.top / 914400
    width = s.width / 914400
    height = s.height / 914400
    stype = s.shape_type
    
    text_preview = ""
    if s.has_text_frame:
        full = s.text_frame.text.replace('\n', ' | ').strip()
        text_preview = full[:100]
        # Get font info from first paragraph
        if s.text_frame.paragraphs:
            p = s.text_frame.paragraphs[0]
            fname = p.font.name or "?"
            fsize = p.font.size
            fsize_pt = fsize / 12700 if fsize else 0
            fbold = p.font.bold
            text_preview = f"[{fname} {fsize_pt:.0f}pt B={fbold}] {text_preview}"
    
    img_info = ""
    if stype == 13:  # picture
        img_info = f" IMG {width:.2f}x{height:.2f}in"
    
    tbl_info = ""
    if stype == 19:  # table
        tbl = s.table
        tbl_info = f" TABLE {tbl.rows.__len__()}x{len(tbl.columns)}"
    
    print(f"S{i:3d} [{left:6.2f},{top:6.2f}] {width:6.2f}x{height:5.2f} T={str(stype):20s}{img_info}{tbl_info} {text_preview}")
