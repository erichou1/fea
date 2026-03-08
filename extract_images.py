#!/usr/bin/env python3
"""
Extract all images from the existing PPTX so we can reuse them in the rebuild.
"""
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
import os

OUT = os.path.join(os.path.dirname(__file__), "poster_images_extracted")
os.makedirs(OUT, exist_ok=True)

prs = Presentation("SASTO_ISEF_Poster_v2.pptx")
slide = prs.slides[0]

for j, shape in enumerate(slide.shapes):
    try:
        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE or shape.shape_type == 13:
            img = shape.image
            ext = "png" if "png" in img.content_type else "jpg"
            fname = f"{shape.name}.{ext}"
            with open(os.path.join(OUT, fname), "wb") as f:
                f.write(img.blob)
            left = round(shape.left / 914400, 2)
            top = round(shape.top / 914400, 2)
            w = round(shape.width / 914400, 2)
            h = round(shape.height / 914400, 2)
            print(f"  Extracted: {fname:30s}  ({w}x{h} in, {len(img.blob)} bytes)  pos=({left},{top})")
    except Exception as e:
        print(f"  Skip shape {j}: {e}")

print(f"\nImages saved to: {OUT}")
